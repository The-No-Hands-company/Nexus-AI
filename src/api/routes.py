import os
import uuid
import json
import asyncio
import threading
import time
import hmac
import secrets
import hashlib
import base64
import re
from urllib import parse as _urlparse
from urllib import request as _urlrequest
from urllib import error as _urlerror
import jwt as _jwt
from datetime import datetime, timezone
from fastapi import Request, HTTPException, APIRouter
from fastapi.responses import FileResponse, StreamingResponse, HTMLResponse, JSONResponse, Response

router = APIRouter()
from ..agent import (run_agent_task, stream_agent_task, get_providers_list, get_provider_health, get_free_provider_diagnostics, get_config, update_config, call_llm_with_fallback, get_session_dir, get_system_resources, _config, PERSONAS, activity_log, _MAX_ACTIVITY, get_session_safety_profile, _push_safety_event, AllProvidersExhausted)
from ..approvals import list_tool_approvals, decide_tool_approval
from ..auth import JWT_SECRET, JWT_ALGO, JWT_EXPIRE_H, MULTI_USER
from ..scheduler import (
    schedule_job,
    list_jobs,
    cancel_job,
    job_to_dict,
    set_run_function,
    restore_from_db,
)
from ..db import (load_chats as db_load_chats, load_chat as db_load_chat, get_pinned_chats, get_usage_stats, get_usage_daily, save_custom_persona as db_save_persona, load_pref as db_load_pref, save_pref as db_save_pref, get_user as db_get_user, list_api_keys as db_list_api_keys, get_api_key_by_hash as db_get_api_key_by_hash, touch_api_key as db_touch_api_key, create_fine_tuning_job as db_create_fine_tuning_job, get_fine_tuning_job as db_get_fine_tuning_job, update_fine_tuning_job as db_update_fine_tuning_job, create_fine_tuning_job_event as db_create_fine_tuning_job_event, save_execution_trace as db_save_execution_trace, load_execution_trace as db_load_execution_trace, list_execution_traces as db_list_execution_traces, delete_execution_trace as db_delete_execution_trace, save_autonomy_trace as db_save_autonomy_trace, load_autonomy_trace as db_load_autonomy_trace)
from ..db import list_strict_clone_bypass_events as db_list_strict_clone_bypass_events, count_strict_clone_bypass_events as db_count_strict_clone_bypass_events
from ..personas import list_personas, set_persona, get_active_persona_name
from ..memory import (
    export_memory_bundle as export_memory_bundle,
    import_memory_bundle as import_memory_bundle,
)
from ..autonomy import Orchestrator, PlanningSystem, classify_subtask
from ..safety import GuardrailViolation, check_user_task
from ..safety_pipeline import (
    SAFETY_POLICY_PROFILES, get_safety_policy,
)
from ..execution_trace import (
    save_checkpoint as _save_checkpoint,
    list_traces as _list_traces,
    load_checkpoints as _load_checkpoints,
    get_latest_checkpoint as _get_latest_checkpoint,
    delete_trace as _delete_trace,
    save_file_diff as _save_file_diff,
    get_file_diffs as _get_file_diffs,
    get_file_diff_detail as _get_file_diff_detail,
)
from ..ensemble import get_ensemble_enabled
from ..profile_loader import inspect_profile_pack
from .schemas import *
from .state import (
    run_results,
    sessions,
    _session_requests,
    _reactions,
    _active_streams,
    autonomy_traces,
    execution_traces,
    get_rag_system,
)

_DEFAULT_AGENT_REQUEST_TIMEOUT_S = max(30.0, float(os.getenv("AGENT_REQUEST_TIMEOUT_S", "180")))


def _resolve_request_timeout_seconds(data: dict) -> float:
    """Resolve per-request timeout with sane bounds and env-configurable default."""
    raw = data.get("request_timeout_s")
    if raw is None:
        return _DEFAULT_AGENT_REQUEST_TIMEOUT_S
    try:
        return max(10.0, min(float(raw), 600.0))
    except (TypeError, ValueError):
        return _DEFAULT_AGENT_REQUEST_TIMEOUT_S

# ── API helpers ─────────────────────────────────────────────────────────────

def _api_error(message: str, code: str = "invalid_request", status_code: int = 400):
    return JSONResponse({"error": message, "type": code}, status_code=status_code)


def _v1_error(message: str, err_type: str = "invalid_request_error", status_code: int = 400, code: str = "invalid_request"):
    return JSONResponse(
        {
            "error": {
                "message": message,
                "type": err_type,
                "code": code,
                "status": status_code,
            },
            "message": message,
            "type": err_type,
            "code": code,
        },
        status_code=status_code,
    )


async def _read_json_body(request: Request, err_message: str = "invalid JSON body") -> dict:
    """Parse request JSON body and raise HTTPException on malformed payloads."""
    try:
        data = await request.json()
    except Exception:
        raise HTTPException(status_code=400, detail=err_message)
    if not isinstance(data, dict):
        raise HTTPException(status_code=400, detail=err_message)
    return data


def _normalize_response_format(response_format):
    if not response_format:
        return {"mode": None, "schema": None}

    if isinstance(response_format, str):
        normalized = response_format.strip().lower()
        if normalized == "json":
            return {"mode": "json", "schema": None}
        if normalized == "json_object":
            return {"mode": "json", "schema": None}
        raise ValueError("response_format must be 'json' or an object with type='json_schema'")

    if not isinstance(response_format, dict):
        raise ValueError("response_format must be a string or object")

    fmt_type = str(response_format.get("type", "")).strip().lower()
    if fmt_type == "json_object":
        return {"mode": "json", "schema": None}

    if fmt_type != "json_schema":
        raise ValueError("response_format.type must be 'json_object' or 'json_schema'")

    json_schema_cfg = response_format.get("json_schema")
    if not isinstance(json_schema_cfg, dict):
        raise ValueError("response_format.json_schema must be an object")

    schema = json_schema_cfg.get("schema")
    if not isinstance(schema, dict):
        raise ValueError("response_format.json_schema.schema must be an object")

    return {"mode": "json", "schema": schema}


def _builtin_chat_fallback(task: str, reason: str = "timeout") -> str:
    if reason == "provider_unavailable":
        return (
            "No model provider could be reached for this turn. "
            "Please check your provider keys in Settings and retry."
        )
    return (
        "The request timed out before a response was received. "
        "Please retry — if the problem persists, try a shorter message or check provider availability."
    )


def _builtin_coding_fallback(task: str, reason: str = "timeout") -> str:
    return ""


def _light_provider_precheck() -> dict:
    """Cheap availability snapshot used before first-turn execution."""
    providers = get_providers_list() or []
    ready = [p for p in providers if p.get("available")]
    cooling = [p for p in providers if p.get("rate_limited")]
    no_key = [p for p in providers if (not p.get("has_key") and not p.get("keyless"))]
    return {
        "total": len(providers),
        "ready": len(ready),
        "cooling": len(cooling),
        "no_key": len(no_key),
        "ready_labels": [str(p.get("label") or p.get("id") or "provider") for p in ready],
        "ts": time.time(),
    }


def _precheck_status_message(precheck: dict | None) -> str:
    if not precheck:
        return "Processing request..."
    ready = int(precheck.get("ready", 0) or 0)
    cooling = int(precheck.get("cooling", 0) or 0)
    if ready <= 0:
        return "Provider precheck: no providers currently ready. Request may fallback."
    return f"Provider precheck: {ready} ready" + (f", {cooling} cooling" if cooling else "") + "."


def _run_scheduled_task(task: str) -> str:
    """Execute a scheduled background task and return short result text."""
    sid = f"sched_{uuid.uuid4().hex[:8]}"
    result = run_agent_task(task, history=[], sid=sid)
    return str(result.get("result", ""))[:1200]


def _json_type_matches(value, expected_type: str) -> bool:
    if expected_type == "object":
        return isinstance(value, dict)
    if expected_type == "array":
        return isinstance(value, list)
    if expected_type == "string":
        return isinstance(value, str)
    if expected_type == "number":
        return isinstance(value, (int, float)) and not isinstance(value, bool)
    if expected_type == "integer":
        return isinstance(value, int) and not isinstance(value, bool)
    if expected_type == "boolean":
        return isinstance(value, bool)
    if expected_type == "null":
        return value is None
    return True


def _validate_json_schema_value(value, schema: dict, path: str = "$"):
    schema_type = schema.get("type")
    if schema_type and not _json_type_matches(value, str(schema_type)):
        raise ValueError(f"{path} expected type '{schema_type}'")

    if "enum" in schema and value not in schema.get("enum", []):
        raise ValueError(f"{path} must be one of the enum values")

    if schema_type == "object":
        properties = schema.get("properties", {})
        required = schema.get("required", [])
        if not isinstance(properties, dict):
            raise ValueError(f"{path} schema properties must be an object")
        if not isinstance(required, list):
            raise ValueError(f"{path} schema required must be an array")

        for key in required:
            if key not in value:
                raise ValueError(f"{path}.{key} is required")

        additional = schema.get("additionalProperties", True)
        if additional is False:
            unknown = [k for k in value.keys() if k not in properties]
            if unknown:
                raise ValueError(f"{path} has unknown keys: {', '.join(sorted(unknown))}")

        for key, child_schema in properties.items():
            if key in value and isinstance(child_schema, dict):
                _validate_json_schema_value(value[key], child_schema, f"{path}.{key}")

    if schema_type == "array":
        item_schema = schema.get("items")
        if isinstance(item_schema, dict):
            for idx, item in enumerate(value):
                _validate_json_schema_value(item, item_schema, f"{path}[{idx}]")


def _validate_json_output(text: str, schema: dict | None = None):
    candidate = (text or "").strip()
    if candidate.startswith("```"):
        lines = candidate.splitlines()
        if lines and lines[0].startswith("```"):
            lines = lines[1:]
        if lines and lines[-1].strip() == "```":
            lines = lines[:-1]
        candidate = "\n".join(lines).strip()

    try:
        parsed = json.loads(candidate)
        if schema:
            _validate_json_schema_value(parsed, schema)
        return parsed
    except Exception as exc:
        initial_error = exc

    starts = [idx for idx, ch in enumerate(candidate) if ch in "[{"]
    for start in starts:
        stack = []
        in_string = False
        escaped = False
        for idx in range(start, len(candidate)):
            ch = candidate[idx]
            if in_string:
                if escaped:
                    escaped = False
                elif ch == "\\":
                    escaped = True
                elif ch == '"':
                    in_string = False
                continue
            if ch == '"':
                in_string = True
                continue
            if ch in "[{":
                stack.append(ch)
                continue
            if ch in "]}":
                if not stack:
                    break
                opener = stack.pop()
                if (opener, ch) not in (("{", "}"), ("[", "]")):
                    break
                if not stack:
                    snippet = candidate[start:idx + 1]
                    try:
                        parsed = json.loads(snippet)
                        if schema:
                            _validate_json_schema_value(parsed, schema)
                        return parsed
                    except Exception:
                        break

    raise ValueError(str(initial_error))


def _provider_capability_flags(provider: dict) -> dict:
    provider_id = str(provider.get("id", "")).lower()
    model = str(provider.get("model", "")).lower()
    openai_compat = bool(provider.get("openai_compat", False))

    vision = provider_id in {"gemini", "ollama"} or any(token in model for token in (
        "vision", "llava", "bakllava", "gpt-4o", "gemini"
    ))
    embeddings = openai_compat
    json_mode = openai_compat
    reasoning = provider_id in {"ollama", "claude", "grok", "gemini"} or any(token in model for token in (
        "r1", "reason", "think"
    ))
    tools = True
    return {
        "tools": tools,
        "vision": vision,
        "embeddings": embeddings,
        "json_mode": json_mode,
        "reasoning": reasoning,
    }


def _provider_capabilities_list(flags: dict) -> list[str]:
    return [name for name, enabled in flags.items() if enabled]


def _principal_from_request(request: Request, sid: str = "", payload_user: str = "") -> str:
    """Resolve the best-effort caller identity for quota accounting."""
    token_user = _read_token(request)
    if token_user:
        return f"user:{token_user}"
    if payload_user:
        return f"openai_user:{str(payload_user).strip()[:80]}"
    if sid:
        return f"session:{sid}"
    forwarded = (request.headers.get("x-forwarded-for") or "").split(",")[0].strip()
    if forwarded:
        return f"ip:{forwarded}"
    if request.client and request.client.host:
        return f"ip:{request.client.host}"
    return "anonymous"


def _load_rate_limit_settings() -> dict:
    """Read persisted rate-limit policy with safe defaults."""
    defaults = {"mode": "soft", "per_minute": 60, "per_day": 2500}
    raw = db_load_pref("rate_limit_settings", "")
    if not raw:
        return dict(defaults)
    try:
        parsed = json.loads(raw)
    except Exception:
        return dict(defaults)

    mode = str(parsed.get("mode", defaults["mode"])).lower().strip()
    per_minute = int(parsed.get("per_minute", defaults["per_minute"]))
    per_day = int(parsed.get("per_day", defaults["per_day"]))
    if mode not in ("soft", "hard"):
        mode = defaults["mode"]
    per_minute = max(1, min(per_minute, 100000))
    per_day = max(1, min(per_day, 10000000))
    return {"mode": mode, "per_minute": per_minute, "per_day": per_day}


_rate_limit_settings = _load_rate_limit_settings()
_rate_limit_lock = threading.Lock()


def _evaluate_rate_limit(principal: str) -> dict:
    """Evaluate and record quota usage for a principal.

    Returns shape:
      {"allowed": bool, "mode": "soft|hard", "limit_type": "per_minute|per_day|", ...}
    """
    now = time.time()
    minute_limit = int(_rate_limit_settings.get("per_minute", 60))
    day_limit = int(_rate_limit_settings.get("per_day", 2500))
    mode = _rate_limit_settings.get("mode", "soft")
    minute_bucket = str(int(now // 60))
    day_bucket = str(int(now // 86400))

    try:
        from ..redis_state import get_rate_counter, incr_rate_counter

        current_minute = get_rate_counter(principal, f"m:{minute_bucket}")
        current_day = get_rate_counter(principal, f"d:{day_bucket}")
        minute_over = current_minute >= minute_limit
        day_over = current_day >= day_limit

        if minute_over or day_over:
            if mode == "hard":
                return {
                    "allowed": False,
                    "mode": mode,
                    "principal": principal,
                    "limit_type": "per_minute" if minute_over else "per_day",
                    "limit": minute_limit if minute_over else day_limit,
                    "used": current_minute if minute_over else current_day,
                    "retry_after_seconds": 60 if minute_over else 3600,
                }

            # Soft mode records pressure and still counts usage.
            new_minute = incr_rate_counter(principal, f"m:{minute_bucket}", 120)
            new_day = incr_rate_counter(principal, f"d:{day_bucket}", 172800)
            return {
                "allowed": True,
                "mode": mode,
                "principal": principal,
                "limit_type": "per_minute" if minute_over else "per_day",
                "limit": minute_limit if minute_over else day_limit,
                "used": new_minute if minute_over else new_day,
                "retry_after_seconds": 0,
            }

        # Allowed path: atomically record request in both windows.
        incr_rate_counter(principal, f"m:{minute_bucket}", 120)
        incr_rate_counter(principal, f"d:{day_bucket}", 172800)
        return {
            "allowed": True,
            "mode": mode,
            "principal": principal,
            "limit_type": "",
            "limit": 0,
            "used": 0,
            "retry_after_seconds": 0,
        }
    except Exception:
        # Fallback path keeps single-process behavior if Redis is unavailable.
        minute_window = 60.0
        day_window = 86400.0
        with _rate_limit_lock:
            entry = _session_requests.get(principal, {"minute": [], "day": []})
            minute = [t for t in entry.get("minute", []) if now - t < minute_window]
            day = [t for t in entry.get("day", []) if now - t < day_window]
            minute_over = len(minute) >= minute_limit
            day_over = len(day) >= day_limit
            is_over = minute_over or day_over

            if not is_over:
                minute.append(now)
                day.append(now)
                _session_requests[principal] = {"minute": minute, "day": day}
                return {
                    "allowed": True,
                    "mode": mode,
                    "principal": principal,
                    "limit_type": "",
                    "limit": 0,
                    "used": 0,
                    "retry_after_seconds": 0,
                }

            if mode == "soft":
                minute.append(now)
                day.append(now)
                _session_requests[principal] = {"minute": minute, "day": day}
                return {
                    "allowed": True,
                    "mode": mode,
                    "principal": principal,
                    "limit_type": "per_minute" if minute_over else "per_day",
                    "limit": minute_limit if minute_over else day_limit,
                    "used": len(minute) if minute_over else len(day),
                    "retry_after_seconds": 0,
                }

            retry_after = 1
            if minute_over and minute:
                retry_after = max(1, int(minute_window - (now - minute[0])))
            elif day_over and day:
                retry_after = max(1, int(day_window - (now - day[0])))
            _session_requests[principal] = {"minute": minute, "day": day}
            return {
                "allowed": False,
                "mode": mode,
                "principal": principal,
                "limit_type": "per_minute" if minute_over else "per_day",
                "limit": minute_limit if minute_over else day_limit,
                "used": len(minute) if minute_over else len(day),
                "retry_after_seconds": retry_after,
            }


def _quota_error_response(limit_result: dict) -> JSONResponse:
    body = {
        "error": "Quota exceeded for this user",
        "type": "quota_exceeded",
        "quota": {
            "mode": limit_result.get("mode", "hard"),
            "principal": limit_result.get("principal", "anonymous"),
            "limit_type": limit_result.get("limit_type", "per_minute"),
            "limit": limit_result.get("limit", 0),
            "used": limit_result.get("used", 0),
            "retry_after_seconds": limit_result.get("retry_after_seconds", 1),
        },
    }
    retry = str(limit_result.get("retry_after_seconds", 1))
    limit = str(limit_result.get("limit", 0))
    remaining = str(max(0, int(limit_result.get("limit", 0)) - int(limit_result.get("used", 0))))
    headers = {
        "Retry-After": retry,
        "X-RateLimit-Limit": limit,
        "X-RateLimit-Remaining": remaining,
        "X-RateLimit-Reset": str(int(time.time()) + int(limit_result.get("retry_after_seconds", 1))),
        "X-RateLimit-Policy": limit_result.get("limit_type", "per_minute"),
    }
    return JSONResponse(body, status_code=429, headers=headers)


JWT_REFRESH_EXPIRE_D = int(os.getenv("JWT_REFRESH_EXPIRE_D", "14"))

# ── Token revocation / session state backed by Redis with in-process fallback ─
# These in-memory structures are the fallback for when Redis is unavailable.
# Redis-backed equivalents are used when Redis is reachable, enabling correctness
# across Gunicorn workers and cold-start survival (zero-downtime rolling deploys).
_revoked_access_tokens: set[str] = set()
_refresh_tokens: dict[str, dict] = {}
_active_user_sessions: dict[str, list[dict]] = {}
_active_user_sessions_lock = threading.Lock()


def _redis_revoke_token(token: str, ttl_seconds: int = 0) -> None:
    """Mark a token as revoked in Redis (with expiry = token TTL)."""
    try:
        from ..redis_state import get_redis
        r = get_redis()
        key = f"nexus:revoked:{hashlib.sha256(token.encode()).hexdigest()[:32]}"
        if ttl_seconds > 0:
            r.set(key, "1", ex=ttl_seconds)
        else:
            r.set(key, "1", ex=86400 * 14)  # default 14 days
    except Exception:
        _revoked_access_tokens.add(token)


def _redis_is_revoked(token: str) -> bool:
    """Check if a token is revoked, consulting Redis first then in-process set."""
    try:
        from ..redis_state import get_redis
        r = get_redis()
        key = f"nexus:revoked:{hashlib.sha256(token.encode()).hexdigest()[:32]}"
        if r.get(key):
            return True
    except Exception:
        pass
    return token in _revoked_access_tokens


def _redis_save_refresh(token: str, data: dict, ttl_days: int = 14) -> None:
    try:
        from ..redis_state import get_redis
        r = get_redis()
        r.set(f"nexus:refresh:{token}", json.dumps(data), ex=ttl_days * 86400)
    except Exception:
        _refresh_tokens[token] = data


def _redis_get_refresh(token: str) -> dict | None:
    try:
        from ..redis_state import get_redis
        r = get_redis()
        raw = r.get(f"nexus:refresh:{token}")
        if raw:
            return json.loads(raw)
    except Exception:
        pass
    return _refresh_tokens.get(token)


def _redis_delete_refresh(token: str) -> None:
    try:
        from ..redis_state import get_redis
        r = get_redis()
        r.set(f"nexus:refresh:{token}", "", ex=1)
    except Exception:
        pass
    _refresh_tokens.pop(token, None)


def _redis_track_session(username: str, record: dict, max_sessions: int) -> list[str]:
    """
    Track an active session in Redis sorted set nexus:sessions:{username}.
    Returns list of revoked access tokens (oldest sessions trimmed to max_sessions).
    """
    revoked_tokens: list[str] = []
    try:
        from ..redis_state import get_redis
        import json as _json
        r = get_redis()
        key = f"nexus:sessions:{username}"
        score = float(record.get("issued_at", time.time()))
        r.set(f"{key}:v:{score}", _json.dumps(record), ex=JWT_REFRESH_EXPIRE_D * 86400)
        # We use a simple list stored in Redis for session tracking
        members_key = f"nexus:session_list:{username}"
        members_raw = r.get(members_key)
        sessions: list[dict] = _json.loads(members_raw) if members_raw else []
        sessions.append(record)
        sessions.sort(key=lambda x: float(x.get("issued_at", 0)))
        while len(sessions) > max_sessions:
            oldest = sessions.pop(0)
            old_access = str(oldest.get("access") or "")
            old_refresh = str(oldest.get("refresh") or "")
            if old_access:
                revoked_tokens.append(old_access)
                _redis_revoke_token(old_access)
            if old_refresh:
                _redis_delete_refresh(old_refresh)
        r.set(members_key, _json.dumps(sessions), ex=JWT_REFRESH_EXPIRE_D * 86400)
        return revoked_tokens
    except Exception:
        pass
    # Fallback to in-process
    with _active_user_sessions_lock:
        sessions = _active_user_sessions.get(username, [])
        sessions.append(record)
        sessions.sort(key=lambda x: float(x.get("issued_at", 0.0)))
        while len(sessions) > max_sessions:
            oldest = sessions.pop(0)
            old_access = str(oldest.get("access") or "")
            old_refresh = str(oldest.get("refresh") or "")
            if old_access:
                revoked_tokens.append(old_access)
                _revoked_access_tokens.add(old_access)
            if old_refresh:
                _refresh_tokens.pop(old_refresh, None)
        _active_user_sessions[username] = sessions
    return revoked_tokens


def _detect_suspicious_login(username: str, device_hash: str, ip: str) -> bool:
    """
    Detect suspicious logins by comparing current device/IP against known values.
    Returns True if the login looks suspicious (new device + new IP subnet).
    Records the current device/IP as known after the check.
    """
    try:
        from ..redis_state import get_redis
        r = get_redis()
        known_key = f"nexus:known_devices:{username}"
        last_ip_key = f"nexus:last_ip:{username}"

        known_devices_raw = r.get(known_key)
        known_devices: set[str] = set(json.loads(known_devices_raw)) if known_devices_raw else set()
        last_ip = (r.get(last_ip_key) or "").strip()

        is_new_device = device_hash not in known_devices
        # IP subnet change: compare first two octets for IPv4
        def _subnet(addr: str) -> str:
            parts = addr.split(".")
            return ".".join(parts[:2]) if len(parts) == 4 else addr
        is_new_subnet = bool(last_ip) and _subnet(last_ip) != _subnet(ip)

        suspicious = is_new_device and is_new_subnet

        # Update known devices (keep last 20)
        known_devices.add(device_hash)
        if len(known_devices) > 20:
            known_devices = set(list(known_devices)[-20:])
        r.set(known_key, json.dumps(list(known_devices)), ex=86400 * 90)
        r.set(last_ip_key, ip, ex=86400 * 90)
        return suspicious
    except Exception:
        return False  # graceful fallback: don't flag if Redis unavailable


def _v1_quota_error_response(limit_result: dict) -> JSONResponse:
    message = "Quota exceeded for this user"
    retry = str(limit_result.get("retry_after_seconds", 1))
    limit = str(limit_result.get("limit", 0))
    remaining = str(max(0, int(limit_result.get("limit", 0)) - int(limit_result.get("used", 0))))
    headers = {
        "Retry-After": retry,
        "X-RateLimit-Limit": limit,
        "X-RateLimit-Remaining": remaining,
        "X-RateLimit-Reset": str(int(time.time()) + int(limit_result.get("retry_after_seconds", 1))),
        "X-RateLimit-Policy": limit_result.get("limit_type", "per_minute"),
    }
    return JSONResponse(
        {
            "error": {
                "message": message,
                "type": "quota_exceeded",
                "code": "quota_exceeded",
                "status": 429,
            },
            "message": message,
            "type": "quota_exceeded",
            "code": "quota_exceeded",
            "quota": {
                "mode": limit_result.get("mode", "hard"),
                "principal": limit_result.get("principal", "anonymous"),
                "limit_type": limit_result.get("limit_type", "per_minute"),
                "limit": limit_result.get("limit", 0),
                "used": limit_result.get("used", 0),
                "retry_after_seconds": limit_result.get("retry_after_seconds", 1),
            },
        },
        status_code=429,
        headers=headers,
    )


# ── auth helpers ──────────────────────────────────────────────────────────────
def _hash_pw(password: str, salt: str = "") -> str:
    s = salt or secrets.token_hex(16)
    import binascii
    dk = hashlib.pbkdf2_hmac("sha256", password.encode(), (s + "nexus_ai_salt").encode(), 200000)
    return s + "$" + binascii.hexlify(dk).decode()

def _verify_pw(password: str, stored: str) -> bool:
    try:
        parts = stored.split("$")
        if len(parts) != 2:
            return False
        salt, _ = parts
        return secrets.compare_digest(stored, _hash_pw(password, salt))
    except Exception:
        return False

def _make_token(username: str, role: str | None = None) -> str:
    from time import time as _t
    user = db_get_user(username)
    role_value = str(role or (user.get("role", "user") if user else "user"))
    payload = {
        "sub": username,
        "role": role_value,
        "exp": int(_t()) + JWT_EXPIRE_H * 3600,
        "type": "access",
        "jti": secrets.token_hex(8),
    }
    return _jwt.encode(payload, JWT_SECRET, algorithm=JWT_ALGO)


def _make_refresh_token(username: str) -> str:
    from time import time as _t
    payload = {
        "sub": username,
        "exp": int(_t()) + JWT_REFRESH_EXPIRE_D * 86400,
        "type": "refresh",
        "jti": secrets.token_hex(8),
    }
    token = _jwt.encode(payload, JWT_SECRET, algorithm=JWT_ALGO)
    _redis_save_refresh(token, {"username": username, "exp": payload["exp"]}, ttl_days=JWT_REFRESH_EXPIRE_D)
    return token

def _orchestrator_llm(prompt: str, task: str = "") -> str:
    result, _pid = call_llm_with_fallback([{"role":"user","content":prompt}], task)
    if isinstance(result, dict):
        return result.get("content", str(result))
    return str(result)


def _save_autonomy_checkpoint(trace_id: str, step_idx: int, goal: str, events: list[dict]) -> None:
    """Persist autonomy execution checkpoint snapshots for replay/resume."""
    try:
        _save_checkpoint(trace_id, step_idx, {"task": goal, "trace_id": trace_id, "events": list(events)})
    except Exception:
        # Checkpointing should not fail the user-visible autonomy request.
        pass


def _hash_api_key(raw: str) -> str:
    return hashlib.sha256(raw.encode()).hexdigest()


def _client_ip(request: Request) -> str:
    forwarded = (request.headers.get("X-Forwarded-For") or "").split(",")[0].strip()
    if forwarded:
        return forwarded
    if request.client and request.client.host:
        return str(request.client.host)
    return "unknown"


def _device_hash(request: Request) -> str:
    ua = request.headers.get("User-Agent", "")
    ip = _client_ip(request)
    raw = f"{ua}|{ip}"
    return hashlib.sha256(raw.encode()).hexdigest()


def _register_user_session(username: str, access_token: str, refresh_token: str, request: Request) -> None:
    """Track active sessions per user; revoke oldest when over configured cap.
    Sessions tracked in Redis for cross-worker correctness; in-process dict as fallback.
    """
    max_sessions = int(os.getenv("MAX_SESSIONS_PER_USER", "5"))
    if max_sessions < 1:
        return
    record = {
        "access": access_token,
        "refresh": refresh_token,
        "issued_at": time.time(),
        "ip": _client_ip(request),
        "device_hash": _device_hash(request),
    }
    _redis_track_session(username, record, max_sessions)


def _read_token(request: Request) -> str | None:
    # 1. X-API-Key header
    raw_api_key = request.headers.get("X-API-Key", "").strip()
    if not raw_api_key:
        # Also allow Bearer nxk_... as API key
        hdr = request.headers.get("Authorization", "")
        if hdr.startswith("Bearer nxk_"):
            raw_api_key = hdr[7:]
    if raw_api_key and raw_api_key.startswith("nxk_"):
        key_hash = _hash_api_key(raw_api_key)
        key_record = db_get_api_key_by_hash(key_hash)
        if key_record:
            db_touch_api_key(key_record["id"], time.time())
            return key_record["username"]
        return None

    # 2. Standard JWT Bearer
    header = request.headers.get("Authorization", "")
    if not header.startswith("Bearer "):
        return None
    token = header[7:]
    if _redis_is_revoked(token):
        return None
    try:
        payload = _jwt.decode(token, JWT_SECRET, algorithms=[JWT_ALGO])
        if payload.get("type") not in (None, "access"):
            return None
        return payload.get("sub")
    except Exception:
        return None


def _get_request_api_key_scopes(request: Request) -> list[str]:
    """Return scopes for the API key used in this request, or [] if JWT auth."""
    raw_api_key = request.headers.get("X-API-Key", "").strip()
    if not raw_api_key:
        hdr = request.headers.get("Authorization", "")
        if hdr.startswith("Bearer nxk_"):
            raw_api_key = hdr[7:]
    if raw_api_key and raw_api_key.startswith("nxk_"):
        key_hash = _hash_api_key(raw_api_key)
        key_record = db_get_api_key_by_hash(key_hash)
        if key_record:
            return list(key_record.get("scopes") or [])
    return ["*"]  # JWT auth = all scopes

def require_auth(request: Request) -> str:
    from fastapi import HTTPException
    username = _read_token(request)
    if not username:
        raise HTTPException(status_code=401, detail="Unauthorized — valid Bearer token required")
    return username


def _get_token_role(request: Request) -> str:
    if not MULTI_USER:
        return "admin"
    header = request.headers.get("Authorization", "")
    if not header.startswith("Bearer "):
        return "guest"
    token = header[7:]
    try:
        payload = _jwt.decode(token, JWT_SECRET, algorithms=[JWT_ALGO])
        return payload.get("role", "user")
    except Exception:
        return "guest"


def require_admin(request: Request) -> str:
    if not MULTI_USER:
        return "nexus_admin"
    username = require_auth(request)
    role = _get_token_role(request)
    if role != "admin":
        raise HTTPException(status_code=403, detail="Admin access required")
    return username


def _require_admin(request: Request) -> str:
    return require_admin(request)

@router.get("/")
def home(): return FileResponse("static/index.html")


@router.get("/playground")
def api_playground():
        html = """
        <!doctype html>
        <html>
            <head>
                <meta charset=\"utf-8\" />
                <meta name=\"viewport\" content=\"width=device-width,initial-scale=1\" />
                <title>Nexus API Playground</title>
            </head>
            <body style=\"font-family: ui-monospace, SFMono-Regular, Menlo, monospace; margin: 20px;\">
                <h1>Nexus API Playground</h1>
                <p>Run quick requests against chat, agent, and autonomy endpoints from one page.</p>
                <label>Base URL <input id=\"base\" value=\"\" placeholder=\"leave blank for same origin\" style=\"width: 420px\"></label>
                <label style=\"margin-left: 8px\">Bearer token <input id=\"token\" placeholder=\"optional\" style=\"width: 340px\"></label>
                <div style=\"margin-top: 12px;\">
                    <button onclick=\"runModels()\">GET /v1/models</button>
                    <button onclick=\"runChat()\">POST /v1/chat/completions</button>
                    <button onclick=\"runAgent()\">POST /v1/agent</button>
                    <button onclick=\"runPlan()\">POST /v1/autonomy/plan</button>
                </div>
                <h3>Request</h3>
                <textarea id=\"req\" style=\"width: 100%; height: 190px;\">{
    \"model\": \"nexus-ai/auto\",
    \"messages\": [{\"role\": \"user\", \"content\": \"Say hello in one sentence.\"}],
    \"stream\": false
}</textarea>
                <h3>Response</h3>
                <pre id=\"out\" style=\"background:#111;color:#d8f2c2;padding:12px;overflow:auto;min-height:180px;\"></pre>
                <script>
                    const out = document.getElementById('out');
                    const reqEl = document.getElementById('req');
                    const baseEl = document.getElementById('base');
                    const tokenEl = document.getElementById('token');
                      function base() { return (baseEl.value || '').trim().replace(/\\/$/, ''); }
                    function headers() {
                        const h = {'Content-Type':'application/json'};
                        const t = (tokenEl.value || '').trim();
                        if (t) h.Authorization = `Bearer ${t}`;
                        return h;
                    }
                    async function run(method, path, body) {
                        const url = `${base()}${path}`;
                        try {
                            const res = await fetch(url, {
                                method,
                                headers: headers(),
                                body: body ? JSON.stringify(body) : undefined,
                            });
                            const text = await res.text();
                            let parsed;
                            try { parsed = JSON.parse(text); } catch (_) { parsed = text; }
                            out.textContent = JSON.stringify({status: res.status, path, body: parsed}, null, 2);
                        } catch (e) {
                            out.textContent = String(e);
                        }
                    }
                    function parseRequest() {
                        try { return JSON.parse(reqEl.value || '{}'); }
                        catch (e) { out.textContent = `Invalid JSON: ${e}`; return null; }
                    }
                    async function runModels() { await run('GET', '/v1/models'); }
                    async function runChat() {
                        const body = parseRequest(); if (!body) return;
                        await run('POST', '/v1/chat/completions', body);
                    }
                    async function runAgent() {
                        const body = { task: 'Summarize this repository architecture in 3 bullets.', history: [] };
                        await run('POST', '/v1/agent', body);
                    }
                    async function runPlan() {
                        const body = { goal: 'Plan a safe production rollout with checkpoints.', max_subtasks: 6 };
                        await run('POST', '/v1/autonomy/plan', body);
                    }
                </script>
            </body>
        </html>
        """
        return HTMLResponse(html)


def _dev_sandbox_enabled() -> bool:
        return os.getenv("NEXUS_DEV_SANDBOX", "0").strip().lower() in {"1", "true", "yes", "on"}


@router.get("/dev/sandbox/status")
def dev_sandbox_status():
        return {
                "enabled": _dev_sandbox_enabled(),
                "mode": "mock-llm",
                "notes": "Set NEXUS_DEV_SANDBOX=1 to enable deterministic mock responses for test environments.",
        }


@router.post("/dev/sandbox/chat")
async def dev_sandbox_chat(request: Request):
        if not _dev_sandbox_enabled():
                return _api_error("developer sandbox is disabled", "feature_disabled", 403)
        body = await _read_json_body(request)
        prompt = str(body.get("prompt") or "").strip()
        if not prompt:
                return _api_error("prompt is required", "validation_error", 422)

        lower_prompt = prompt.lower()
        if "json" in lower_prompt:
                content = {"sandbox": True, "echo": prompt[:200], "status": "ok"}
        elif "risk" in lower_prompt or "safety" in lower_prompt:
                content = "Sandbox response: risk identified, recommend HITL approval before high-impact actions."
        else:
                content = f"Sandbox response: {prompt[:240]}"

        return {
                "id": f"sandbox-{secrets.token_hex(6)}",
                "object": "chat.completion",
                "created": int(time.time()),
                "model": "nexus-sandbox/mock-llm",
                "choices": [
                        {
                                "index": 0,
                                "message": {"role": "assistant", "content": content},
                                "finish_reason": "stop",
                        }
                ],
                "usage": {
                        "prompt_tokens": max(1, len(prompt.split())),
                        "completion_tokens": max(1, len(str(content).split())),
                        "total_tokens": max(2, len(prompt.split()) + len(str(content).split())),
                },
                "sandbox": True,
        }


@router.get("/graphql")
def graphql_help():
    return {
        "endpoint": "/graphql",
        "method": "POST",
        "body": {"query": "{ health models providers usage }"},
        "notes": [
            "Supported root fields: health, models, providers, usage",
            "usage requires admin role in multi-user mode",
        ],
    }


def _extract_graphql_fields(query: str) -> set[str]:

    allowed = {"health", "models", "providers", "usage"}
    tokens = set(re.findall(r"\b[a-zA-Z_][a-zA-Z0-9_]*\b", str(query or "")))
    return {token for token in tokens if token in allowed}


@router.post("/graphql")
async def graphql_query(request: Request):
    body = await _read_json_body(request)
    query = str(body.get("query") or "").strip()
    if not query:
        return JSONResponse({"errors": [{"message": "query is required"}]}, status_code=400)

    wanted = _extract_graphql_fields(query)
    if not wanted:
        return JSONResponse(
            {"errors": [{"message": "No supported fields requested (health, models, providers, usage)"}]},
            status_code=400,
        )

    data: dict[str, Any] = {}
    errors: list[dict[str, str]] = []

    if "health" in wanted:
        data["health"] = {
            "status": "ok",
            "timestamp": datetime.now(timezone.utc).isoformat(),
            "version": "v1",
        }
    if "models" in wanted:
        from ..routes.v1 import _v1_models_catalog
        data["models"] = _v1_models_catalog()
    if "providers" in wanted:
        data["providers"] = get_provider_health()
    if "usage" in wanted:
        try:
            require_admin(request)
            data["usage"] = {
                "summary": get_usage_stats(7),
                "daily": get_usage_daily(7),
            }
        except HTTPException:
            errors.append({"message": "usage field requires admin role"})

    payload: dict[str, Any] = {"data": data}
    if errors:
        payload["errors"] = errors
    return payload


# ── Webhook trigger ─────────────────────────────────────────────────────────────
# POST /webhook/trigger  { "task": "fix the login bug", "repo": "owner/repo" }
# Runs the agent asynchronously and streams back SSE or returns a run_id for polling.
# Optional header: X-Webhook-Secret: <secret>  (validated against WEBHOOK_SECRET env var)

WEBHOOK_SECRET = os.getenv("WEBHOOK_SECRET", "")
WEBHOOK_HMAC_SECRET = os.getenv("WEBHOOK_HMAC_SECRET", "")


def _webhook_allowed_events() -> set[str]:
    raw = os.getenv("WEBHOOK_ALLOWED_EVENTS", "").strip()
    if not raw:
        return set()
    return {part.strip().lower() for part in raw.split(",") if part.strip()}


def _safe_int(value, default: int, min_value: int, max_value: int) -> int:
    try:
        parsed = int(value)
    except Exception:
        parsed = default
    return max(min_value, min(max_value, parsed))

@router.post("/webhook/trigger")
async def webhook_trigger(request: Request):
    raw_body = await request.body()
    try:
        body = json.loads(raw_body.decode("utf-8")) if raw_body else {}
    except Exception:
        return JSONResponse({"error": "invalid JSON body"}, status_code=400)
    if not isinstance(body, dict):
        return JSONResponse({"error": "invalid JSON body"}, status_code=400)

    secret = request.headers.get("x-webhook-secret", "")
    if WEBHOOK_SECRET and not hmac.compare_digest(secret, WEBHOOK_SECRET):
        return JSONResponse({"error": "invalid webhook secret"}, status_code=403)

    if WEBHOOK_HMAC_SECRET:
        signature_header = request.headers.get("x-webhook-signature-256", "").strip()
        if not signature_header.startswith("sha256="):
            return JSONResponse({"error": "missing or invalid webhook signature"}, status_code=403)
        provided_sig = signature_header.split("=", 1)[1].strip().lower()
        expected_sig = hmac.new(
            WEBHOOK_HMAC_SECRET.encode("utf-8"),
            raw_body,
            digestmod=hashlib.sha256,
        ).hexdigest()
        if not hmac.compare_digest(provided_sig, expected_sig):
            # Compatibility fallback for senders that sign canonicalized JSON.
            canonical_body = json.dumps(body, separators=(",", ":"), sort_keys=True).encode("utf-8")
            canonical_sig = hmac.new(
                WEBHOOK_HMAC_SECRET.encode("utf-8"),
                canonical_body,
                digestmod=hashlib.sha256,
            ).hexdigest()
            if not hmac.compare_digest(provided_sig, canonical_sig):
                return JSONResponse({"error": "invalid webhook signature"}, status_code=403)

    task = body.get("task", "")
    if not task:
        return JSONResponse({"error": "task field is required"}, status_code=400)
    try:
        task = check_user_task(task, policy_profile=_config.get("safety_profile", "standard"))
    except GuardrailViolation as exc:
        return _api_error(exc.reason, exc.code, 422)

    event_type = str(body.get("event_type") or request.headers.get("x-webhook-event") or "generic").strip().lower()
    allowed_events = _webhook_allowed_events()
    if allowed_events and event_type not in allowed_events:
        return JSONResponse(
            {
                "accepted": False,
                "ignored": True,
                "reason": "event_type_not_allowed",
                "event_type": event_type,
                "allowed_events": sorted(allowed_events),
            },
            status_code=202,
        )

    max_retries_default = _safe_int(os.getenv("WEBHOOK_MAX_RETRIES", "2"), default=2, min_value=0, max_value=5)
    max_retries = _safe_int(body.get("max_retries", max_retries_default), default=max_retries_default, min_value=0, max_value=5)
    backoff_secs_default = _safe_int(os.getenv("WEBHOOK_RETRY_BACKOFF_SECONDS", "1"), default=1, min_value=1, max_value=60)
    retry_backoff_secs = _safe_int(body.get("retry_backoff_secs", backoff_secs_default), default=backoff_secs_default, min_value=1, max_value=60)

    repo = body.get("repo", "")
    run_id = "run_" + secrets.token_hex(8)
    run_results[run_id] = {
        "status": "running",
        "result": None,
        "error": None,
        "event_type": event_type,
        "repo": repo,
        "attempt": 0,
        "max_retries": max_retries,
        "retry_backoff_secs": retry_backoff_secs,
        "errors": [],
    }

    def _run():
        attempt = 0
        errors = []
        try:
            while True:
                try:
                    attempt += 1
                    run_results[run_id]["status"] = "running" if attempt == 1 else "retrying"
                    run_results[run_id]["attempt"] = attempt
                    from ..agent import run_agent_task
                    result = run_agent_task(task, [], sid=run_id)
                    run_results[run_id] = {
                        "status": "done",
                        "result": result,
                        "error": None,
                        "event_type": event_type,
                        "repo": repo,
                        "attempt": attempt,
                        "max_retries": max_retries,
                        "retry_backoff_secs": retry_backoff_secs,
                        "errors": errors,
                    }
                    break
                except Exception as e:
                    err = str(e)
                    errors.append(err)
                    if attempt > max_retries:
                        run_results[run_id] = {
                            "status": "error",
                            "result": None,
                            "error": err,
                            "event_type": event_type,
                            "repo": repo,
                            "attempt": attempt,
                            "max_retries": max_retries,
                            "retry_backoff_secs": retry_backoff_secs,
                            "errors": errors,
                        }
                        break
                    delay = retry_backoff_secs * (2 ** (attempt - 1))
                    run_results[run_id]["next_retry_in_secs"] = delay
                    time.sleep(delay)
        except Exception as e:
            run_results[run_id] = {
                "status": "error",
                "result": None,
                "error": str(e),
                "event_type": event_type,
                "repo": repo,
                "attempt": attempt,
                "max_retries": max_retries,
                "retry_backoff_secs": retry_backoff_secs,
                "errors": errors,
            }

    threading.Thread(target=_run, daemon=True).start()
    return {
        "run_id": run_id,
        "status": "https://github.com/The-No-Hands-company/Nexus-AI#webhook-triggers",
        "event_type": event_type,
        "max_retries": max_retries,
        "retry_backoff_secs": retry_backoff_secs,
    }

@router.get("/webhook/status/{run_id}")
async def webhook_status(run_id: str):
    result = run_results.get(run_id)
    if not result: return JSONResponse({"error": "run_id not found"}, status_code=404)
    return result


def _http_json_request(
    method: str,
    url: str,
    payload: dict | None = None,
    headers: dict | None = None,
    timeout: int = 15,
) -> dict:
    data = None
    req_headers = {"Accept": "application/json"}
    if headers:
        req_headers.update({str(k): str(v) for k, v in headers.items()})
    if payload is not None:
        data = json.dumps(payload).encode("utf-8")
        req_headers.setdefault("Content-Type", "application/json")
    req = _urlrequest.Request(url=url, data=data, method=method.upper(), headers=req_headers)
    try:
        with _urlrequest.urlopen(req, timeout=timeout) as resp:
            body = resp.read().decode("utf-8", errors="replace")
            parsed = {}
            if body:
                try:
                    parsed = json.loads(body)
                except Exception:
                    parsed = {"raw": body}
            return {"ok": True, "status": int(resp.getcode()), "json": parsed, "text": body}
    except _urlerror.HTTPError as exc:
        body = exc.read().decode("utf-8", errors="replace") if hasattr(exc, "read") else ""
        parsed = {}
        if body:
            try:
                parsed = json.loads(body)
            except Exception:
                parsed = {"raw": body}
        return {"ok": False, "status": int(exc.code), "json": parsed, "text": body}
    except Exception as exc:
        return {"ok": False, "status": 0, "json": {}, "text": str(exc)}


def _is_truthy_env(name: str, default: bool = False) -> bool:
    val = os.getenv(name, "").strip().lower()
    if not val:
        return default
    return val in {"1", "true", "yes", "on"}


def _is_safe_callback_url(url: str) -> bool:
    try:
        parsed = _urlparse.urlparse(url)
    except Exception:
        return False
    if parsed.scheme == "https":
        return True
    if parsed.scheme != "http":
        return False
    host = (parsed.hostname or "").lower()
    return host in {"127.0.0.1", "localhost"}


def _start_integration_run(
    task: str,
    source: str,
    metadata: dict | None = None,
    callback=None,
) -> str:
    run_id = "ext_" + secrets.token_hex(8)
    run_results[run_id] = {
        "status": "running",
        "result": None,
        "error": None,
        "source": source,
        "metadata": dict(metadata or {}),
    }

    def _run() -> None:
        status_payload = run_results.get(run_id, {})
        try:
            result = run_agent_task(task, [], sid=run_id)
            status_payload = {
                "status": "done",
                "result": result,
                "error": None,
                "source": source,
                "metadata": dict(metadata or {}),
            }
        except Exception as exc:
            status_payload = {
                "status": "error",
                "result": None,
                "error": str(exc),
                "source": source,
                "metadata": dict(metadata or {}),
            }
        run_results[run_id] = status_payload
        if callback:
            try:
                callback(run_id, status_payload)
            except Exception:
                pass

    threading.Thread(target=_run, daemon=True).start()
    return run_id


def _slack_verify_signature(raw_body: bytes, timestamp: str, signature: str) -> bool:
    secret = os.getenv("SLACK_SIGNING_SECRET", "").strip()
    if not secret:
        return True
    if not timestamp or not signature:
        return False
    try:
        ts = int(timestamp)
    except Exception:
        return False
    if abs(int(time.time()) - ts) > 60 * 5:
        return False
    base = f"v0:{timestamp}:{raw_body.decode('utf-8', errors='replace')}"
    expected = "v0=" + hmac.new(secret.encode("utf-8"), base.encode("utf-8"), hashlib.sha256).hexdigest()
    return hmac.compare_digest(expected, signature)


def _github_verify_signature(raw_body: bytes, signature_header: str) -> bool:
    secret = os.getenv("GITHUB_WEBHOOK_SECRET", "").strip()
    if not secret:
        return True
    sig = (signature_header or "").strip().lower()
    if not sig.startswith("sha256="):
        return False
    provided = sig.split("=", 1)[1]
    expected = hmac.new(secret.encode("utf-8"), raw_body, hashlib.sha256).hexdigest()
    return hmac.compare_digest(provided, expected)


def _slack_send_message(channel: str, text: str, thread_ts: str = "") -> dict:
    token = os.getenv("SLACK_BOT_TOKEN", "").strip()
    if not token or not channel:
        return {"ok": False, "error": "slack_not_configured"}
    payload = {"channel": channel, "text": text[:3900]}
    if thread_ts:
        payload["thread_ts"] = thread_ts
    resp = _http_json_request(
        "POST",
        "https://slack.com/api/chat.postMessage",
        payload=payload,
        headers={"Authorization": f"Bearer {token}"},
        timeout=15,
    )
    api_ok = bool(resp.get("json", {}).get("ok"))
    return {"ok": api_ok and resp.get("ok", False), "status": resp.get("status", 0), "raw": resp}


def _discord_send_message(channel_id: str, content: str) -> dict:
    token = os.getenv("DISCORD_BOT_TOKEN", "").strip()
    if not token or not channel_id:
        return {"ok": False, "error": "discord_not_configured"}
    resp = _http_json_request(
        "POST",
        f"https://discord.com/api/v10/channels/{channel_id}/messages",
        payload={"content": content[:1800]},
        headers={"Authorization": f"Bot {token}"},
        timeout=15,
    )
    return {"ok": bool(resp.get("ok")), "status": resp.get("status", 0), "raw": resp}


def _extract_result_text(payload: dict) -> str:
    if payload.get("status") != "done":
        return f"Nexus run failed: {payload.get('error') or 'unknown error'}"
    result = payload.get("result")
    if isinstance(result, dict):
        text = result.get("result") or result.get("content") or json.dumps(result, ensure_ascii=False)
        return str(text)[:3900]
    return str(result)[:3900]


@router.post("/integrations/slack/events")
async def integrations_slack_events(request: Request):
    raw_body = await request.body()
    timestamp = request.headers.get("x-slack-request-timestamp", "")
    signature = request.headers.get("x-slack-signature", "")
    if not _slack_verify_signature(raw_body, timestamp, signature):
        return _api_error("invalid slack signature", "unauthorized", 403)

    try:
        body = json.loads(raw_body.decode("utf-8")) if raw_body else {}
    except Exception:
        return _api_error("invalid JSON body", "validation_error", 400)
    if not isinstance(body, dict):
        return _api_error("invalid JSON body", "validation_error", 400)

    if body.get("type") == "url_verification":
        return {"challenge": body.get("challenge", "")}

    event = body.get("event") if isinstance(body.get("event"), dict) else {}
    if body.get("type") != "event_callback" or event.get("type") != "message" or event.get("bot_id"):
        return {"accepted": True, "ignored": True}

    text = str(event.get("text") or "").strip()
    if not text:
        return {"accepted": True, "ignored": True, "reason": "empty_message"}
    channel = str(event.get("channel") or "").strip()
    thread_ts = str(event.get("thread_ts") or event.get("ts") or "").strip()

    def _callback(_run_id: str, result_payload: dict) -> None:
        response_text = _extract_result_text(result_payload)
        _slack_send_message(channel, response_text, thread_ts=thread_ts)

    run_id = _start_integration_run(
        task=text,
        source="slack",
        metadata={"channel": channel, "thread_ts": thread_ts},
        callback=_callback,
    )
    return {"accepted": True, "run_id": run_id}


@router.post("/integrations/discord/messages")
async def integrations_discord_messages(request: Request):
    try:
        body = await _read_json_body(request)
    except HTTPException as exc:
        return _api_error(str(exc.detail), "validation_error", exc.status_code)

    expected_secret = os.getenv("DISCORD_INBOUND_SECRET", "").strip()
    provided_secret = request.headers.get("x-discord-secret", "")
    if expected_secret and not hmac.compare_digest(expected_secret, provided_secret):
        return _api_error("invalid discord secret", "unauthorized", 403)

    if bool(body.get("author_bot", False)):
        return {"accepted": True, "ignored": True}

    text = str(body.get("task") or body.get("content") or "").strip()
    if not text:
        return _api_error("content or task is required", "validation_error", 422)
    channel_id = str(body.get("channel_id") or "").strip()

    def _callback(_run_id: str, result_payload: dict) -> None:
        response_text = _extract_result_text(result_payload)
        _discord_send_message(channel_id, response_text)

    run_id = _start_integration_run(
        task=text,
        source="discord",
        metadata={"channel_id": channel_id, "author_id": str(body.get("author_id") or "")},
        callback=_callback,
    )
    return {"accepted": True, "run_id": run_id}


@router.post("/integrations/github-actions/event")
async def integrations_github_actions_event(request: Request):
    raw_body = await request.body()
    if not _github_verify_signature(raw_body, request.headers.get("x-hub-signature-256", "")):
        return _api_error("invalid github signature", "unauthorized", 403)

    try:
        body = json.loads(raw_body.decode("utf-8")) if raw_body else {}
    except Exception:
        return _api_error("invalid JSON body", "validation_error", 400)
    if not isinstance(body, dict):
        return _api_error("invalid JSON body", "validation_error", 400)

    event_name = str(request.headers.get("x-github-event") or body.get("event") or "").strip().lower()
    if event_name not in {"push", "pull_request"}:
        return {"accepted": False, "ignored": True, "reason": "event_not_supported", "event": event_name}

    override_task = str(body.get("task") or "").strip()
    task = override_task
    if not task:
        repo_name = str((body.get("repository") or {}).get("full_name") or "unknown/repo")
        if event_name == "push":
            ref = str(body.get("ref") or "")
            head_commit = body.get("head_commit") if isinstance(body.get("head_commit"), dict) else {}
            message = str(head_commit.get("message") or "").strip()
            task = f"GitHub push event for {repo_name} on {ref}. Summarize impact and propose next action. Commit message: {message}"
        else:
            pr = body.get("pull_request") if isinstance(body.get("pull_request"), dict) else {}
            title = str(pr.get("title") or "")
            draft = bool(pr.get("draft", False))
            action = str(body.get("action") or "updated")
            task = f"GitHub pull_request event ({action}) for {repo_name}. Title: {title}. Draft={draft}. Provide review guidance."

    run_id = _start_integration_run(
        task=task,
        source="github_actions",
        metadata={"event": event_name},
    )
    return {"accepted": True, "run_id": run_id, "event": event_name}


@router.post("/integrations/automation/webhook")
async def integrations_automation_webhook(request: Request):
    try:
        body = await _read_json_body(request)
    except HTTPException as exc:
        return _api_error(str(exc.detail), "validation_error", exc.status_code)

    expected_secret = os.getenv("AUTOMATION_WEBHOOK_SECRET", "").strip()
    provided_secret = request.headers.get("x-automation-secret", "")
    if expected_secret and not hmac.compare_digest(expected_secret, provided_secret):
        return _api_error("invalid automation secret", "unauthorized", 403)

    task = str(body.get("task") or body.get("prompt") or body.get("text") or "").strip()
    if not task:
        return _api_error("task, prompt, or text is required", "validation_error", 422)

    callback_url = str(body.get("callback_url") or "").strip()
    callback_bearer = str(body.get("callback_bearer") or "").strip()
    if callback_url and not _is_safe_callback_url(callback_url):
        return _api_error("callback_url must be https or localhost http", "validation_error", 422)

    def _callback(run_id: str, result_payload: dict) -> None:
        if not callback_url:
            return
        headers = {}
        if callback_bearer:
            headers["Authorization"] = f"Bearer {callback_bearer}"
        _http_json_request(
            "POST",
            callback_url,
            payload={
                "run_id": run_id,
                "status": result_payload.get("status"),
                "error": result_payload.get("error"),
                "result": result_payload.get("result"),
                "source": "automation_webhook",
            },
            headers=headers,
            timeout=20,
        )

    run_id = _start_integration_run(
        task=task,
        source="automation_webhook",
        metadata={"has_callback": bool(callback_url)},
        callback=_callback if callback_url else None,
    )
    return {"accepted": True, "run_id": run_id}


@router.get("/integrations/channels/status")
def integrations_channels_status():
    return {
        "slack": {
            "inbound_signature": bool(os.getenv("SLACK_SIGNING_SECRET", "").strip()),
            "outbound_bot_token": bool(os.getenv("SLACK_BOT_TOKEN", "").strip()),
        },
        "discord": {
            "inbound_secret": bool(os.getenv("DISCORD_INBOUND_SECRET", "").strip()),
            "outbound_bot_token": bool(os.getenv("DISCORD_BOT_TOKEN", "").strip()),
        },
        "github_actions": {
            "webhook_signature": bool(os.getenv("GITHUB_WEBHOOK_SECRET", "").strip()),
        },
        "automation": {
            "webhook_secret": bool(os.getenv("AUTOMATION_WEBHOOK_SECRET", "").strip()),
        },
    }


def _load_mcp_tools_from_env() -> list[dict]:
    raw = os.getenv("MCP_TOOLS", "").strip()
    if not raw:
        return []
    try:
        parsed = json.loads(raw)
    except Exception:
        return []
    if not isinstance(parsed, list):
        return []
    tools = []
    for item in parsed:
        if not isinstance(item, dict):
            continue
        name = str(item.get("name") or "").strip()
        url = str(item.get("url") or "").strip()
        if not name or not url:
            continue
        method = str(item.get("method") or "POST").strip().upper()
        headers = item.get("headers") if isinstance(item.get("headers"), dict) else {}
        timeout_s = _safe_int(item.get("timeout_s", 15), default=15, min_value=1, max_value=60)
        tools.append(
            {
                "name": name,
                "url": url,
                "method": method if method in {"GET", "POST", "PUT", "PATCH", "DELETE"} else "POST",
                "headers": {str(k): str(v) for k, v in headers.items()},
                "timeout_s": timeout_s,
            }
        )
    return tools


def _call_mcp_tool(name: str, args: dict | None = None) -> dict:
    args = args if isinstance(args, dict) else {}
    tools = _load_mcp_tools_from_env()
    selected = next((t for t in tools if t.get("name") == name), None)
    if not selected:
        return {"ok": False, "error": "mcp_tool_not_found", "available": [t.get("name") for t in tools]}

    method = selected.get("method", "POST")
    url = selected.get("url", "")
    payload = None
    if method == "GET" and args:
        query = _urlparse.urlencode({k: str(v) for k, v in args.items()})
        url = url + ("&" if "?" in url else "?") + query
    elif args:
        payload = {"arguments": args}

    resp = _http_json_request(
        method=method,
        url=url,
        payload=payload,
        headers=selected.get("headers") or {},
        timeout=int(selected.get("timeout_s", 15)),
    )
    return {
        "ok": bool(resp.get("ok")),
        "status": resp.get("status", 0),
        "tool": name,
        "response": resp.get("json") if resp.get("json") else resp.get("text"),
    }


def _mcp_server_allowed_tools() -> set[str]:
    from ..tools_builtin import list_tool_schemas

    raw = os.getenv("MCP_SERVER_ALLOWED_TOOLS", "").strip()
    all_tools = set(list_tool_schemas().keys())
    high_risk = {
        "run_command",
        "write_file",
        "delete_file",
        "commit_push",
        "create_repo",
        "git_checkout",
        "git_pull",
    }
    if raw:
        return {part.strip() for part in raw.split(",") if part.strip()}
    if _is_truthy_env("MCP_SERVER_ALLOW_HIGH_RISK", default=False):
        return all_tools
    return all_tools - high_risk


def _mcp_jsonrpc_success(req_id, result: dict):
    return {"jsonrpc": "2.0", "id": req_id, "result": result}


def _mcp_jsonrpc_error(req_id, code: int, message: str, data=None):
    payload = {"jsonrpc": "2.0", "id": req_id, "error": {"code": code, "message": message}}
    if data is not None:
        payload["error"]["data"] = data
    return payload

@router.get("/health")
def health(): return {"status":"healthy","provider":get_config()["provider"]}

@router.get("/api/system/resources")
def system_resources(): return get_system_resources()

@router.get("/providers")
def providers(): return {"providers":get_providers_list()}


@router.get("/providers/free-diagnostics")
def providers_free_diagnostics():
    return get_free_provider_diagnostics()


@router.get("/providers/health")
def providers_health():
    """Return per-provider health status including rate limits, capabilities, and benchmarks."""
    return get_provider_health()


@router.get("/providers/status")
def providers_status():
    """Alias for /providers/health - same as /providers but with detailed health info."""
    return get_provider_health()


# ── Swarm View ────────────────────────────────────────────────────────────────

@router.get("/swarm/activity")
def swarm_activity(limit: int = 50, action: str = "", since_ts: float = 0.0):
    """Return the most recent swarm activity events (capped at _MAX_ACTIVITY)."""
    limit = max(1, min(limit, _MAX_ACTIVITY))
    action_filter = str(action or "").strip()
    if action_filter == "strict_clone_bootstrap_exception":
        events = db_list_strict_clone_bypass_events(limit=limit, since_ts=since_ts)
        return {
            "events": list(reversed(events)),
            "total": db_count_strict_clone_bypass_events(since_ts=since_ts),
            "unfiltered_total": db_count_strict_clone_bypass_events(),
            "source": "persistent",
            "filters": {
                "action": action_filter,
                "since_ts": since_ts if since_ts > 0 else None,
            },
        }
    events = activity_log
    if action_filter:
        events = [event for event in events if str(event.get("action") or "") == action_filter]
    if since_ts > 0:
        events = [event for event in events if float(event.get("ts") or 0.0) >= since_ts]
    return {
        "events": events[-limit:],
        "total": len(events),
        "unfiltered_total": len(activity_log),
        "source": "activity_log",
        "filters": {
            "action": action_filter or None,
            "since_ts": since_ts if since_ts > 0 else None,
        },
    }


@router.get("/swarm/live")
async def swarm_live_stream():
    """SSE stream of swarm activity events as they are pushed to activity_log."""
    import asyncio as _aio
    import json as _json

    async def _generate():
        last_seen = len(activity_log)
        # Send existing events first (catch-up)
        for ev in activity_log[-30:]:
            yield f"data: {_json.dumps(ev)}\n\n"
        # Then stream new ones
        while True:
            await _aio.sleep(0.5)
            current = len(activity_log)
            if current > last_seen:
                for ev in activity_log[last_seen:current]:
                    yield f"data: {_json.dumps(ev)}\n\n"
                last_seen = current

    return StreamingResponse(
        _generate(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )


# ── Scheduler API ─────────────────────────────────────────────────────────────

@router.get("/scheduler/jobs")
def scheduler_jobs():
    jobs = [job_to_dict(j) for j in list_jobs()]
    return {"jobs": jobs, "total": len(jobs)}


@router.post("/scheduler/jobs")
async def scheduler_create_job(request: Request):
    try:
        body = await _read_json_body(request)
    except HTTPException as exc:
        return _api_error(str(exc.detail), "validation_error", exc.status_code)
    name = (body.get("name") or "background-task").strip()
    task = (body.get("task") or "").strip()
    schedule = (body.get("schedule") or "5m").strip()
    if not task:
        return _api_error("task is required", "validation_error", 422)
    try:
        safe_task = check_user_task(task)
    except GuardrailViolation as exc:
        return JSONResponse(
            {
                "error": exc.reason,
                "type": exc.code,
                "safety": {"blocked": True, "code": exc.code, "message": exc.reason},
            },
            status_code=422,
        )
    try:
        job = schedule_job(
            name=name,
            task=safe_task,
            schedule=schedule,
            max_retries=int(body.get("max_retries", 0)),
            retry_backoff_secs=int(body.get("retry_backoff_secs", 60)),
        )
        return {"job": job_to_dict(job)}
    except Exception as exc:
        return _api_error(f"Failed to create job: {exc}", "validation_error", 422)


@router.post("/scheduler/jobs/{job_id}/cancel")
def scheduler_cancel_job(job_id: str):
    if cancel_job(job_id):
        return {"ok": True, "job_id": job_id}
    return _api_error("job not found", "not_found", 404)


@router.get("/scheduler/jobs/{job_id}/history")
def scheduler_job_history(job_id: str, limit: int = 50):
    """Return past execution records for a scheduled job."""
    job = next((j for j in list_jobs() if j.id == job_id), None)
    if job is None:
        return _api_error("job not found", "not_found", 404)
    history = getattr(job, "history", [])
    return {
        "job_id": job_id,
        "job_name": getattr(job, "name", job_id),
        "history": list(history[-max(1, min(int(limit), 500)):]),
        "total": len(history),
    }


@router.post("/scheduler/webhook/{job_id}")
async def scheduler_webhook_trigger(job_id: str, request: Request):
    """Immediately trigger a scheduled job via webhook."""
    try:
        body = await _read_json_body(request)
    except HTTPException:
        body = {}
    job = next((j for j in list_jobs() if j.id == job_id), None)
    if job is None:
        return _api_error("job not found", "not_found", 404)
    try:
        from ..scheduler import run_job_now
        result = run_job_now(job_id)
        return {"ok": True, "job_id": job_id, "result": str(result)[:200]}
    except Exception as exc:
        return _api_error(str(exc), "server_error", 500)


@router.post("/settings/domain-guards")
async def safety_domain_guards_update(request: Request):
    """Update domain-guard policy rules."""
    try:
        body = await _read_json_body(request)
    except HTTPException as exc:
        return _api_error(str(exc.detail), "validation_error", exc.status_code)
    if not isinstance(body, dict):
        return _api_error("body must be a JSON object", "validation_error", 422)
    # Validate allowed keys only
    valid_keys = {"blocked_domains", "allowed_categories", "block_adult", "block_malware"}
    unknown = set(body.keys()) - valid_keys
    if unknown:
        return _api_error(f"Unknown fields: {sorted(unknown)}", "validation_error", 422)
    existing = _config.get("domain_guards") or {}
    existing.update(body)
    _config["domain_guards"] = existing
    return {"ok": True, "domain_guards": existing}


@router.post("/generation/video")
async def generation_video(request: Request):
    """Generate local video bytes and return base64 payload."""
    try:
        body = await _read_json_body(request)
    except HTTPException as exc:
        return _api_error(str(exc.detail), "validation_error", exc.status_code)

    prompt = (body.get("prompt") or "").strip()
    if not prompt:
        return _api_error("prompt is required", "validation_error", 422)

    try:
        from ..generation import generate_video

        backend_name = str(body.get("backend") or "auto")
        video_bytes = generate_video(
            prompt=prompt,
            duration_seconds=float(body.get("duration_seconds") or 4.0),
            fps=int(body.get("fps") or 8),
            width=int(body.get("width") or 512),
            height=int(body.get("height") or 512),
            backend=backend_name,
        )
    except ValueError as exc:
        return _api_error(str(exc), "validation_error", 422)
    except Exception as exc:
        return _api_error(str(exc), "generation_error", 500)

    return {
        "mime_type": "video/mp4",
        "backend": backend_name,
        "model": "sd3" if "sd3" in backend_name else "flux",
        "duration_seconds": float(body.get("duration_seconds") or 4.0),
        "video_b64": base64.b64encode(video_bytes).decode("ascii"),
    }


@router.post("/generation/video/stream")
async def generation_video_stream(request: Request):
    """Generate video with realtime SSE progress events and final payload."""
    try:
        body = await _read_json_body(request)
    except HTTPException as exc:
        return _api_error(str(exc.detail), "validation_error", exc.status_code)

    prompt = (body.get("prompt") or "").strip()
    if not prompt:
        return _api_error("prompt is required", "validation_error", 422)

    duration_seconds = float(body.get("duration_seconds") or 4.0)
    fps = int(body.get("fps") or 8)
    width = int(body.get("width") or 512)
    height = int(body.get("height") or 512)
    backend_name = str(body.get("backend") or "auto")

    from ..generation import generate_video_stream as _generate_video_stream

    async def _gen():
        for event in _generate_video_stream(
            prompt=prompt,
            duration_seconds=duration_seconds,
            fps=fps,
            width=width,
            height=height,
            backend=backend_name,
            include_frame_payload=False,
        ):
            payload = dict(event)
            if payload.get("type") == "done" and payload.get("video_bytes"):
                payload["video_b64"] = base64.b64encode(bytes(payload.pop("video_bytes"))).decode("ascii")
            yield f"data: {json.dumps(payload)}\n\n"
        yield "data: [DONE]\n\n"

    return StreamingResponse(
        _gen(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )


@router.post("/vision/understand")
async def vision_understand(request: Request):
    """Describe an image using local vision flow."""
    try:
        body = await _read_json_body(request)
    except HTTPException as exc:
        return _api_error(str(exc.detail), "validation_error", exc.status_code)

    image_b64 = str(body.get("image_b64") or "").strip()
    image_url = str(body.get("image_url") or "").strip()
    prompt = str(body.get("prompt") or "Describe this image in detail.")
    mime_type = str(body.get("mime_type") or "image/png")

    if not image_b64 and not image_url:
        return _api_error("image_b64 or image_url is required", "validation_error", 422)

    try:
        from ..vision import describe_image, capture_screenshot

        if image_b64:
            image_bytes = base64.b64decode(image_b64)
        else:
            image_bytes = capture_screenshot(image_url)
            mime_type = "image/png"

        description = describe_image(image_bytes=image_bytes, mime_type=mime_type, prompt=prompt)
    except ValueError as exc:
        return _api_error(str(exc), "validation_error", 422)
    except Exception as exc:
        return _api_error(str(exc), "vision_error", 500)

    return {
        "description": description,
        "mime_type": mime_type,
        "prompt": prompt,
    }


@router.post("/integrations/tunnel")
async def integrations_tunnel(request: Request):
    try:
        body = await _read_json_body(request)
    except HTTPException as exc:
        return _api_error(str(exc.detail), "validation_error", exc.status_code)

    try:
        from ..integrations import NexusTunnelConfig, connect_tunnel, get_tunnel_status
        config = NexusTunnelConfig(
            endpoint=str(body.get("endpoint") or "").strip(),
            auth_token=str(body.get("auth_token") or "").strip(),
            local_port=int(body.get("local_port") or 8000),
            subdomain=str(body.get("subdomain") or "").strip(),
            auto_reconnect=bool(body.get("auto_reconnect", True)),
        )
        public_url = connect_tunnel(config)
        return {"public_url": public_url, "status": get_tunnel_status()}
    except ValueError as exc:
        return _api_error(str(exc), "validation_error", 422)


@router.post("/integrations/guardian")
async def integrations_guardian(request: Request):
    try:
        body = await _read_json_body(request)
    except HTTPException as exc:
        return _api_error(str(exc.detail), "validation_error", exc.status_code)

    try:
        from ..integrations import GuardianConfig, get_guardian_status, register_with_guardian
        config = GuardianConfig(
            endpoint=str(body.get("endpoint") or "").strip(),
            api_key=str(body.get("api_key") or "").strip(),
            organisation_id=str(body.get("organisation_id") or "").strip(),
            enforce_policies=bool(body.get("enforce_policies", True)),
        )
        instance_id = register_with_guardian(config)
        return {"instance_id": instance_id, "status": get_guardian_status()}
    except ValueError as exc:
        return _api_error(str(exc), "validation_error", 422)


@router.post("/integrations/edge")
async def integrations_edge(request: Request):
    try:
        body = await _read_json_body(request)
    except HTTPException as exc:
        return _api_error(str(exc.detail), "validation_error", exc.status_code)

    try:
        from ..integrations import EdgeNodeConfig, get_edge_status, register_edge_node
        config = EdgeNodeConfig(
            node_id=str(body.get("node_id") or "").strip(),
            orchestrator_url=str(body.get("orchestrator_url") or "").strip(),
            model_ids=list(body.get("model_ids") or []),
            heartbeat_interval_s=int(body.get("heartbeat_interval_s") or 30),
            max_concurrent_requests=int(body.get("max_concurrent_requests") or 4),
            api_key=str(body.get("api_key") or "").strip(),
        )
        node_id = register_edge_node(config)
        return {"node_id": node_id, "status": get_edge_status()}
    except ValueError as exc:
        return _api_error(str(exc), "validation_error", 422)





@router.post("/benchmark/eval-suite")
async def benchmark_eval_suite(request: Request):
    try:
        body = await _read_json_body(request)
    except HTTPException as exc:
        return _api_error(str(exc.detail), "validation_error", exc.status_code)

    model = str(body.get("model") or "").strip()
    if not model:
        return _api_error("model is required", "validation_error", 422)

    try:
        from ..eval_pipeline import run_eval_suite
        result = run_eval_suite(
            model=model,
            provider=str(body.get("provider") or "ollama"),
            suites=list(body.get("suites") or []),
            n_samples=int(body.get("n_samples") or 20),
            adapter_id=body.get("adapter_id"),
        )
        jobs = []
        for job in result.get("jobs") or []:
            if isinstance(job, dict):
                jobs.append(job)
            else:
                jobs.append(getattr(job, "__dict__", {"value": str(job)}))
        return {
            **{k: v for k, v in result.items() if k != "jobs"},
            "jobs": jobs,
        }
    except ValueError as exc:
        return _api_error(str(exc), "validation_error", 422)


@router.post("/benchmark/regression")
async def benchmark_regression(request: Request):
    try:
        body = await _read_json_body(request)
    except HTTPException as exc:
        return _api_error(str(exc.detail), "validation_error", exc.status_code)

    model = str(body.get("model") or "").strip()
    if not model:
        return _api_error("model is required", "validation_error", 422)

    try:
        from ..eval_pipeline import run_regression_benchmark
        return run_regression_benchmark(
            model=model,
            provider=str(body.get("provider") or "ollama"),
            suites=list(body.get("suites") or []),
            threshold=float(body.get("threshold") or 0.05),
            n_samples=int(body.get("n_samples") or 20),
        )
    except ValueError as exc:
        return _api_error(str(exc), "validation_error", 422)


@router.post("/validation/program/run")
async def validation_program_run(request: Request):
    try:
        body = await _read_json_body(request)
    except HTTPException as exc:
        return _api_error(str(exc.detail), "validation_error", exc.status_code)

    try:
        from ..validation_program import run_validation_program

        return run_validation_program(
            domains=body.get("domains") if isinstance(body.get("domains"), list) else None,
            update_baseline=bool(body.get("update_baseline", False)),
            alert_on_regression=bool(body.get("alert_on_regression", True)),
        )
    except ValueError as exc:
        return _api_error(str(exc), "validation_error", 422)


@router.get("/validation/program/reports")
def validation_program_reports(program: str = "core_quality", limit: int = 20):
    from ..validation_program import get_validation_overview

    return get_validation_overview(program=program, limit=limit)


@router.get("/validation/program/baselines")
def validation_program_baselines():
    from ..db import load_validation_baselines

    return {"baselines": load_validation_baselines()}


@router.post("/validation/program/baselines")
async def validation_program_update_baselines(request: Request):
    body = await _read_json_body(request, "invalid JSON body")
    from ..db import save_validation_baselines
    from ..validation_program import run_validation_program

    if isinstance(body.get("baselines"), dict):
        saved = save_validation_baselines(body.get("baselines") or {})
        return {"ok": True, "baselines": saved, "mode": "direct"}

    report = run_validation_program(
        domains=body.get("domains") if isinstance(body.get("domains"), list) else None,
        update_baseline=True,
        alert_on_regression=bool(body.get("alert_on_regression", False)),
    )
    return {"ok": True, "mode": "recomputed", "report": report, "baselines": report.get("baselines", {})}


@router.post("/compute/contributed/register")
async def contributed_register(request: Request):
    try:
        body = await _read_json_body(request)
    except HTTPException as exc:
        return _api_error(str(exc.detail), "validation_error", exc.status_code)

    try:
        from ..hardware import register_contributed_compute
        registration = register_contributed_compute(
            endpoint=str(body.get("endpoint") or "").strip(),
            api_key=str(body.get("api_key") or "").strip(),
            max_concurrent=int(body.get("max_concurrent") or 1),
        )
        return registration
    except ValueError as exc:
        return _api_error(str(exc), "validation_error", 422)


@router.get("/compute/contributed/nodes")
def contributed_nodes():
    from ..hardware import list_contributed_compute_nodes
    nodes = list_contributed_compute_nodes()
    return {"nodes": nodes, "total": len(nodes)}


@router.delete("/compute/contributed/{node_id}")
def contributed_deregister(node_id: str):
    from ..hardware import deregister_contributed_compute
    removed = deregister_contributed_compute(node_id)
    if not removed:
        return _api_error("node not found", "not_found", 404)
    return {"ok": True, "node_id": node_id}


# ── OpenAI-compatible API (v1) ────────────────────────────────────────────────
# Allows Nexusclaw, Nexus Computer, and any OpenAI-compatible client to use
# Nexus AI as a drop-in API engine.  Set base_url to http://<host>:<port>/v1.
#
# Endpoints:
#   GET  /v1/models                  – list available models
#   GET  /v1/models/{model}          – retrieve model metadata
#   POST /v1/chat/completions        – synchronous or streaming chat


def _v1_models_catalog() -> list[dict]:
    providers = get_providers_list()
    provider_models = []
    for provider in providers:
        if isinstance(provider, dict):
            provider_id = str(provider.get("id", "")).strip()
        else:
            provider_id = str(provider).strip()
        if not provider_id:
            continue
        provider_models.append(
            {
                "id": f"nexus-ai/{provider_id}",
                "object": "model",
                "created": 0,
                "owned_by": "nexus-systems",
            }
        )

    return [
        {"id": "nexus-ai", "object": "model", "created": 0, "owned_by": "nexus-systems"},
        {"id": "nexus-ai/auto", "object": "model", "created": 0, "owned_by": "nexus-systems"},
    ] + provider_models


def _normalize_embeddings_input(raw_input):
    if isinstance(raw_input, str):
        text = raw_input
        return [text], max(1, len(text.split()))

    if not isinstance(raw_input, list) or not raw_input:
        raise ValueError("input is required")

    if all(isinstance(item, str) for item in raw_input):
        texts = raw_input
        token_count = sum(max(1, len(item.split())) for item in texts)
        return texts, token_count

    if all(isinstance(item, int) for item in raw_input):
        token_ids = raw_input
        return [" ".join(str(token) for token in token_ids)], max(1, len(token_ids))

    if all(isinstance(item, list) and all(isinstance(token, int) for token in item) for item in raw_input):
        token_batches = raw_input
        texts = [" ".join(str(token) for token in token_batch) for token_batch in token_batches]
        token_count = sum(max(1, len(token_batch)) for token_batch in token_batches)
        return texts, token_count

    raise ValueError("input must be a string, list of strings, token array, or list of token arrays")


def _estimate_text_tokens(text: str) -> int:
    normalized = str(text or "").strip()
    if not normalized:
        return 0
    return len(normalized.split())


def _resolve_provider_order_from_model(model: str) -> list[str] | None:
    raw = str(model or "").strip()
    if not raw:
        return None
    normalized = raw
    if normalized.startswith("nexus-ai/"):
        normalized = normalized.split("/", 1)[1].strip()
    elif normalized == "nexus-ai":
        normalized = "auto"

    if normalized in {"", "auto"}:
        return None

    available = {str(p.get("id") or "").strip() for p in get_providers_list() if isinstance(p, dict)}
    if normalized in available:
        return [normalized]
    return None


# ── settings ──────────────────────────────────────────────────────────────────
@router.get("/manifest.json")
def manifest():
    return FileResponse("static/manifest.json", media_type="application/manifest+json")

@router.get("/sw.js")
def service_worker():
    return FileResponse("static/sw.js", media_type="application/javascript",
                        headers={"Service-Worker-Allowed": "/"})

@router.get("/personas")
def get_personas():
    return {"personas": list_personas(), "active": get_active_persona_name()}

@router.post("/personas/{persona_id}")
async def switch_persona(persona_id: str, request: Request):
    # Route order means POST /personas/custom can match this path first.
    # Support custom persona creation here to preserve the public contract.
    if persona_id == "custom":
        data = await request.json()
        pid = data.get("id") or str(uuid.uuid4())
        db_save_persona(
            pid,
            data.get("name", "Custom"),
            data.get("icon", "🤖"),
            data.get("description", ""),
            data.get("prompt_prefix", ""),
            data.get("color", "#7c6af7"),
            float(data.get("temperature", 0.2)),
            data.get("tier", "medium"),
        )
        return {"id": pid}

    p = set_persona(persona_id)
    return {"active": persona_id, "persona": p}

@router.get("/settings")
def get_settings(): return get_config()

@router.post("/settings")
async def post_settings(request: Request):
    try:
        data = await _read_json_body(request)
    except HTTPException as exc:
        return _api_error(str(exc.detail), "validation_error", exc.status_code)
    prev_profile = _config.get("safety_profile", "standard")
    result = update_config(provider=data.get("provider"),
                           model=data.get("model"),
                           temperature=data.get("temperature"),
                           persona=data.get("persona"),
                           safety_profile=data.get("safety_profile"),
                           strict_mode_profile=data.get("strict_mode_profile"),
                           strict_no_guess_mode=data.get("strict_no_guess_mode"),
                           strict_confidence_threshold=data.get("strict_confidence_threshold"),
                           strict_evidence_threshold=data.get("strict_evidence_threshold"))
    new_profile = _config.get("safety_profile", "standard")
    if data.get("safety_profile") and new_profile != prev_profile:
        _push_safety_event("profile_change", {"scope": "global", "from": prev_profile, "to": new_profile})
    return result


@router.get("/settings/safety")
def get_safety_settings():
    profile = _config.get("safety_profile", "standard")
    return {
        "safety_profile": profile,
        "policy": get_safety_policy(profile),
        "available_profiles": sorted(SAFETY_POLICY_PROFILES.keys()),
    }


@router.post("/settings/safety")
async def update_safety_settings(request: Request):
    try:
        data = await _read_json_body(request)
    except HTTPException as exc:
        return _api_error(str(exc.detail), "validation_error", exc.status_code)
    profile = str(data.get("safety_profile", "standard")).lower().strip()
    if profile not in SAFETY_POLICY_PROFILES:
        allowed = ", ".join(sorted(SAFETY_POLICY_PROFILES.keys()))
        return _api_error(f"safety_profile must be one of: {allowed}", "validation_error", 422)
    prev = _config.get("safety_profile", "standard")
    update_config(safety_profile=profile)
    if profile != prev:
        _push_safety_event("profile_change", {"scope": "global", "from": prev, "to": profile})
    return {
        "safety_profile": _config.get("safety_profile", "standard"),
        "policy": get_safety_policy(_config.get("safety_profile", "standard")),
        "available_profiles": sorted(SAFETY_POLICY_PROFILES.keys()),
    }


@router.get("/settings/rate-limits")
def get_rate_limit_settings():
    return dict(_rate_limit_settings)


@router.post("/settings/rate-limits")
async def update_rate_limit_settings(request: Request):
    try:
        data = await _read_json_body(request)
    except HTTPException as exc:
        return _api_error(str(exc.detail), "validation_error", exc.status_code)
    mode = str(data.get("mode", _rate_limit_settings.get("mode", "soft"))).lower().strip()
    per_minute = data.get("per_minute", _rate_limit_settings.get("per_minute", 60))
    per_day = data.get("per_day", _rate_limit_settings.get("per_day", 2500))

    if mode not in ("soft", "hard"):
        return _api_error("mode must be one of: soft, hard", "validation_error", 422)
    try:
        per_minute = int(per_minute)
        per_day = int(per_day)
    except Exception:
        return _api_error("per_minute and per_day must be integers", "validation_error", 422)
    if per_minute < 1 or per_day < 1:
        return _api_error("per_minute and per_day must be >= 1", "validation_error", 422)

    updated = {
        "mode": mode,
        "per_minute": min(per_minute, 100000),
        "per_day": min(per_day, 10000000),
    }
    _rate_limit_settings.update(updated)
    db_save_pref("rate_limit_settings", json.dumps(updated))
    return dict(_rate_limit_settings)


_SEVERITY_ORDER = {
    "none": 0,
    "low": 1,
    "medium": 2,
    "high": 3,
    "critical": 4,
}


def _event_severity(event: dict) -> str:
    verdict = event.get("verdict") or {}
    issue_levels = []
    for issue in verdict.get("issues", []) or []:
        level = str(issue.get("severity") or issue.get("threat") or "").lower().strip()
        if level in _SEVERITY_ORDER:
            issue_levels.append(level)

    if issue_levels:
        return max(issue_levels, key=lambda lvl: _SEVERITY_ORDER.get(lvl, 0))

    event_type = str(event.get("type") or "")
    if event_type == "block":
        return "high"
    if event_type == "pii_scrub":
        return "medium"
    if event_type == "profile_change":
        return "low"
    return "none"

@router.get("/personas/legacy")
def list_personas_legacy():
    active = _config["persona"]
    return {"personas": [
        {"id": k, "label": v["label"], "emoji": v["emoji"],
         "description": v["description"], "active": k == active}
        for k, v in PERSONAS.items()
    ]}


# ── Benchmark endpoints ────────────────────────────────────────────────────────
_BENCHMARK_PROBES = [
    ("arithmetic",  "What is 17 * 23?"),
    ("reasoning",   "If all roses are flowers and some flowers fade quickly, can we conclude that some roses fade quickly?"),
    ("coding",      "Write a one-line Python expression to reverse a string."),
]

@router.post("/benchmark/run")
async def benchmark_run(request: Request):
    """Run a lightweight probe suite against all available providers and store results."""
    from ..benchmark import run_benchmark_probes
    body = {}
    try:
        body = await request.json()
    except Exception:
        pass
    return run_benchmark_probes(requested_providers=body.get("providers") or [])


@router.get("/benchmark/results")
def benchmark_results():
    """Return stored benchmark results (most recent first)."""
    from ..db import load_benchmark_results
    return {"results": load_benchmark_results()}


@router.get("/benchmark/history")
def benchmark_history(provider: str = "", model: str = "", task_type: str = "", limit: int = 500):
    """Return benchmark history with optional provider/model/task filters."""
    from ..benchmark import get_benchmark_history
    return get_benchmark_history(provider=provider, model=model, task_type=task_type, limit=limit)


@router.get("/benchmark/leaderboard")
def benchmark_leaderboard(sort_by: str = "task_type", limit: int = 20, provider: str = ""):
    """Return ranked benchmark leaderboard entries."""
    from dataclasses import asdict
    from ..leaderboard import get_leaderboard

    entries = [asdict(item) for item in get_leaderboard(limit=limit)]
    if provider:
        entries = [row for row in entries if str(row.get("provider", "")) == provider]
    return {"entries": entries, "sort_by": sort_by}


@router.get("/benchmark/tradeoff")
def benchmark_tradeoff(days: int = 14, limit: int = 2000):
    """Return per-model cost-quality-latency tradeoff aggregates for dashboard visualization."""
    from ..benchmark import get_benchmark_tradeoff
    return get_benchmark_tradeoff(days=days, limit=limit)


def _extract_citation_urls(text: str) -> list[str]:
    import re as _re

    source = text or ""
    urls = []

    # Markdown links: [label](https://...)
    for m in _re.finditer(r"\[[^\]]+\]\((https?://[^)\s]+)\)", source):
        urls.append(m.group(1))

    # Bare URLs
    for m in _re.finditer(r"\bhttps?://[^\s)]+", source):
        urls.append(m.group(0))

    # Preserve order and uniqueness
    unique = []
    seen = set()
    for u in urls:
        if u not in seen:
            seen.add(u)
            unique.append(u)
    return unique


def _score_citation_confidence(answer: str, expected_sources: list[str] | None = None) -> dict:
    from urllib.parse import urlparse as _urlparse

    citations = _extract_citation_urls(answer)
    expected = [str(s).strip() for s in (expected_sources or []) if str(s).strip()]

    if not citations:
        return {
            "score": 0.1 if expected else 0.25,
            "citations": [],
            "matched_expected_sources": [],
            "expected_source_coverage": 0.0,
        }

    if not expected:
        score = min(0.9, 0.35 + 0.12 * len(citations))
        return {
            "score": round(score, 3),
            "citations": citations,
            "matched_expected_sources": [],
            "expected_source_coverage": None,
        }

    expected_domains = {(_urlparse(u).netloc or u).lower() for u in expected}
    citation_domains = [(_urlparse(u).netloc or u).lower() for u in citations]

    matched = []
    for domain in expected_domains:
        if any(domain in cd or cd in domain for cd in citation_domains):
            matched.append(domain)

    coverage = len(matched) / max(1, len(expected_domains))
    score = 0.35 + 0.65 * coverage
    return {
        "score": round(min(1.0, score), 3),
        "citations": citations,
        "matched_expected_sources": matched,
        "expected_source_coverage": round(coverage, 3),
    }


# ── Adaptive confidence routing helper ────────────────────────────────────

_ADAPTIVE_ROUTING_DEFAULTS: Dict[str, Any] = {
    "enabled": True,
    "confidence_threshold": 0.6,
    "escalation_tries": 2,
}
_adaptive_routing_config: Dict[str, Any] = dict(_ADAPTIVE_ROUTING_DEFAULTS)

# High-quality provider tiers used for escalation — ordered strongest-first.
_ESCALATION_PROVIDERS = ["claude", "openai", "groq", "cerebras", "gemini", "mistral"]


def _call_llm_adaptive(messages: List[Dict], task: str = "") -> tuple:
    """Call LLM with confidence-aware adaptive escalation.

    If the first response has confidence below the configured threshold,
    retry up to ``escalation_tries`` times, preferring stronger providers.

    Returns:
        (result_dict, provider_id, escalated: bool, final_confidence: float)
    """
    cfg = _adaptive_routing_config
    threshold = float(cfg.get("confidence_threshold", 0.6))
    tries = int(cfg.get("escalation_tries", 2))
    enabled = bool(cfg.get("enabled", True))

    result, provider = call_llm_with_fallback(messages, task)

    # Extract confidence from result if present
    confidence = 1.0
    content = result.get("content", "")
    if isinstance(content, str):
        try:
            parsed = json.loads(content)
            confidence = float(parsed.get("confidence", 1.0))
        except Exception:
            confidence = 1.0

    if not enabled or confidence >= threshold or tries <= 0:
        return result, provider, False, confidence

    # Escalate to stronger providers
    escalated = False
    for attempt in range(tries):
        # Try escalation providers not already used
        escalation_order = [p for p in _ESCALATION_PROVIDERS if p != provider]
        if not escalation_order:
            break
        try:
            better_result, better_provider = call_llm_with_fallback(
                messages,
                task,
            )
            better_content = better_result.get("content", "")
            better_confidence = 1.0
            if isinstance(better_content, str):
                try:
                    parsed2 = json.loads(better_content)
                    better_confidence = float(parsed2.get("confidence", 1.0))
                except Exception:
                    better_confidence = 1.0

            if better_confidence > confidence:
                result, provider, confidence = better_result, better_provider, better_confidence
                escalated = True

            if confidence >= threshold:
                break
        except Exception:
            break

    return result, provider, escalated, confidence


def _rag_retrieval_confidence(results: list[dict]) -> float:
    if not results:
        return 0.0
    scores = []
    for r in results:
        try:
            score = float(r.get("score", 0.0))
        except Exception:
            score = 0.0
        scores.append(max(0.0, min(1.0, score)))
    return round(sum(scores) / max(len(scores), 1), 4)


def _build_rag_citations(results: list[dict]) -> list[dict]:
    citations: list[dict] = []
    for idx, item in enumerate(results, start=1):
        meta = item.get("metadata", {}) if isinstance(item.get("metadata"), dict) else {}
        source = meta.get("source_url") or meta.get("source") or f"document-{idx}"
        chunk_ref = meta.get("chunk_index", meta.get("section", idx - 1))
        citations.append(
            {
                "rank": idx,
                "source": source,
                "chunk_ref": chunk_ref,
                "score": float(item.get("score", 0.0) or 0.0),
                "id": item.get("id", ""),
            }
        )
    return citations


def _rag_answer_with_critic(query: str, results: list[dict]) -> dict:
    from ..rag.critic import CriticAgent

    if not results:
        return {
            "answer": "No relevant documents were found for this query.",
            "model_confidence": 0.0,
            "retrieval_confidence": 0.0,
            "calibrated_confidence": 0.0,
            "critique": None,
        }

    context = "\n\n".join(
        f"[{i + 1}] {r.get('document', '')}"
        for i, r in enumerate(results[:5])
    )

    prompt = (
        "Answer the question using only the provided retrieval context. "
        "Cite sources inline as [1], [2], etc. If uncertain, explicitly say so.\n\n"
        f"Question: {query}\n\nContext:\n{context}"
    )

    answer_text = ""
    model_conf = 0.55
    try:
        llm_resp, _provider = call_llm_with_fallback(
            [{"role": "user", "content": prompt}],
            "rag_query",
        )
        answer_text = (llm_resp.get("content") if isinstance(llm_resp, dict) else str(llm_resp) or "").strip()
    except Exception:
        answer_text = " ".join(str(r.get("document", "")).strip() for r in results[:3])[:1200]

    retrieval_conf = _rag_retrieval_confidence(results)
    critic = CriticAgent()
    critique = critic.critique(query, answer_text, results)
    model_conf = round(max(0.0, min(1.0, float(critique.overall_score))), 4)
    calibrated = round((0.45 * retrieval_conf) + (0.55 * model_conf), 4)

    return {
        "answer": answer_text,
        "model_confidence": model_conf,
        "retrieval_confidence": retrieval_conf,
        "calibrated_confidence": calibrated,
        "critique": critique.to_dict(),
    }


# ── Diff viewer endpoints ─────────────────────────────────────────────────────

def _compute_diff_stats(original: str, modified: str, filename: str = "file") -> dict:
    """Compute structured diff stats between two text blobs."""
    import difflib
    orig_lines = original.splitlines(keepends=True)
    mod_lines  = modified.splitlines(keepends=True)
    diff = list(difflib.unified_diff(
        orig_lines, mod_lines,
        fromfile=f"a/{filename}", tofile=f"b/{filename}", lineterm="",
    ))
    unified = "\n".join(diff)
    additions = sum(1 for l in diff if l.startswith("+") and not l.startswith("+++"))
    deletions = sum(1 for l in diff if l.startswith("-") and not l.startswith("---"))
    chunks = sum(1 for l in diff if l.startswith("@@"))
    return {
        "filename": filename,
        "additions": additions,
        "deletions": deletions,
        "chunks": chunks,
        "unchanged": max(0, len(orig_lines) - deletions),
        "unified_diff": unified if len(unified) <= 8000 else unified[:8000] + "\n… (truncated)",
        "has_changes": bool(diff),
    }


@router.post("/diff")
async def compute_diff(request: Request):
    """Compute a structured unified diff between two text blobs.

    POST body:
      { "original": "...", "modified": "...", "filename": "app.py", "trace_id": "" }

    If trace_id is provided the diff is persisted to the file diff history.
    """
    body = {}
    try:
        body = await request.json()
    except Exception:
        pass
    original = body.get("original") or ""
    modified = body.get("modified") or ""
    filename = (body.get("filename") or "file").strip()
    trace_id = (body.get("trace_id") or "").strip()

    stats = _compute_diff_stats(original, modified, filename)

    saved = None
    if trace_id and stats["has_changes"]:
        saved = _save_file_diff(trace_id, filename, original, modified)

    return {**stats, "saved": saved}


@router.get("/diff/history")
def diff_history(trace_id: str = "", limit: int = 50):
    """Return file diff history, optionally filtered by trace_id."""
    try:
        limit = max(1, min(int(limit), 200))
    except Exception:
        limit = 50
    diffs = _get_file_diffs(trace_id=trace_id, limit=limit)
    return {"diffs": diffs, "total": len(diffs)}


@router.get("/diff/{diff_id}")
def diff_detail(diff_id: str):
    """Return full diff detail (including before/after text and unified diff) by id."""
    record = _get_file_diff_detail(diff_id)
    if not record:
        return JSONResponse({"error": "diff not found"}, status_code=404)
    return record


# ── Self-improvement loop endpoints ──────────────────────────────────────────

def _build_self_review_prompt(traces: list[dict]) -> str:
    """Build an LLM prompt for analyzing agent traces and suggesting improvements."""
    if not traces:
        return "No execution traces available to review."
    lines = [
        "You are an AI self-improvement analyst. The following are summaries of recent agent "
        "execution traces. Analyze them and respond with a JSON object:\n"
        '{\n'
        '  "insights": ["<insight 1>", ...],\n'
        '  "suggestions": ["<improvement suggestion 1>", ...]\n'
        '}\n\n'
        "Include 3-7 insights (patterns you observed) and 3-5 actionable suggestions for improving "
        "agent prompts, tool selection, or task handling. Be specific.\n\n"
        "Traces:\n"
    ]
    for t in traces[:10]:
        lines.append(
            f"- trace_id={t.get('trace_id','?')} steps={t.get('steps',0)} "
            f"task={str(t.get('task',''))[:120]} started={t.get('started_at','?')[:19]}"
        )
    return "\n".join(lines)


# ── Document understanding endpoints ─────────────────────────────────────────

def _extract_document_text(path: str, file_type: str = "") -> tuple[str, str]:
    """Extract text from a document file. Returns (text, detected_type)."""
    from ..tools_builtin import tool_read_pdf, tool_read_docx, tool_read_xlsx, tool_read_pptx

    ext = file_type.lower().lstrip(".") if file_type else ""
    if not ext and path:
        ext = path.rsplit(".", 1)[-1].lower() if "." in path else ""

    workdir = os.path.dirname(path) if os.path.isabs(path) else "/tmp"

    if ext in ("pdf",):
        return tool_read_pdf(path, workdir=workdir), "pdf"
    elif ext in ("docx", "doc"):
        return tool_read_docx(path, workdir=workdir), "docx"
    elif ext in ("xlsx", "xls"):
        return tool_read_xlsx(path, workdir=workdir), "xlsx"
    elif ext in ("pptx", "ppt"):
        return tool_read_pptx(path, workdir=workdir), "pptx"
    elif ext in ("txt", "md", "csv", "json", "yaml", "yml", "toml", "ini"):
        try:
            with open(path, "r", encoding="utf-8", errors="replace") as f:
                return f.read()[:8000], ext
        except Exception as e:
            return f"❌ Failed to read {path}: {e}", "text"
    return f"❌ Unsupported file type: {ext or 'unknown'}", ext


def _extract_document_segments(
    path: str,
    file_type: str = "",
    filename: str = "",
) -> list[dict]:
    """Extract document text as layout-aware segments with per-page/section metadata.

    Returns a list of dicts: {"text": str, "metadata": {"source": ..., "type": ..., ...}}
    For PDFs, each segment is one page.  For DOCX, one segment per section.
    For XLSX, one segment per sheet.  For PPTX, one segment per slide.
    """
    ext = file_type.lower().lstrip(".") if file_type else ""
    if not ext and path:
        ext = path.rsplit(".", 1)[-1].lower() if "." in path else ""

    name = filename or (path.split("/")[-1] if path else "document")
    base_meta: dict = {"source": name, "type": ext or "text"}

    if ext == "pdf":
        try:
            import pypdf
            full = path if os.path.isabs(path) else os.path.join("/tmp", path)
            reader = pypdf.PdfReader(full)
            total = len(reader.pages)
            segments = []
            form_fields = []
            try:
                fields = reader.get_fields() or {}
                form_fields = sorted(str(k) for k in fields.keys())
            except Exception:
                form_fields = []

            table_rows_total = 0
            for i, page in enumerate(reader.pages):
                text = (page.extract_text() or "").strip()
                lines = [ln for ln in text.splitlines() if ln.strip()]
                table_rows = []
                for ln in lines:
                    if "  " in ln:
                        cols = [c.strip() for c in ln.split("  ") if c.strip()]
                        if len(cols) >= 2:
                            table_rows.append(" | ".join(cols))
                table_rows_total += len(table_rows)

                if table_rows:
                    text += "\n\n[Extracted Table]\n" + "\n".join(table_rows[:60])

                if text:
                    segments.append({
                        "text": text,
                        "metadata": {
                            **base_meta,
                            "page": i + 1,
                            "total_pages": total,
                            "table_rows": len(table_rows),
                            "form_fields": form_fields,
                        },
                    })
            if not segments:
                # OCR fallback for scanned PDFs if optional dependencies are available.
                try:
                    from pdf2image import convert_from_path
                    from io import BytesIO
                    from ..vision import ocr_image_bytes

                    images = convert_from_path(full, first_page=1, last_page=min(total, 5), dpi=200)
                    ocr_segments = []
                    for idx, image in enumerate(images):
                        buf = BytesIO()
                        image.save(buf, format="PNG")
                        ocr_text = (ocr_image_bytes(buf.getvalue(), mime_type="image/png") or "").strip()
                        if ocr_text:
                            ocr_segments.append({
                                "text": ocr_text,
                                "metadata": {
                                    **base_meta,
                                    "page": idx + 1,
                                    "total_pages": total,
                                    "ocr_used": True,
                                    "form_fields": form_fields,
                                },
                            })
                    if ocr_segments:
                        return ocr_segments
                except Exception:
                    pass

                return [{"text": "❌ No extractable text found (may be a scanned PDF)", "metadata": {**base_meta, "form_fields": form_fields}}]

            for seg in segments:
                seg["metadata"]["table_rows_total"] = table_rows_total
            return segments
        except ImportError:
            return [{"text": "❌ pypdf not installed. Run: pip install pypdf", "metadata": base_meta}]
        except Exception as e:
            return [{"text": f"❌ PDF read failed: {e}", "metadata": base_meta}]

    if ext in ("docx", "doc"):
        try:
            from docx import Document
            full = path if os.path.isabs(path) else os.path.join("/tmp", path)
            doc = Document(full)
            segments: list[dict] = []
            current_parts: list[str] = []
            current_heading: str | None = None
            section_idx = 0

            def _flush(heading: str | None, parts: list[str], idx: int) -> None:
                text = "\n\n".join(parts).strip()
                if not text:
                    return
                meta: dict = {**base_meta, "section": idx}
                if heading:
                    meta["heading"] = heading
                segments.append({"text": text, "metadata": meta})

            for para in doc.paragraphs:
                style = para.style.name if para.style else ""
                if style.startswith("Heading ") and para.text.strip():
                    _flush(current_heading, current_parts, section_idx)
                    section_idx += 1
                    current_parts = []
                    current_heading = para.text.strip()
                    try:
                        level = int(style.split()[-1])
                    except (ValueError, IndexError):
                        level = 1
                    current_parts.append(f"{'#' * min(level, 4)} {para.text.strip()}")
                elif para.text.strip():
                    current_parts.append(para.text.strip())

            for tbl_idx, table in enumerate(doc.tables):
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        rows.append(" | ".join(cells))
                if rows:
                    current_parts.append(f"[Table {tbl_idx + 1}]\n" + "\n".join(rows[:100]))

            _flush(current_heading, current_parts, section_idx)
            if not segments:
                return [{"text": "❌ No extractable text found in document.", "metadata": base_meta}]
            return segments
        except ImportError:
            return [{"text": "❌ python-docx not installed. Run: pip install python-docx", "metadata": base_meta}]
        except Exception as e:
            return [{"text": f"❌ DOCX read failed: {e}", "metadata": base_meta}]

    if ext in ("xlsx", "xls"):
        try:
            import openpyxl
            full = path if os.path.isabs(path) else os.path.join("/tmp", path)
            wb = openpyxl.load_workbook(full, read_only=True, data_only=True)
            segments = []
            for sheet_name in wb.sheetnames[:5]:
                ws = wb[sheet_name]
                lines = [f"### Sheet: {sheet_name}"]
                for row in ws.iter_rows(max_row=200, values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    if any(c.strip() for c in cells):
                        lines.append(" | ".join(cells))
                if ws.max_row and ws.max_row > 200:
                    lines.append(f"… ({ws.max_row - 200} more rows)")
                segments.append({"text": "\n".join(lines), "metadata": {**base_meta, "sheet": sheet_name}})
            wb.close()
            return segments if segments else [{"text": "❌ No data found in workbook.", "metadata": base_meta}]
        except ImportError:
            return [{"text": "❌ openpyxl not installed. Run: pip install openpyxl", "metadata": base_meta}]
        except Exception as e:
            return [{"text": f"❌ XLSX read failed: {e}", "metadata": base_meta}]

    if ext in ("pptx", "ppt"):
        try:
            from pptx import Presentation
            full = path if os.path.isabs(path) else os.path.join("/tmp", path)
            prs = Presentation(full)
            total = len(prs.slides)
            segments = []
            for i, slide in enumerate(prs.slides):
                texts = [s.text.strip() for s in slide.shapes if hasattr(s, "text") and s.text.strip()]
                if texts:
                    segments.append({
                        "text": "\n".join(texts),
                        "metadata": {**base_meta, "slide": i + 1, "total_slides": total},
                    })
            return segments if segments else [{"text": "❌ No text found in presentation.", "metadata": base_meta}]
        except ImportError:
            return [{"text": "❌ python-pptx not installed. Run: pip install python-pptx", "metadata": base_meta}]
        except Exception as e:
            return [{"text": f"❌ PPTX read failed: {e}", "metadata": base_meta}]

    if ext in ("txt", "md", "csv", "json", "yaml", "yml", "toml", "ini"):
        try:
            full = path if os.path.isabs(path) else os.path.join("/tmp", path)
            with open(full, "r", encoding="utf-8", errors="replace") as f:
                raw = f.read()
            if len(raw) <= 4000:
                return [{"text": raw, "metadata": base_meta}]
            # Split into overlapping sections for large text files
            chunk, overlap = 3000, 200
            return [
                {"text": raw[i: i + chunk], "metadata": {**base_meta, "section": idx}}
                for idx, i in enumerate(range(0, len(raw), chunk - overlap))
                if raw[i: i + chunk].strip()
            ]
        except Exception as e:
            return [{"text": f"❌ Failed to read {path}: {e}", "metadata": base_meta}]

    return [{"text": f"❌ Unsupported file type: {ext or 'unknown'}", "metadata": base_meta}]


@router.post("/documents/ingest")
async def documents_ingest(request: Request):
    """Extract text from a document and ingest it into the RAG store.

    POST body:
      { "path": "/tmp/report.pdf", "file_type": "pdf", "metadata": {} }
      OR
      { "text": "...", "filename": "report.pdf", "metadata": {} }
    """
    body = {}
    try:
        body = await request.json()
    except Exception:
        pass

    text  = (body.get("text") or "").strip()
    path  = (body.get("path") or "").strip()
    ftype = (body.get("file_type") or body.get("type") or "").strip()
    filename = (body.get("filename") or (path.split("/")[-1] if path else "document")).strip()
    metadata = body.get("metadata") or {}

    if not text and not path:
        return JSONResponse({"error": "path or text is required"}, status_code=400)

    rag = get_rag_system()
    chunks_stored = 0

    if not text and path:
        # Layout-aware extraction: each page/section ingested with positional metadata
        segments = _extract_document_segments(path, ftype, filename=filename)
        errors = [s for s in segments if s["text"].startswith("❌")]
        valid = [s for s in segments if not s["text"].startswith("❌")]
        if not valid:
            return JSONResponse({"error": errors[0]["text"] if errors else "No content extracted"}, status_code=400)
        detected_type = valid[0]["metadata"].get("type", ftype)
        total_chars = sum(len(s["text"]) for s in valid)
        extraction_meta = {
            "ocr_used": any(bool(s.get("metadata", {}).get("ocr_used")) for s in valid),
            "table_rows_total": sum(int(s.get("metadata", {}).get("table_rows", 0) or 0) for s in valid),
            "form_fields": sorted(
                {
                    field
                    for s in valid
                    for field in (s.get("metadata", {}).get("form_fields") or [])
                }
            ),
        }
        for seg in valid:
            seg_meta = {**seg["metadata"], **metadata}
            try:
                chunks_stored += rag.ingest(seg["text"], metadata=seg_meta, doc_id_prefix=filename)
            except Exception:
                pass
        return {
            "filename": filename,
            "type": detected_type,
            "ingested_chunks": chunks_stored,
            "char_count": total_chars,
            "segments": len(valid),
            "extraction": extraction_meta,
            "status": "ok",
        }

    # Text-direct path
    if len(text) > 500_000:
        text = text[:500_000]
    doc_meta = {"source": filename, "type": ftype or "text", **metadata}
    try:
        chunks_stored = rag.ingest(text, metadata=doc_meta, doc_id_prefix=filename)
        return {
            "filename": filename,
            "type": ftype or "text",
            "ingested_chunks": chunks_stored,
            "char_count": len(text),
            "segments": 1,
            "status": "ok",
        }
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)


@router.post("/documents/understand")
async def documents_understand(request: Request):
    """Extract text from a document and answer a question about it using an LLM.

    POST body:
      { "path": "/tmp/report.pdf", "question": "What are the key findings?", "file_type": "pdf" }
      OR
      { "text": "...", "question": "Summarise this." }

    For documents larger than ~8 000 chars, the endpoint uses a temporary RAG
    pipeline (ingest → query) to surface the most relevant sections before
    answering, rather than truncating to the first 12 000 chars.
    """
    body = {}
    try:
        body = await request.json()
    except Exception:
        pass

    text     = (body.get("text") or "").strip()
    path     = (body.get("path") or "").strip()
    question = (body.get("question") or "Summarise this document.").strip()
    ftype    = (body.get("file_type") or body.get("type") or "").strip()
    filename = (body.get("filename") or (path.split("/")[-1] if path else "document")).strip()

    if not text and not path:
        return JSONResponse({"error": "path or text is required"}, status_code=400)
    if not question:
        return JSONResponse({"error": "question is required"}, status_code=400)

    try:
        question = check_user_task(question, policy_profile=_config.get("safety_profile", "standard"))
    except GuardrailViolation as exc:
        return _api_error(exc.reason, exc.code, 422)

    detected_type = ftype
    if not text and path:
        text, detected_type = _extract_document_text(path, ftype)
        if text.startswith("❌"):
            return JSONResponse({"error": text}, status_code=400)

    _UNDERSTAND_RAG_THRESHOLD = 8000

    rag_backed = False
    excerpt = ""

    if len(text) > _UNDERSTAND_RAG_THRESHOLD:
        # Use a fresh in-memory RAG instance so we don't pollute the global store
        try:
            from ..rag.rag_system import RAGSystem as _RAGSystem
            _tmp_rag = _RAGSystem()
            chunk, overlap = 3000, 200
            for idx, offset in enumerate(range(0, min(len(text), 90000), chunk - overlap)):
                chunk_text = text[offset: offset + chunk]
                if chunk_text.strip():
                    _tmp_rag.ingest(
                        chunk_text,
                        metadata={"source": filename, "type": detected_type or "text", "section": idx},
                        doc_id_prefix=f"_und_{filename}",
                    )
            results = _tmp_rag.query(question, top_k=5)
            if results:
                excerpt = "\n\n---\n\n".join(r["document"] for r in results)
                rag_backed = True
        except Exception:
            pass  # fall through to direct path

    if not excerpt:
        excerpt = text[:12000]

    context_label = "relevant excerpts from" if rag_backed else "the content of"
    prompt = (
        f"The following is {context_label} a document "
        f"({detected_type or 'text'}) named '{filename}':\n\n"
        f"---\n{excerpt}\n---\n\n"
        f"Question: {question}\n\n"
        "Please answer based only on the document content above."
    )

    try:
        resp, provider = call_llm_with_fallback(
            [{"role": "user", "content": prompt}], question
        )
    except AllProvidersExhausted as exc:
        return _api_error(str(exc), "no_providers", 503)
    answer = resp.get("content") or str(resp)

    return {
        "filename": filename,
        "type": detected_type,
        "question": question,
        "answer": answer,
        "provider": provider,
        "excerpt_chars": len(excerpt),
        "total_chars": len(text),
        "rag_backed": rag_backed,
    }


@router.post("/autonomy/plan")
async def autonomy_plan(request: Request):
    """Decompose a goal into a structured subtask plan without executing it."""
    data = await request.json()
    goal = (data.get("goal") or "").strip()
    if not goal:
        return JSONResponse({"error": "goal field is required"}, status_code=400)
    try:
        goal = check_user_task(goal)
    except GuardrailViolation as exc:
        _push_safety_event("block", {
            "scope": "input",
            "tool": "autonomy_plan",
            "label": goal[:120],
            "profile": _config.get("safety_profile", "standard"),
            "verdict": {"allowed": False, "reason": exc.reason, "code": exc.code, "detail": exc.detail},
        })
        return _api_error(exc.reason, exc.code, 422)
    try:
        max_subtasks = int(data.get("max_subtasks", 6))
    except Exception:
        max_subtasks = 6
    trace_id = secrets.token_hex(8)
    try:
        planner = PlanningSystem()
        steps_list: list[dict[str, Any]] = []
        if hasattr(planner, "decompose"):
            tasks = planner.decompose(goal, max_subtasks)
            steps_list = [
                {
                    "id": t.task_id,
                    "name": t.name,
                    "description": t.description,
                    "priority": t.priority,
                    "dependencies": t.dependencies,
                    "estimated_hours": t.estimated_hours,
                    "agent": classify_subtask(t.description),
                }
                for t in tasks
            ]
        else:
            raw_steps = planner.plan(goal)[:max_subtasks]
            steps_list = [
                {
                    "id": f"t{i}",
                    "name": step,
                    "description": step,
                    "priority": 1,
                    "dependencies": [],
                    "estimated_hours": 1,
                    "agent": classify_subtask(step),
                }
                for i, step in enumerate(raw_steps, 1)
            ]
        plan = {"trace_id": trace_id, "goal": goal, "steps": steps_list}
        autonomy_traces[trace_id] = {"type": "plan", "status": "ready", **plan}
        db_save_autonomy_trace(trace_id, autonomy_traces[trace_id])
        return plan
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)


@router.post("/autonomy/execute")
async def autonomy_execute(request: Request):
    data = await request.json()
    goal = (data.get("goal") or "").strip()
    if not goal:
        return JSONResponse({"error": "goal field is required"}, status_code=400)
    try:
        goal = check_user_task(goal)
    except GuardrailViolation as exc:
        _push_safety_event("block", {
            "scope": "input",
            "tool": "autonomy_execute",
            "label": goal[:120],
            "profile": _config.get("safety_profile", "standard"),
            "verdict": {"allowed": False, "reason": exc.reason, "code": exc.code, "detail": exc.detail},
        })
        return _api_error(exc.reason, exc.code, 422)
    strategy = data.get("strategy", "parallel")
    try:
        max_subtasks = int(data.get("max_subtasks", 6))
    except Exception:
        max_subtasks = 6
    trace_id = secrets.token_hex(8)
    checkpoint_events: List[Dict[str, Any]] = [{"type": "autonomy_start", "goal": goal, "trace_id": trace_id}]
    _save_autonomy_checkpoint(trace_id, 0, goal, checkpoint_events)
    try:
        orchestrator = Orchestrator(_orchestrator_llm, max_parallel=2)
        result = orchestrator.execute(goal, {"strategy": strategy, "max_subtasks": max_subtasks})
        result["trace_id"] = trace_id
        try:
            from ..agent_lineage import record_lineage_edge

            previous_subtask_id = ""
            for idx, subtask in enumerate(result.get("subtasks", []), 1):
                subtask_id = str(subtask.get("id") or f"{trace_id}:subtask:{idx}")
                record_lineage_edge(
                    parent_task_id=trace_id,
                    child_task_id=subtask_id,
                    relation="spawned_by",
                    source="autonomy_execute",
                    metadata={
                        "goal": goal,
                        "strategy": strategy,
                        "subtask_name": str(subtask.get("name") or ""),
                    },
                )
                if previous_subtask_id:
                    record_lineage_edge(
                        parent_task_id=previous_subtask_id,
                        child_task_id=subtask_id,
                        relation="next",
                        source="autonomy_execute",
                        metadata={"goal": goal},
                    )
                previous_subtask_id = subtask_id
        except Exception:
            pass
        for idx, subtask in enumerate(result.get("subtasks", []), 1):
            checkpoint_events.append({"type": "subtask_done", "trace_id": trace_id, "subtask": subtask})
            _save_autonomy_checkpoint(trace_id, idx, goal, checkpoint_events)
        checkpoint_events.append(
            {
                "type": "autonomy_done",
                "trace_id": trace_id,
                "result": result.get("result", ""),
                "plan_summary": result.get("plan_summary", ""),
                "execution_time": result.get("execution_time", 0),
            }
        )
        _save_autonomy_checkpoint(trace_id, len(checkpoint_events), goal, checkpoint_events)
        autonomy_traces[trace_id] = {"type": "execution", "goal": goal, "status": "done", **result}
        db_save_autonomy_trace(trace_id, autonomy_traces[trace_id])
        return result
    except Exception as e:
        checkpoint_events.append({"type": "error", "trace_id": trace_id, "error": str(e)})
        _save_autonomy_checkpoint(trace_id, len(checkpoint_events), goal, checkpoint_events)
        return JSONResponse({"error": str(e)}, status_code=500)


@router.get("/autonomy/trace/{trace_id}")
async def autonomy_trace(trace_id: str):
    trace = autonomy_traces.get(trace_id) or db_load_autonomy_trace(trace_id)
    if not trace:
        return JSONResponse({"error": "trace not found"}, status_code=404)
    return trace


# ── chat history ──────────────────────────────────────────────────────────────
_TITLE_SKIP_PREFIXES = (
    "[MEMORY",
    "Tool result:",
    "Continue.",
    "Noted —",
    "You have reached",
)

def _auto_title(history: list) -> str:
    for msg in history:
        if msg.get("role") != "user":
            continue
        content = str(msg.get("content") or "").strip()
        if not content:
            continue
        # Skip injected context messages (memory, KG, tool scaffolding)
        if any(content.startswith(p) for p in _TITLE_SKIP_PREFIXES):
            continue
        head = content.split("\n", 1)[0].strip()
        return (head[:77] + "...") if len(head) > 80 else head
    return "New Chat"


def _extract_markdown_messages(markdown_text: str) -> list[dict]:
    messages: list[dict] = []
    lines = (markdown_text or "").splitlines()
    current_role: str | None = None
    current_content: list[str] = []

    def _flush():
        nonlocal current_role, current_content
        if current_role is None:
            return
        text = "\n".join(current_content).strip()
        if text:
            messages.append({"role": current_role, "content": text})
        current_role = None
        current_content = []

    for line in lines:
        token = line.strip().upper()
        if token in {"## USER", "**YOU:**"}:
            _flush()
            current_role = "user"
            continue
        if token in {"## ASSISTANT", "**ASSISTANT:**"}:
            _flush()
            current_role = "assistant"
            continue
        if current_role is not None:
            current_content.append(line)

    _flush()
    return messages


def _load_project_collaborators(pid: str) -> list[dict]:
    raw = db_load_pref(f"project_collaborators:{pid}", "[]")
    try:
        data = json.loads(raw)
        return data if isinstance(data, list) else []
    except Exception:
        return []


def _save_project_collaborators(pid: str, collaborators: list[dict]) -> None:
    db_save_pref(f"project_collaborators:{pid}", json.dumps(collaborators))


def _load_instruction_history() -> list[dict]:
    raw = db_load_pref("instructions_history_v1", "[]")
    try:
        data = json.loads(raw)
        return data if isinstance(data, list) else []
    except Exception:
        return []


def _save_instruction_history(entries: list[dict]):
    # Keep only recent history to avoid unbounded pref growth.
    db_save_pref("instructions_history_v1", json.dumps(entries[-200:], separators=(",", ":")))


def _append_instruction_version(previous: str, current: str, project_id: str = ""):
    if previous == current:
        return
    now = datetime.now(timezone.utc).isoformat().replace("+00:00", "Z")
    history = _load_instruction_history()
    history.append(
        {
            "id": str(uuid.uuid4()),
            "project_id": project_id or None,
            "previous": previous,
            "current": current,
            "changed_at": now,
        }
    )
    _save_instruction_history(history)


# ── provider health ────────────────────────────────────────────────────────────
@router.get("/providers/health")
async def provider_health():
    """Quick pre-flight check on each provider."""
    import time
    from ..agent import PROVIDERS, _has_key, _is_rate_limited, _cooldowns

    results = {}
    for pid, cfg in PROVIDERS.items():
        has_key = _has_key(cfg)
        cooling = _is_rate_limited(pid)
        cd_left = max(0, int(_cooldowns.get(pid, 0) - time.time()))
        results[pid] = {
            "label":     cfg["label"],
            "has_key":   has_key,
            "cooling":   cooling,
            "cd_left":   cd_left,
            "available": has_key and not cooling,
        }
    return {"health": results, "ts": time.time()}


# ── message reactions ─────────────────────────────────────────────────────────

@router.post("/reactions")
async def add_reaction(request: Request):
    data = await request.json()
    rid  = str(uuid.uuid4())[:8]
    _reactions[rid] = {
        "chat_id":  data.get("chat_id"),
        "msg_idx":  data.get("msg_idx"),
        "reaction": data.get("reaction"),   # "up" or "down"
        "text":     data.get("text","")[:200],
    }
    return {"id": rid}

@router.get("/reactions")
def get_reactions(chat_id: str = ""):
    if chat_id:
        return {"reactions": {k:v for k,v in _reactions.items() if v.get("chat_id")==chat_id}}
    return {"reactions": _reactions}


# ── pins ──────────────────────────────────────────────────────────────────────
_pins: set = set(get_pinned_chats())


# ── Nexus Cloud registration ──────────────────────────────────────────────────
async def _register_with_nexus_cloud():
    """Register this AI node with Nexus Cloud (non-blocking, best-effort)."""
    import httpx
    cloud_url = os.getenv("NEXUS_CLOUD_URL", "").rstrip("/")
    if not cloud_url:
        return
    api_key   = os.getenv("NEXUS_CLOUD_API_KEY", "")
    public_url = os.getenv("PUBLIC_URL", f"http://localhost:{os.getenv('PORT', '8000')}")
    payload = {
        "id": "nexus-ai",
        "name": "Nexus AI",
        "description": "Autonomous AI assistant with multi-provider fallback, memory, and RAG",
        "upstreamUrl": public_url,
        "mode": "standalone",
        "exposed": True,
        "health": "healthy",
        "capabilities": ["ai", "chat", "rag", "memory", "autonomy", "multi-provider"],
    }
    headers = {"Content-Type": "application/json", "Accept": "application/json"}
    if api_key:
        headers["X-Api-Key"] = api_key
    try:
        async with httpx.AsyncClient(timeout=10) as client:
            res = await client.post(f"{cloud_url}/api/v1/tools", json=payload, headers=headers)
            if res.is_success:
                print(f"[nexus-cloud] Registered with Nexus Cloud at {cloud_url}")
            else:
                print(f"[nexus-cloud] Registration rejected: {res.status_code}")
    except Exception as e:
        print(f"[nexus-cloud] Could not reach Nexus Cloud — continuing ({e})")

async def _heartbeat_loop():
    """Send a heartbeat every 30 s so Cloud knows this node is alive."""
    import httpx
    cloud_url = os.getenv("NEXUS_CLOUD_URL", "").rstrip("/")
    if not cloud_url:
        return
    api_key    = os.getenv("NEXUS_CLOUD_API_KEY", "")
    public_url = os.getenv("PUBLIC_URL", f"http://localhost:{os.getenv('PORT', '8000')}")
    headers    = {"Content-Type": "application/json"}
    if api_key:
        headers["X-Api-Key"] = api_key
    while True:
        await asyncio.sleep(30)
        try:
            async with httpx.AsyncClient(timeout=5) as client:
                await client.post(
                    f"{cloud_url}/api/v1/tools/nexus-ai/heartbeat",
                    json={"health": "healthy", "upstreamUrl": public_url},
                    headers=headers,
                )
        except Exception:
            pass


# ── Sprint E: per-message feedback ────────────────────────────────────────────

_FEEDBACK_TRACE_CONSENT_KEY = "feedback_trace_opt_in"
_FEEDBACK_TRACE_STORE_KEY = "feedback_trace_events"


def _feedback_trace_opt_in_enabled() -> bool:
    raw = db_load_pref(_FEEDBACK_TRACE_CONSENT_KEY, "false")
    if isinstance(raw, bool):
        return raw
    return str(raw).strip().lower() in {"1", "true", "yes", "on"}


def _set_feedback_trace_opt_in(enabled: bool) -> None:
    db_save_pref(_FEEDBACK_TRACE_CONSENT_KEY, bool(enabled))


def _load_feedback_trace_events(limit: int = 5000) -> list[dict]:
    raw = db_load_pref(_FEEDBACK_TRACE_STORE_KEY, [])
    rows: list[dict] = []
    if isinstance(raw, list):
        rows = raw
    elif isinstance(raw, str) and raw.strip():
        try:
            parsed = json.loads(raw)
            if isinstance(parsed, list):
                rows = parsed
        except Exception:
            rows = []
    safe_limit = max(1, min(int(limit or 5000), 50000))
    return list(reversed(rows))[:safe_limit]


def _append_feedback_trace_event(payload: dict) -> dict:
    rows = db_load_pref(_FEEDBACK_TRACE_STORE_KEY, [])
    if isinstance(rows, str) and rows.strip():
        try:
            rows = json.loads(rows)
        except Exception:
            rows = []
    if not isinstance(rows, list):
        rows = []

    row = {
        "id": f"ftrace_{uuid.uuid4().hex[:12]}",
        "created_at": datetime.now(timezone.utc).isoformat().replace("+00:00", "Z"),
        "chat_id": str(payload.get("chat_id") or ""),
        "message_idx": int(payload.get("message_idx") or 0),
        "prompt": str(payload.get("prompt") or ""),
        "response": str(payload.get("response") or ""),
        "provider": str(payload.get("provider") or ""),
        "model": str(payload.get("model") or ""),
        "persona": str(payload.get("persona") or get_active_persona_name() or ""),
        "session_id": str(payload.get("session_id") or ""),
        "meta": payload.get("meta") if isinstance(payload.get("meta"), dict) else {},
    }
    rows.append(row)
    db_save_pref(_FEEDBACK_TRACE_STORE_KEY, json.dumps(rows[-5000:]))
    return row


def _feedback_rows_to_jsonl(rows: list[dict]) -> str:
    return "\n".join(json.dumps(item, ensure_ascii=False) for item in rows) + ("\n" if rows else "")


def _feedback_rows_to_alpaca(rows: list[dict]) -> list[dict]:
    out: list[dict] = []
    for row in rows:
        out.append(
            {
                "instruction": str(row.get("prompt") or f"Chat {row.get('chat_id')} message {row.get('message_idx')}").strip(),
                "input": "",
                "output": str(row.get("response") or row.get("reaction") or "").strip(),
                "metadata": {
                    "provider": row.get("provider"),
                    "model": row.get("model"),
                    "persona": row.get("persona"),
                    "chat_id": row.get("chat_id"),
                    "message_idx": row.get("message_idx"),
                },
            }
        )
    return out


def _feedback_rows_to_sharegpt(rows: list[dict]) -> list[dict]:
    out: list[dict] = []
    for row in rows:
        out.append(
            {
                "id": str(row.get("id") or f"{row.get('chat_id')}:{row.get('message_idx')}"),
                "conversations": [
                    {"from": "human", "value": str(row.get("prompt") or "").strip()},
                    {"from": "gpt", "value": str(row.get("response") or row.get("reaction") or "").strip()},
                ],
                "meta": {
                    "provider": row.get("provider"),
                    "model": row.get("model"),
                    "persona": row.get("persona"),
                    "chat_id": row.get("chat_id"),
                    "message_idx": row.get("message_idx"),
                },
            }
        )
    return out


def _derive_test_tags(prompt: str, response: str) -> list[str]:
    text = f"{prompt} {response}".lower()
    tags: list[str] = []
    if any(k in text for k in ("code", "python", "bug", "refactor", "function")):
        tags.append("coding")
    if any(k in text for k in ("reason", "logic", "why", "explain", "math")):
        tags.append("reasoning")
    if any(k in text for k in ("policy", "safe", "guardrail", "compliance")):
        tags.append("safety")
    if any(k in text for k in ("rag", "retrieve", "document", "source")):
        tags.append("retrieval")
    return tags or ["general"]


def _keywords_from_response(response: str, limit: int = 6) -> list[str]:
    counts: dict[str, int] = {}
    for token in str(response or "").lower().replace("\n", " ").split(" "):
        clean = "".join(ch for ch in token if ch.isalnum())
        if len(clean) < 4:
            continue
        counts[clean] = counts.get(clean, 0) + 1
    ranked = sorted(counts.items(), key=lambda kv: kv[1], reverse=True)
    return [word for word, _ in ranked[: max(1, limit)]]


def _build_auto_test_cases(rows: list[dict], max_cases: int = 100) -> list[dict]:
    cases: list[dict] = []
    seen: set[str] = set()
    for row in rows:
        prompt = str(row.get("prompt") or "").strip()
        response = str(row.get("response") or "").strip()
        if len(prompt) < 12 or len(response) < 12:
            continue
        key = f"{prompt[:160]}::{response[:160]}"
        if key in seen:
            continue
        seen.add(key)
        expected_keywords = _keywords_from_response(response, limit=5)
        case_id = f"atc_{uuid.uuid4().hex[:10]}"
        cases.append(
            {
                "id": case_id,
                "created_at": datetime.now(timezone.utc).isoformat().replace("+00:00", "Z"),
                "prompt": prompt[:2000],
                "expected_keywords": expected_keywords,
                "tags": _derive_test_tags(prompt, response),
                "source": {
                    "trace_id": row.get("id"),
                    "chat_id": row.get("chat_id"),
                    "message_idx": row.get("message_idx"),
                    "provider": row.get("provider"),
                    "model": row.get("model"),
                },
            }
        )
        if len(cases) >= max_cases:
            break
    return cases


_AUTO_TEST_CASES_KEY = "auto_generated_test_cases"


@router.get("/debug/profile-pack")
def debug_profile_pack(persona: str = "", sid: str = "", use_session_dir: bool = False):
    """Debug view for runtime-resolved profile layers and precedence."""
    effective_persona = (persona or get_active_persona_name() or "general").strip()
    base_dir = get_session_dir(sid) if (sid and use_session_dir) else None
    inspected = inspect_profile_pack(base_dir=base_dir, persona_name=effective_persona)
    return {
        "persona": effective_persona,
        "session_id": sid or None,
        "session_safety_profile": get_session_safety_profile(sid) if sid else None,
        "using_session_dir": bool(sid and use_session_dir),
        "base_dir": base_dir,
        "profile": inspected,
    }


@router.get("/feedback/consent")
def feedback_trace_consent_status():
    return {
        "trace_opt_in": _feedback_trace_opt_in_enabled(),
    }


@router.post("/feedback/consent")
async def set_feedback_trace_consent(request: Request):
    body = await _read_json_body(request, "invalid JSON body")
    enabled = bool(body.get("trace_opt_in", False))
    _set_feedback_trace_opt_in(enabled)
    return {"trace_opt_in": enabled}


@router.post("/feedback/trace")
async def save_feedback_trace(request: Request):
    """Store opt-in interaction traces for training/export pipelines."""
    if not _feedback_trace_opt_in_enabled():
        return _api_error("trace collection is disabled (opt-in required)", "consent_required", 403)

    body = await _read_json_body(request, "invalid JSON body")
    row = _append_feedback_trace_event(body)
    return {"saved": True, "trace": row}

@router.post("/feedback/{chat_id}/{message_idx}")
async def save_message_feedback(chat_id: str, message_idx: int, request: Request):
    """Store a 👍/👎 reaction for a specific message.

    POST body: {"reaction": "thumbs_up" | "thumbs_down", "provider": "...", "model": "..."}
    """
    from ..db import save_feedback as db_save_feedback
    body: dict = {}
    try:
        body = await request.json()
    except Exception:
        pass
    reaction = (body.get("reaction") or "").strip()
    if reaction not in ("thumbs_up", "thumbs_down"):
        return _api_error("reaction must be 'thumbs_up' or 'thumbs_down'", "validation_error", 422)
    db_save_feedback(
        chat_id=chat_id,
        message_idx=message_idx,
        reaction=reaction,
        provider=body.get("provider", ""),
        model=body.get("model", ""),
    )
    return {"saved": True, "chat_id": chat_id, "message_idx": message_idx, "reaction": reaction}


@router.get("/feedback/export")
def feedback_export(limit: int = 5000, format: str = "json", include_trace: bool = True):
    """Export message feedback/traces for training pipelines.

    Supported formats: json, jsonl, alpaca, sharegpt.
    """
    from ..db import load_feedback_export, get_feedback_stats
    data = load_feedback_export(limit)
    stats = get_feedback_stats()
    trace_rows = _load_feedback_trace_events(limit=limit) if include_trace else []

    fmt = (format or "json").strip().lower()
    combined_rows = list(trace_rows)
    if not combined_rows:
        # Backward-compatible fallback when only reaction-level feedback exists.
        for item in data:
            combined_rows.append(
                {
                    "id": f"feedback_{item.get('chat_id')}:{item.get('message_idx')}",
                    "chat_id": item.get("chat_id"),
                    "message_idx": item.get("message_idx"),
                    "reaction": item.get("reaction"),
                    "provider": item.get("provider"),
                    "model": item.get("model"),
                    "created_at": item.get("ts"),
                    "prompt": "",
                    "response": str(item.get("reaction") or ""),
                    "persona": "",
                }
            )

    if fmt == "jsonl":
        return Response(_feedback_rows_to_jsonl(combined_rows), media_type="application/x-ndjson")

    if fmt == "alpaca":
        return {
            "format": "alpaca",
            "count": len(combined_rows),
            "data": _feedback_rows_to_alpaca(combined_rows),
            "stats": stats,
            "trace_opt_in": _feedback_trace_opt_in_enabled(),
        }

    if fmt == "sharegpt":
        return {
            "format": "sharegpt",
            "count": len(combined_rows),
            "data": _feedback_rows_to_sharegpt(combined_rows),
            "stats": stats,
            "trace_opt_in": _feedback_trace_opt_in_enabled(),
        }

    if fmt != "json":
        return _api_error("format must be one of: json, jsonl, alpaca, sharegpt", "validation_error", 422)

    return {
        "format": "json",
        "stats":  stats,
        "count":  len(data),
        "data":   data,
        "trace_opt_in": _feedback_trace_opt_in_enabled(),
        "trace_count": len(trace_rows),
        "trace_data": trace_rows,
    }


@router.get("/feedback/stats")
def feedback_stats():
    """Return aggregate thumbs-up / thumbs-down counts."""
    from ..db import get_feedback_stats
    stats = get_feedback_stats()
    trace_rows = _load_feedback_trace_events(limit=5000)
    stats["trace_opt_in"] = _feedback_trace_opt_in_enabled()
    stats["trace_total"] = len(trace_rows)
    return stats


@router.get("/architecture/hierarchy")
def architecture_hierarchy():
    """Return a live scaffold of the AI-system hierarchy.

    This endpoint is intentionally read-only for now and acts as an architecture
    planning contract for future model/agent/workflow expansion.
    """
    from ..agent import get_providers_list
    from ..agents import list_agents
    from ..architecture import build_runtime_hierarchy

    return build_runtime_hierarchy(
        providers=get_providers_list(),
        specialist_agents=list_agents(),
    )


@router.post("/architecture/blueprints")
async def create_architecture_blueprint(request: Request):
    from ..agent import get_providers_list
    from ..agents import list_agents
    from ..architecture import build_runtime_hierarchy
    from ..db import save_architecture_blueprint

    body: dict = {}
    try:
        body = await request.json()
    except Exception:
        body = {}

    name = str(body.get("name") or "default").strip() or "default"
    notes = str(body.get("notes") or "").strip()
    use_runtime = bool(body.get("use_runtime", True))

    snapshot = body.get("snapshot") if isinstance(body.get("snapshot"), dict) else None
    if snapshot is None and use_runtime:
        snapshot = build_runtime_hierarchy(
            providers=get_providers_list(),
            specialist_agents=list_agents(),
        )

    if snapshot is None:
        return _api_error("snapshot is required when use_runtime=false", "validation_error", 422)

    created = save_architecture_blueprint(name=name, snapshot=snapshot, notes=notes)
    return {"blueprint": created}


@router.get("/architecture/blueprints")
def list_architecture_blueprints(name: str = "", limit: int = 50):
    from ..db import list_architecture_blueprints as db_list_architecture_blueprints

    items = db_list_architecture_blueprints(name=name, limit=limit)
    return {"blueprints": items, "total": len(items)}


@router.get("/architecture/blueprints/{name}")
def get_architecture_blueprint(name: str, version: int = 0):
    from ..db import load_architecture_blueprint

    data = load_architecture_blueprint(name=name, version=version if version > 0 else None)
    if not data:
        return _api_error(f"architecture blueprint '{name}' not found", "not_found", 404)
    return data


@router.get("/architecture/blueprints/{name}/versions")
def list_architecture_blueprint_versions(name: str, limit: int = 100):
    from ..db import list_architecture_blueprints as db_list_architecture_blueprints

    rows = db_list_architecture_blueprints(name=name, limit=max(1, min(int(limit), 500)))
    versions = [r for r in rows if str(r.get("name") or "") == name]
    versions.sort(key=lambda r: int(r.get("version") or 0), reverse=True)
    return {
        "name": name,
        "versions": versions,
        "latest_version": versions[0].get("version") if versions else None,
        "total": len(versions),
    }


@router.post("/architecture/blueprints/{name}/activate")
async def activate_architecture_blueprint_version(name: str, request: Request):
    from ..db import load_architecture_blueprint, save_pref as _save_pref

    body = await _read_json_body(request, "invalid JSON body")
    version = int(body.get("version") or 0)
    if version <= 0:
        return _api_error("version must be a positive integer", "validation_error", 422)

    bp = load_architecture_blueprint(name=name, version=version)
    if not bp:
        return _api_error(f"architecture blueprint '{name}' version {version} not found", "not_found", 404)

    marker = {
        "name": name,
        "version": version,
        "activated_at": datetime.now(timezone.utc).isoformat().replace("+00:00", "Z"),
        "activated_by": str(body.get("actor") or body.get("username") or "system"),
    }
    _save_pref(f"arch_bp:{name}:active", json.dumps(marker))
    return {"ok": True, "active": marker, "blueprint": bp}


@router.get("/architecture/blueprints/{name}/active")
def get_architecture_blueprint_active(name: str):
    raw = db_load_pref(f"arch_bp:{name}:active", "")
    if not raw:
        return {"name": name, "active": None}
    try:
        parsed = json.loads(raw)
    except Exception:
        parsed = None
    return {"name": name, "active": parsed}


@router.get("/architecture/registry/{name}")
def get_architecture_registry(name: str, version: int = 0):
    from ..db import load_architecture_registry

    data = load_architecture_registry(name=name, version=version if version > 0 else None)
    if not data:
        return _api_error(f"architecture registry '{name}' not found", "not_found", 404)

    if "counts" not in data:
        snapshot = data.get("snapshot") if isinstance(data.get("snapshot"), dict) else {}
        nodes = snapshot.get("nodes") if isinstance(snapshot.get("nodes"), list) else []
        edges = snapshot.get("edges") if isinstance(snapshot.get("edges"), list) else []
        data = {
            **data,
            "counts": {
                "nodes": max(1, len(nodes)),
                "edges": max(1, len(edges)),
            },
        }
    return data


# ── Sprint G: Simulate Endpoint ───────────────────────────────────────────────

@router.post("/simulate")
async def run_simulation(request: Request):
    """Run a swarm prediction simulation (MiroFish-inspired).

    POST body:
        topic      — required  e.g. "Will AI replace software engineers by 2030?"
        seed       — optional context / background text
        n_personas — optional int 2-8 (default 5)
        n_rounds   — optional int 1-5 (default 3)

    Returns a SimulationResult dict including prediction, confidence, personas, rounds,
    minority_views and a full Markdown report.
    """
    body: dict = {}
    try:
        body = await request.json()
    except Exception:
        pass

    topic = (body.get("topic") or "").strip()
    if not topic:
        return _api_error("topic is required", "validation_error", 422)

    try:
        check_user_task(topic)
    except GuardrailViolation as exc:
        _push_safety_event("block", {
            "scope": "input",
            "tool": "simulate",
            "label": topic[:120],
            "profile": _config.get("safety_profile", "standard"),
            "verdict": {"allowed": False, "reason": exc.reason, "code": exc.code, "detail": exc.detail},
        })
        return _api_error(exc.reason, exc.code, 422)

    n_personas = max(2, min(int(body.get("n_personas", 5)), 8))
    n_rounds   = max(1, min(int(body.get("n_rounds",   3)), 5))
    seed       = (body.get("seed") or "").strip()

    def _sim_llm(msgs):
        try:
            res, _ = call_llm_with_fallback(msgs, "simulation")
            if isinstance(res, dict):
                if res.get("action") == "respond":
                    return res.get("content", "")
                return json.dumps(res)
            return str(res)
        except Exception as _se:
            return f"error: {_se}"

    try:
        from ..simulation import SimulationEngine
        engine = SimulationEngine(_sim_llm, max_personas=8, max_rounds=5)
        result = engine.run(topic, seed, n_personas, n_rounds)
        return result.to_dict()
    except Exception as exc:
        return _api_error(str(exc), "simulation_error", 500)


@router.delete("/tasks/{trace_id}")
def delete_task_trace(trace_id: str):
    deleted = _delete_trace(trace_id)
    execution_traces.pop(trace_id, None)
    if not deleted:
        return _api_error("trace not found", "not_found", 404)
    return {"deleted": trace_id, "ok": True}


@router.get("/tasks/search")
def search_task_traces(tool: str = "", event_type: str = "", error: str = "", limit: int = 100):
    safe_limit = max(1, min(int(limit), 500))
    tool_filter = (tool or "").strip().lower()
    type_filter = (event_type or "").strip().lower()
    error_filter = (error or "").strip().lower()

    traces = list_tasks(limit=safe_limit * 2).get("traces", [])
    matched: list[dict] = []
    for trace in traces:
        trace_id = str(trace.get("trace_id") or "").strip()
        if not trace_id:
            continue
        events = execution_traces.get(trace_id)
        if events is None:
            events = db_load_execution_trace(trace_id)
        if events is None:
            cp = _get_latest_checkpoint(trace_id)
            events = cp.get("events", []) if cp else []

        event_list = events if isinstance(events, list) else []
        has_tool = (not tool_filter) or any(
            tool_filter in str(evt.get("action", "")).lower() for evt in event_list if isinstance(evt, dict)
        )
        has_type = (not type_filter) or any(
            type_filter == str(evt.get("type", "")).lower() for evt in event_list if isinstance(evt, dict)
        )
        has_error = (not error_filter) or any(
            error_filter in str(evt.get("message", "")).lower() for evt in event_list if isinstance(evt, dict)
        )
        if has_tool and has_type and has_error:
            matched.append(
                {
                    "trace_id": trace_id,
                    "events": len(event_list),
                    "task": trace.get("task", ""),
                    "last_active": trace.get("last_active", trace.get("updated_at", "")),
                }
            )
            if len(matched) >= safe_limit:
                break

    return {"results": matched, "count": len(matched)}


@router.get("/tasks/{trace_id}/export")
def export_task_trace(trace_id: str):
    events = execution_traces.get(trace_id)
    if events is None:
        events = db_load_execution_trace(trace_id)
    checkpoints = _load_checkpoints(trace_id)
    if events is None and not checkpoints:
        return _api_error("trace not found", "not_found", 404)
    payload = {
        "trace_id": trace_id,
        "exported_at": datetime.now(timezone.utc).isoformat(),
        "events": events if isinstance(events, list) else [],
        "checkpoints": checkpoints,
    }
    return payload


@router.get("/tasks/{trace_id}/diff")
def diff_task_traces(trace_id: str, other_trace_id: str = ""):
    import difflib

    if not other_trace_id.strip():
        return _api_error("other_trace_id is required", "validation_error", 422)

    left = execution_traces.get(trace_id)
    if left is None:
        left = db_load_execution_trace(trace_id)
    right = execution_traces.get(other_trace_id)
    if right is None:
        right = db_load_execution_trace(other_trace_id)

    if left is None or right is None:
        return _api_error("one or both traces not found", "not_found", 404)

    left_lines = json.dumps(left, indent=2, sort_keys=True).splitlines()
    right_lines = json.dumps(right, indent=2, sort_keys=True).splitlines()
    diff_lines = list(
        difflib.unified_diff(
            left_lines,
            right_lines,
            fromfile=f"a/{trace_id}",
            tofile=f"b/{other_trace_id}",
            lineterm="",
        )
    )
    return {
        "trace_id": trace_id,
        "other_trace_id": other_trace_id,
        "diff": "\n".join(diff_lines),
        "left_events": len(left) if isinstance(left, list) else 0,
        "right_events": len(right) if isinstance(right, list) else 0,
    }


@router.get("/tasks/anomalies")
def task_trace_anomalies(limit: int = 50):
    safe_limit = max(1, min(int(limit), 500))
    traces = list_tasks(limit=safe_limit).get("traces", [])
    anomalies: list[dict] = []

    for trace in traces:
        trace_id = str(trace.get("trace_id") or "")
        if not trace_id:
            continue
        events = execution_traces.get(trace_id)
        if events is None:
            events = db_load_execution_trace(trace_id)
        event_list = events if isinstance(events, list) else []

        error_count = sum(1 for evt in event_list if isinstance(evt, dict) and str(evt.get("type", "")).lower() == "error")
        tool_count = sum(1 for evt in event_list if isinstance(evt, dict) and str(evt.get("type", "")).lower() in {"tool", "tool_start", "tool_result"})

        reasons: list[str] = []
        if error_count > 0:
            reasons.append(f"errors:{error_count}")
        if len(event_list) > 500:
            reasons.append("high_event_volume")
        if tool_count > 200:
            reasons.append("high_tool_activity")

        if reasons:
            anomalies.append(
                {
                    "trace_id": trace_id,
                    "events": len(event_list),
                    "errors": error_count,
                    "tool_events": tool_count,
                    "reasons": reasons,
                    "last_active": trace.get("last_active", trace.get("updated_at", "")),
                }
            )

    return {"anomalies": anomalies, "count": len(anomalies)}


# ── Execution Trace replay/resume endpoints ──────────────────────────────────

@router.get("/tasks")
def list_tasks(limit: int = 50):
    traces = _list_traces(limit=limit)
    persisted = db_list_execution_traces(limit=limit)
    seen = {t.get("trace_id") for t in traces if isinstance(t, dict)}
    for trace in persisted:
        if trace.get("trace_id") not in seen:
            traces.append(trace)
    return {"traces": traces, "count": len(traces)}


@router.get("/tasks/{trace_id}")
def get_task_trace(trace_id: str):
    # Check in-memory first (live traces), then SQLite checkpoints
    in_memory = execution_traces.get(trace_id)
    if in_memory is None:
        in_memory = db_load_execution_trace(trace_id)
    checkpoints = _load_checkpoints(trace_id)
    if in_memory is None and not checkpoints:
        return _api_error("trace not found", "not_found", 404)
    events = in_memory if in_memory is not None else []
    return {"trace_id": trace_id, "events": events, "checkpoints": len(checkpoints)}


@router.get("/tasks/{trace_id}/replay")
async def replay_task(trace_id: str):
    """Stream stored trace events as SSE with a short delay for visual replay."""
    import asyncio as _asyncio

    # Prefer in-memory events, fall back to last SQLite checkpoint's events
    stored_events = execution_traces.get(trace_id)
    if stored_events is None:
        stored_events = db_load_execution_trace(trace_id)
    if stored_events is None:
        cp = _get_latest_checkpoint(trace_id)
        stored_events = cp["events"] if cp else None
    if stored_events is None:
        return _api_error("trace not found", "not_found", 404)

    events_copy = list(stored_events)

    async def _stream():
        for evt in events_copy:
            yield f"data: {json.dumps(evt)}\n\n"
            await _asyncio.sleep(0.04)

    return StreamingResponse(_stream(), media_type="text/event-stream",
                             headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"})


@router.post("/tasks/{trace_id}/resume")
async def resume_task(trace_id: str, request: Request):
    """Resume a task from its latest checkpoint."""
    data = await request.json()
    sid = data.get("session_id", "")

    cp = _get_latest_checkpoint(trace_id)
    if not cp:
        return _api_error("no checkpoints found for this trace", "not_found", 404)

    task = cp.get("task", "")
    saved_history = cp.get("history", [])

    if not task:
        return _api_error("checkpoint has no task to resume", "invalid_request", 422)

    new_trace_id = str(uuid.uuid4())
    execution_traces[new_trace_id] = []
    db_save_execution_trace(new_trace_id, execution_traces[new_trace_id])
    loop = asyncio.get_event_loop()
    queue: asyncio.Queue = asyncio.Queue()
    stop_evt = threading.Event()
    new_stream_id = str(uuid.uuid4())
    _active_streams[new_stream_id] = stop_evt

    def _run_resume():
        try:
            for event in stream_agent_task(task, saved_history, [], stop_evt,
                                           sid=sid or "", trace_id=new_trace_id):
                if stop_evt.is_set():
                    break
                if event["type"] == "done" and sid:
                    sessions[sid] = event.get("history", saved_history)
                trace_event = {k: v for k, v in event.items() if k not in ("history", "workdir")}
                execution_traces[new_trace_id].append(trace_event)
                db_save_execution_trace(new_trace_id, execution_traces[new_trace_id])
                loop.call_soon_threadsafe(queue.put_nowait, event)
        except Exception as e:
            err_event = {"type": "error", "message": str(e)}
            execution_traces[new_trace_id].append(err_event)
            db_save_execution_trace(new_trace_id, execution_traces[new_trace_id])
            loop.call_soon_threadsafe(queue.put_nowait, err_event)
        finally:
            _active_streams.pop(new_stream_id, None)
            loop.call_soon_threadsafe(queue.put_nowait, None)

    threading.Thread(target=_run_resume, daemon=True).start()

    async def _generate():
        try:
            while True:
                event = await queue.get()
                if event is None:
                    break
                payload = {k: v for k, v in event.items() if k not in ("history", "workdir")}
                yield f"data: {json.dumps(payload)}\n\n"
        except asyncio.CancelledError:
            stop_evt.set()

    return StreamingResponse(_generate(), media_type="text/event-stream",
                             headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no",
                                      "X-Trace-Id": new_trace_id})


@router.delete("/tasks/{trace_id}")
def delete_task_trace(trace_id: str):
    deleted = _delete_trace(trace_id)
    execution_traces.pop(trace_id, None)
    deleted = db_delete_execution_trace(trace_id) or deleted
    if not deleted:
        return _api_error("trace not found", "not_found", 404)
    return {"deleted": trace_id, "ok": True}


# ── Ensemble settings endpoints ──────────────────────────────────────────────

@router.get("/settings/ensemble")
def get_ensemble_settings():
    return {
        "ensemble_mode":      _config.get("ensemble_mode", True),
        "ensemble_threshold": _config.get("ensemble_threshold", 0.4),
        "ensemble_enabled":   get_ensemble_enabled(),
    }


@router.post("/settings/ensemble")
async def update_ensemble_settings(request: Request):
    data = await request.json()
    kwargs = {}
    if "ensemble_mode" in data:
        kwargs["ensemble_mode"] = bool(data["ensemble_mode"])
    if "ensemble_threshold" in data:
        thr = float(data["ensemble_threshold"])
        if not 0.0 <= thr <= 1.0:
            return _api_error("ensemble_threshold must be between 0.0 and 1.0", "validation_error", 422)
        kwargs["ensemble_threshold"] = thr
    if kwargs:
        update_config(**kwargs)
    return {
        "ensemble_mode":      _config.get("ensemble_mode", True),
        "ensemble_threshold": _config.get("ensemble_threshold", 0.4),
        "ensemble_enabled":   get_ensemble_enabled(),
    }


@router.get("/settings/hitl")
def get_hitl_settings():
    return {
        "hitl_approval_mode": _config.get("hitl_approval_mode", "off"),
    }


@router.post("/settings/hitl")
async def update_hitl_settings(request: Request):
    data = await request.json()
    mode = str(data.get("hitl_approval_mode", "off")).lower().strip()
    if mode not in ("off", "warn", "block"):
        return _api_error("hitl_approval_mode must be one of: off, warn, block", "validation_error", 422)
    update_config(hitl_approval_mode=mode)
    return {"hitl_approval_mode": _config.get("hitl_approval_mode", "off")}


@router.get("/approvals")
def get_approvals(session_id: str = ""):
    return {
        "items": list_tool_approvals(session_id),
        "total": len(list_tool_approvals(session_id)),
    }


@router.post("/approvals/{approval_id}")
async def resolve_approval(approval_id: str, request: Request):
    data = await request.json()
    approved = bool(data.get("approved", False))
    note = str(data.get("note", ""))
    resolved = decide_tool_approval(approval_id, approved=approved, note=note)
    if not resolved:
        return _api_error("approval not found", "not_found", 404)
    return resolved


def _run_scheduled_task_extended(task: str) -> str:
    if task == "__internal_quota_cleanup__":
        return _quota_cleanup_task()
    if task.startswith("__finetune_continual__:"):
        raw = task.split(":", 1)[1].strip()
        try:
            payload = json.loads(raw)
        except Exception as exc:
            return f"continual policy parse error: {exc}"
        if not isinstance(payload, dict):
            return "continual policy payload must be an object"
        return _execute_continual_re_tune_task(payload)
    if task == "__internal_validation_program__":
        from ..validation_program import run_validation_program

        report = run_validation_program(update_baseline=False, alert_on_regression=True)
        summary = report.get("summary") if isinstance(report.get("summary"), dict) else {}
        return (
            f"validation success_rate={summary.get('success_rate', 0.0)} "
            f"tool_error_rate={summary.get('tool_error_rate', 0.0)}"
        )
    return _run_scheduled_task(task)


def startup_event() -> None:
    set_run_function(_run_scheduled_task_extended)
    restore_from_db()
    _register_quota_reset_scheduler()
    try:
        from ..benchmark import register_benchmark_schedules
        from ..validation_program import register_validation_program_schedules

        register_benchmark_schedules(_run_scheduled_task_extended)
        register_validation_program_schedules(_run_scheduled_task_extended)
    except Exception:
        pass
    asyncio.create_task(_register_with_nexus_cloud())
    asyncio.create_task(_heartbeat_loop())


@router.get("/quota/me")
def my_quota(request: Request):
    if not MULTI_USER:
        username = "nexus_admin"
    else:
        username = _read_token(request)
        if not username:
            return _api_error("Unauthorized", "unauthorized", 401)
    from ..profiles import get_quota_state
    return {"username": username, **get_quota_state(username)}


# ── API key management ────────────────────────────────────────────────────────

_VALID_SCOPES = {"chat", "read", "admin", "embeddings", "tools"}


def _generate_api_key() -> tuple[str, str, str]:
    """Return (raw_key, key_hash, key_prefix)."""
    raw = "nxk_" + secrets.token_urlsafe(40)
    key_hash = hashlib.sha256(raw.encode()).hexdigest()
    prefix = raw[:12]
    return raw, key_hash, prefix


@router.get("/auth/api-keys")
def list_api_keys_endpoint(request: Request):
    username = require_auth(request)
    keys = db_list_api_keys(username)
    safe = []
    for k in keys:
        safe.append({
            "id": k["id"],
            "key_prefix": k["key_prefix"],
            "name": k["name"],
            "scopes": k["scopes"],
            "created_at": k["created_at"],
            "last_used_at": k.get("last_used_at"),
            "revoked_at": k.get("revoked_at"),
            "active": k.get("revoked_at") is None,
        })
    return {"keys": safe, "total": len(safe)}


# ── email verification ────────────────────────────────────────────────────────

_SMTP_HOST = os.getenv("SMTP_HOST", "")
_SMTP_PORT = int(os.getenv("SMTP_PORT", "587"))
_SMTP_USER = os.getenv("SMTP_USER", "")
_SMTP_PASS = os.getenv("SMTP_PASS", "")
_EMAIL_FROM = os.getenv("EMAIL_FROM", _SMTP_USER or "nexus@localhost")
_APP_URL = os.getenv("APP_URL", "http://localhost:8000")


def _send_verification_email(email: str, token: str, username: str) -> bool:
    link = f"{_APP_URL}/auth/verify-email?token={token}&username={username}"
    body = f"Hello {username},\n\nVerify your email:\n{link}\n\nThis link expires in 24 hours."
    if not _SMTP_HOST:
        print(f"[email-verify] Token for {username}: {token} (SMTP not configured)")
        return True
    try:
        import smtplib
        from email.mime.text import MIMEText
        msg = MIMEText(body)
        msg["Subject"] = "Verify your Nexus AI account"
        msg["From"] = _EMAIL_FROM
        msg["To"] = email
        with smtplib.SMTP(_SMTP_HOST, _SMTP_PORT) as s:
            s.starttls()
            if _SMTP_USER:
                s.login(_SMTP_USER, _SMTP_PASS)
            s.send_message(msg)
        return True
    except Exception as e:
        print(f"[email-verify] SMTP error: {e}")
        return False


# ── OAuth2 / OIDC SSO ─────────────────────────────────────────────────────────

_OAUTH_PROVIDERS: dict[str, dict] = {
    "google": {
        "auth_url": "https://accounts.google.com/o/oauth2/v2/auth",
        "token_url": "https://oauth2.googleapis.com/token",
        "userinfo_url": "https://www.googleapis.com/oauth2/v3/userinfo",
        "client_id_env": "GOOGLE_CLIENT_ID",
        "client_secret_env": "GOOGLE_CLIENT_SECRET",  # pragma: allowlist secret
        "scope": "openid email profile",
    },
    "github": {
        "auth_url": "https://github.com/login/oauth/authorize",
        "token_url": "https://github.com/login/oauth/access_token",
        "userinfo_url": "https://api.github.com/user",
        "client_id_env": "GITHUB_CLIENT_ID",
        "client_secret_env": "GITHUB_CLIENT_SECRET",  # pragma: allowlist secret
        "scope": "read:user user:email",
    },
}


# ── quota reset scheduler ─────────────────────────────────────────────────────

def _quota_cleanup_task() -> str:
    """Purge stale daily quota pref keys older than 8 days."""
    try:
        from ..profiles import cleanup_stale_quota_days

        deleted = cleanup_stale_quota_days(days_to_keep=8)
        return f"quota cleanup: removed {deleted} stale daily keys"
    except Exception as e:
        return f"quota cleanup error: {e}"


def _register_quota_reset_scheduler():
    """Register weekly quota cleanup job if not already registered."""
    existing = list_jobs()
    if any(getattr(j, "name", None) == "quota-weekly-cleanup" for j in existing):
        return
    schedule_job(
        name="quota-weekly-cleanup",
        task="__internal_quota_cleanup__",
        schedule="0 3 * * 0",  # every Sunday at 03:00 UTC
    )


# ─────────────────────────────────────────────────────────────────────────────
# Fine-tuning API  — persisted compatibility lifecycle (/v1/fine-tuning/jobs)
# ─────────────────────────────────────────────────────────────────────────────

def _run_fine_tuning_job(job_id: str):
    """Background lifecycle for compatibility fine-tuning jobs.

    This keeps OpenAI-compatible job states realistic while the backend
    remains a compatibility implementation rather than a real trainer.
    """
    try:
        job = db_get_fine_tuning_job(job_id)
        if not job or job.get("status") != "queued":
            return

        db_update_fine_tuning_job(job_id, status="running")
        db_create_fine_tuning_job_event(job_id, "Job status changed to running", data={"status": "running"})
        time.sleep(0.6)

        job = db_get_fine_tuning_job(job_id)
        if not job or job.get("status") == "cancelled":
            return

        training_meta = _load_file_meta(str(job.get("training_file") or "")) or {}
        trained_tokens = max(128, int(training_meta.get("bytes", 0) // 4))
        ft_model = f"ft:{job.get('model', 'model')}:{job_id[-6:]}"
        db_update_fine_tuning_job(
            job_id,
            status="succeeded",
            fine_tuned_model=ft_model,
            trained_tokens=trained_tokens,
            finished_at=int(time.time()),
            result_files=[str(job.get("training_file") or "")],
        )
        db_create_fine_tuning_job_event(
            job_id,
            "Job completed successfully",
            data={"status": "succeeded", "fine_tuned_model": ft_model, "trained_tokens": trained_tokens},
        )
        try:
            from ..eval_pipeline import run_regression_benchmark

            suites = ["gsm8k", "arc", "safety"]
            hyperparams = job.get("hyperparameters") if isinstance(job.get("hyperparameters"), dict) else {}
            regression_provider = str(hyperparams.get("regression_provider") or "offline").strip() or "offline"
            regression = run_regression_benchmark(
                model=str(job.get("model") or "nexus-prime-base"),
                provider=regression_provider,
                suites=suites,
                threshold=0.05,
                n_samples=8,
            )
            db_create_fine_tuning_job_event(
                job_id,
                "Post-train regression benchmark completed",
                data={
                    "suites": suites,
                    "provider": regression_provider,
                    "overall_regression": bool(regression.get("overall_regression")),
                    "current_avg": regression.get("current_avg"),
                },
            )
        except Exception as exc:
            db_create_fine_tuning_job_event(
                job_id,
                "Post-train regression benchmark skipped",
                level="warning",
                data={"error": str(exc)[:300]},
            )
    except Exception as exc:
        db_update_fine_tuning_job(
            job_id,
            status="failed",
            finished_at=int(time.time()),
            error={"message": str(exc), "code": "fine_tuning_job_error"},
        )
        db_create_fine_tuning_job_event(
            job_id,
            "Job failed",
            level="error",
            data={"status": "failed", "error": str(exc)},
        )


def _training_rows_from_feedback(include_trace: bool = True, limit: int = 5000) -> list[dict]:
    from ..db import load_feedback_export

    feedback_rows = load_feedback_export(limit=limit)
    trace_rows = _load_feedback_trace_events(limit=limit) if include_trace else []

    rows: list[dict] = []
    if trace_rows:
        for item in trace_rows:
            rows.append(
                {
                    "prompt": str(item.get("prompt") or ""),
                    "response": str(item.get("response") or ""),
                    "provider": str(item.get("provider") or ""),
                    "model": str(item.get("model") or ""),
                    "persona": str(item.get("persona") or ""),
                    "chat_id": str(item.get("chat_id") or ""),
                    "message_idx": int(item.get("message_idx") or 0),
                    "created_at": str(item.get("created_at") or ""),
                    "source": "trace",
                }
            )
    else:
        for item in feedback_rows:
            rows.append(
                {
                    "prompt": "",
                    "response": str(item.get("reaction") or ""),
                    "provider": str(item.get("provider") or ""),
                    "model": str(item.get("model") or ""),
                    "persona": "",
                    "chat_id": str(item.get("chat_id") or ""),
                    "message_idx": int(item.get("message_idx") or 0),
                    "created_at": str(item.get("ts") or ""),
                    "source": "reaction",
                }
            )
    return rows


def _dataset_checksum(rows: list[dict]) -> str:
    payload = json.dumps(rows, sort_keys=True, ensure_ascii=False, separators=(",", ":"))
    return hashlib.sha256(payload.encode("utf-8")).hexdigest()


def _create_finetune_job_from_rows(
    model: str,
    rows: list[dict],
    source: str,
    provenance_extra: dict | None = None,
    dataset_format: str = "jsonl",
    auto_start: bool = True,
) -> dict:
    from ..db import save_ft_dataset_version

    checksum = _dataset_checksum(rows)
    dataset_id = f"dsver-{uuid.uuid4().hex[:10]}"
    provenance = {
        "source": source,
        "row_count": len(rows),
        "generated_at": datetime.now(timezone.utc).isoformat().replace("+00:00", "Z"),
    }
    if isinstance(provenance_extra, dict):
        provenance.update(provenance_extra)

    dataset_version = save_ft_dataset_version(
        dataset_id=dataset_id,
        source=source,
        fmt=dataset_format,
        row_count=len(rows),
        provenance=provenance,
        checksum=checksum,
        preview_rows=rows,
    )

    training_file = f"inline://dataset/{dataset_id}"
    job = FineTuningJob(
        model=model,
        training_file=training_file,
        validation_file=None,
        hyperparameters={
            "dataset_version_id": dataset_id,
            "dataset_checksum": checksum,
            "provenance": provenance,
        },
        status="queued",
    ).model_dump()
    db_create_fine_tuning_job(job)
    db_create_fine_tuning_job_event(
        job["id"],
        "Fine-tune created",
        data={"status": "queued", "dataset_version_id": dataset_id, "source": source},
    )
    if auto_start:
        threading.Thread(target=_run_fine_tuning_job, args=(job["id"],), daemon=True).start()

    return {
        "job": {
            "id": job["id"],
            "status": job["status"],
            "model": job["model"],
            "training_file": training_file,
            "object": "finetune.job",
        },
        "dataset_version": dataset_version,
    }


def _run_model_update_regression_gate(
    model: str,
    provider: str,
    suites: list[str] | None,
    threshold: float,
    n_samples: int,
    enforce: bool,
) -> dict:
    """Run regression benchmark for model-update actions and return gate metadata."""
    from ..eval_pipeline import run_regression_benchmark

    selected_suites = [str(s).strip() for s in (suites or ["gsm8k", "arc", "safety"]) if str(s).strip()]
    result = run_regression_benchmark(
        model=model,
        provider=provider,
        suites=selected_suites,
        threshold=float(threshold),
        n_samples=max(1, int(n_samples)),
    )
    rows = result.get("suite_results") if isinstance(result.get("suite_results"), list) else []
    avg_score = 0.0
    if rows:
        avg_score = sum(float(r.get("score") or 0.0) for r in rows) / max(1, len(rows))
    has_regression = not bool(result.get("ok", True))
    return {
        "enforced": bool(enforce),
        "allowed": (not has_regression) or (not bool(enforce)),
        "regression_count": int(result.get("regression_count") or 0),
        "has_regression": has_regression,
        "threshold": float(threshold),
        "provider": provider,
        "suites": selected_suites,
        "n_samples": max(1, int(n_samples)),
        "average_score": round(avg_score, 4),
        "benchmark": result,
    }


def _run_distill_job(job_id: str):
    from ..db import (
        append_distill_job_event,
        get_distill_job,
        save_lora_adapter_version,
        update_distill_job,
    )

    job = get_distill_job(job_id)
    if not job or str(job.get("status") or "") != "queued":
        return

    append_distill_job_event(job_id, "Distillation job started")
    update_distill_job(job_id, status="running")
    teacher = str(job.get("teacher_model") or "")
    student = str(job.get("student_model") or "nexus-prime-base")
    provider = str(job.get("provider") or "auto")
    cfg = job.get("config") if isinstance(job.get("config"), dict) else {}
    prompts = cfg.get("prompts") if isinstance(cfg.get("prompts"), list) else []
    if not prompts:
        prompts = [
            "Explain how adapter hot-swap works in a sovereign deployment.",
            "Design a retry-safe continual fine-tune policy.",
            "Summarize RLHF vs DPO tradeoffs for local models.",
            "Provide robust guardrails for tool execution in an agent loop.",
        ]

    rows: list[dict] = []
    for prompt in prompts[:50]:
        p = str(prompt or "").strip()
        if not p:
            continue
        teacher_answer = ""
        try:
            teacher_messages = [{"role": "user", "content": p}]
            answer, _used_provider = call_llm_with_fallback(teacher_messages, "distillation", provider)
            if isinstance(answer, dict):
                teacher_answer = str(answer.get("content") or answer.get("result") or "").strip()
            else:
                teacher_answer = str(answer or "").strip()
        except Exception:
            teacher_answer = ""
        if not teacher_answer:
            teacher_answer = f"Teacher synthesis for: {p}"
        rows.append({"prompt": p, "response": teacher_answer, "source": "distillation", "teacher_model": teacher})

    if not rows:
        update_distill_job(job_id, status="failed", error={"message": "no distillation rows generated"})
        append_distill_job_event(job_id, "No rows generated", level="error")
        return

    bundle = _create_finetune_job_from_rows(
        model=student,
        rows=rows,
        source="distillation",
        provenance_extra={"teacher_model": teacher, "provider": provider, "distill_job_id": job_id},
    )
    ft_job_id = str(bundle.get("job", {}).get("id") or "")
    append_distill_job_event(job_id, "Fine-tune child job created", data={"fine_tune_job_id": ft_job_id})

    timeout_s = max(30, min(int(cfg.get("timeout_seconds") or 600), 3600))
    started = time.time()
    while time.time() - started <= timeout_s:
        current = get_distill_job(job_id)
        if not current or str(current.get("status") or "") == "cancelled":
            append_distill_job_event(job_id, "Distillation cancelled")
            return
        ft_job = db_get_fine_tuning_job(ft_job_id)
        if ft_job and str(ft_job.get("status") or "") in {"succeeded", "failed", "cancelled"}:
            break
        time.sleep(0.4)

    ft_job = db_get_fine_tuning_job(ft_job_id)
    status = str(ft_job.get("status") or "") if isinstance(ft_job, dict) else "failed"
    if status != "succeeded":
        update_distill_job(job_id, status="failed", error={"message": f"child fine-tune ended with {status}"})
        append_distill_job_event(job_id, "Child fine-tune failed", level="error", data={"status": status})
        return

    adapter_id = f"distill-{student.replace('/', '-') }"
    version = f"v{int(time.time())}"
    adapter = save_lora_adapter_version(
        adapter_id=adapter_id,
        version=version,
        base_model=student,
        checkpoint_uri=f"inline://adapters/{adapter_id}/{version}",
        metrics={"trained_tokens": int(ft_job.get("trained_tokens") or 0)},
        provenance={
            "distill_job_id": job_id,
            "teacher_model": teacher,
            "dataset_version_id": str(bundle.get("dataset_version", {}).get("dataset_id") or ""),
        },
        tags=["distillation", "student"],
        status="ready",
    )
    result = {
        "fine_tune_job_id": ft_job_id,
        "dataset_version": bundle.get("dataset_version", {}),
        "adapter": adapter,
    }
    update_distill_job(job_id, status="succeeded", result=result, error=None)
    append_distill_job_event(job_id, "Distillation completed", data=result)


def _nexus_prime_alpha_wire_key() -> str:
    return "finetune.persona.nexus_prime_alpha.v1"


@router.get("/models/{model_id}/card")
def get_model_card(model_id: str, markdown: bool = True):
    from ..eval_pipeline import generate_model_card, list_eval_jobs

    jobs = list_eval_jobs()
    card_md = generate_model_card(model=model_id, eval_results=jobs)
    if markdown:
        return {"model": model_id, "model_card": card_md, "format": "markdown"}

    rows = []
    for j in jobs:
        if str(j.model) != model_id:
            continue
        rows.append(
            {
                "task_id": j.task_id,
                "suite": j.suite,
                "score": j.score,
                "status": j.status,
                "regression": j.regression,
                "adapter_id": j.adapter_id,
            }
        )
    return {"model": model_id, "format": "json", "evaluations": rows}


@router.get("/models/{model_id}/transparency")
def get_model_transparency_report(model_id: str):
    from ..db import list_adapter_proof_reports, list_ft_dataset_versions, list_lora_adapter_versions
    from ..eval_pipeline import list_eval_jobs

    card = get_model_card(model_id=model_id, markdown=True)
    eval_rows = [
        {
            "task_id": j.task_id,
            "suite": j.suite,
            "score": j.score,
            "regression": j.regression,
            "created_at": j.created_at,
        }
        for j in list_eval_jobs()
        if str(j.model) == model_id
    ]
    datasets = [
        d
        for d in list_ft_dataset_versions(limit=500)
        if model_id in json.dumps(d.get("provenance") or {})
        or str(d.get("source") or "") in {"feedback_trace", "synthetic"}
    ]
    adapters = [a for a in list_lora_adapter_versions() if str(a.get("base_model") or "") == model_id]
    proof_reports = []
    for adapter in adapters[:100]:
        proof_reports.extend(
            list_adapter_proof_reports(
                adapter_id=str(adapter.get("adapter_id") or ""),
                adapter_version=str(adapter.get("version") or ""),
                limit=5,
            )
        )

    return {
        "model": model_id,
        "generated_at": datetime.now(timezone.utc).isoformat().replace("+00:00", "Z"),
        "model_card": card.get("model_card", ""),
        "evaluation_summary": {
            "count": len(eval_rows),
            "rows": eval_rows[:100],
        },
        "training_data_provenance": {
            "dataset_versions": datasets[:100],
            "count": len(datasets),
        },
        "adapter_lineage": {
            "versions": adapters[:100],
            "count": len(adapters),
        },
        "adapter_proof_reports": {
            "reports": proof_reports[:100],
            "count": len(proof_reports),
        },
        "limitations": [
            "Transparency report is derived from local registries and may omit external trainer state.",
            "Adapter application at runtime depends on provider backend support.",
        ],
    }


def _run_rlhf_dpo_job(job_id: str):
    from ..db import (
        append_rlhf_dpo_job_event,
        get_ft_dataset_version,
        get_rlhf_dpo_job,
        save_lora_adapter_version,
        update_rlhf_dpo_job,
    )
    from ..eval_pipeline import score_response
    from ..rlhf_dpo import (
        create_dpo_job as _create_native_dpo_job,
        create_rlhf_job as _create_native_rlhf_job,
        run_dpo_training as _run_native_dpo_training,
        run_rlhf_training as _run_native_rlhf_training,
    )

    def _build_preference_pair(row: dict) -> dict | None:
        prompt = str(row.get("prompt") or row.get("instruction") or "").strip()
        if not prompt:
            return None
        chosen = str(row.get("chosen") or row.get("response") or row.get("output") or "").strip()
        if not chosen:
            return None
        rejected = str(row.get("rejected") or "").strip()
        synthetic_negative = False
        if not rejected:
            synthetic_negative = True
            words = chosen.split()
            if len(words) > 6:
                rejected = " ".join(words[: max(3, len(words) // 3)])
            else:
                rejected = "Insufficient or incomplete answer."
        return {
            "prompt": prompt,
            "chosen": chosen,
            "rejected": rejected,
            "synthetic_negative": synthetic_negative,
        }

    def _score_preference_pairs(pairs: list[dict]) -> dict:
        scored_rows: list[dict] = []
        chosen_total = 0.0
        rejected_total = 0.0
        margin_total = 0.0
        synthetic_negative_count = 0
        for idx, pair in enumerate(pairs, start=1):
            prompt = str(pair.get("prompt") or "")
            chosen = str(pair.get("chosen") or "")
            rejected = str(pair.get("rejected") or "")
            synthetic_negative = bool(pair.get("synthetic_negative"))
            if synthetic_negative:
                synthetic_negative_count += 1
            chosen_score = float(score_response(prompt, chosen, reference=chosen).get("score") or 0.0)
            rejected_score = float(score_response(prompt, rejected, reference=chosen).get("score") or 0.0)
            margin = round(chosen_score - rejected_score, 6)
            chosen_total += chosen_score
            rejected_total += rejected_score
            margin_total += margin
            scored_rows.append(
                {
                    "pair_id": idx,
                    "prompt": prompt,
                    "chosen_score": round(chosen_score, 6),
                    "rejected_score": round(rejected_score, 6),
                    "margin": margin,
                    "synthetic_negative": synthetic_negative,
                    "passed": margin > 0,
                }
            )

        pair_count = len(scored_rows)
        return {
            "pair_count": pair_count,
            "synthetic_negative_count": synthetic_negative_count,
            "chosen_avg": round(chosen_total / max(1, pair_count), 6),
            "rejected_avg": round(rejected_total / max(1, pair_count), 6),
            "preference_alignment_score": round(margin_total / max(1, pair_count), 6),
            "rows": scored_rows,
        }

    job = get_rlhf_dpo_job(job_id)
    if not job:
        return
    if str(job.get("status") or "") != "queued":
        return

    append_rlhf_dpo_job_event(job_id, "RLHF/DPO job started")
    update_rlhf_dpo_job(job_id, status="running", error=None)

    latest = get_rlhf_dpo_job(job_id)
    if not latest:
        return
    if str(latest.get("status") or "") == "cancelled":
        append_rlhf_dpo_job_event(job_id, "RLHF/DPO job cancelled before execution")
        return

    method = str(latest.get("method") or "dpo")
    base_model = str(latest.get("base_model") or "nexus-prime-base")
    dataset_version_id = str(latest.get("dataset_version_id") or "")
    config = latest.get("config") if isinstance(latest.get("config"), dict) else {}
    training_backend = str(config.get("training_backend") or "orchestration").strip().lower()
    if training_backend not in {"orchestration", "native"}:
        training_backend = "orchestration"
    telemetry_gates = {
        "min_pair_count": max(1, int(config.get("gate_min_pair_count") or 8)),
        "min_alignment_score": float(config.get("gate_min_alignment_score") or 0.02),
        "max_synthetic_negative_ratio": float(config.get("gate_max_synthetic_negative_ratio") or 0.75),
    }
    dataset = get_ft_dataset_version(dataset_version_id)
    if dataset is None:
        update_rlhf_dpo_job(job_id, status="failed", error={"message": "dataset_version_id not found"})
        append_rlhf_dpo_job_event(job_id, "Dataset not found", level="error", data={"dataset_version_id": dataset_version_id})
        return

    preview_rows = dataset.get("preview_rows") if isinstance(dataset.get("preview_rows"), list) else []
    preference_pairs = []
    for row in preview_rows[:200]:
        if isinstance(row, dict):
            pair = _build_preference_pair(row)
            if pair is not None:
                preference_pairs.append(pair)
    if not preference_pairs:
        update_rlhf_dpo_job(job_id, status="failed", error={"message": "dataset_version_id contains no usable preference rows"})
        append_rlhf_dpo_job_event(job_id, "No usable preference rows found", level="error", data={"dataset_version_id": dataset_version_id})
        return

    preference_metrics = _score_preference_pairs(preference_pairs)
    synthetic_ratio = float(preference_metrics["synthetic_negative_count"]) / max(1, float(preference_metrics["pair_count"]))
    gate_fail_reasons: list[str] = []
    if int(preference_metrics["pair_count"]) < int(telemetry_gates["min_pair_count"]):
        gate_fail_reasons.append(
            f"pair_count={preference_metrics['pair_count']} < min_pair_count={telemetry_gates['min_pair_count']}"
        )
    if float(preference_metrics["preference_alignment_score"]) < float(telemetry_gates["min_alignment_score"]):
        gate_fail_reasons.append(
            "preference_alignment_score="
            f"{preference_metrics['preference_alignment_score']} < min_alignment_score={telemetry_gates['min_alignment_score']}"
        )
    if synthetic_ratio > float(telemetry_gates["max_synthetic_negative_ratio"]):
        gate_fail_reasons.append(
            f"synthetic_negative_ratio={round(synthetic_ratio, 6)} > "
            f"max_synthetic_negative_ratio={telemetry_gates['max_synthetic_negative_ratio']}"
        )

    append_rlhf_dpo_job_event(
        job_id,
        "RLHF/DPO telemetry gates evaluated",
        data={
            "training_backend": training_backend,
            "telemetry_gates": telemetry_gates,
            "pair_count": preference_metrics["pair_count"],
            "preference_alignment_score": preference_metrics["preference_alignment_score"],
            "synthetic_negative_ratio": round(synthetic_ratio, 6),
            "gate_fail_reasons": gate_fail_reasons,
        },
    )

    if gate_fail_reasons:
        update_rlhf_dpo_job(
            job_id,
            status="failed",
            error={"message": "telemetry gates failed", "reasons": gate_fail_reasons},
            result={
                "training_backend": training_backend,
                "telemetry_gates": telemetry_gates,
                "pair_count": preference_metrics["pair_count"],
                "preference_alignment_score": preference_metrics["preference_alignment_score"],
                "synthetic_negative_ratio": round(synthetic_ratio, 6),
            },
        )
        append_rlhf_dpo_job_event(job_id, "RLHF/DPO telemetry gates failed", level="error", data={"reasons": gate_fail_reasons})
        return

    native_training_result: dict | None = None

    if training_backend == "native":
        import tempfile as _tempfile

        tmp_ds = os.path.join(_tempfile.gettempdir(), f"nexus_rlhf_dpo_{job_id}.jsonl")
        with open(tmp_ds, "w", encoding="utf-8") as f:
            for pair in preference_pairs:
                f.write(
                    json.dumps(
                        {
                            "prompt": pair["prompt"],
                            "chosen": pair["chosen"],
                            "rejected": pair["rejected"],
                            "margin": max(0.0, float(preference_metrics["preference_alignment_score"])),
                            "source": "rlhf_dpo_experiment",
                        },
                        ensure_ascii=False,
                    )
                    + "\n"
                )

        append_rlhf_dpo_job_event(
            job_id,
            "Native RLHF/DPO backend selected",
            data={"dataset_path": str(tmp_ds), "method": method},
        )

        if method == "dpo":
            native_job = _create_native_dpo_job(
                base_model=base_model,
                dataset_path=str(tmp_ds),
                adapter_name=f"{base_model.replace('/', '-')}-dpo-native",
                config=config,
            )
            _run_native_dpo_training(native_job)
            if native_job.status != "completed":
                update_rlhf_dpo_job(
                    job_id,
                    status="failed",
                    error={"message": f"native DPO backend failed: {native_job.error or native_job.status}"},
                )
                append_rlhf_dpo_job_event(job_id, "Native DPO backend failed", level="error", data={"job_id": native_job.job_id, "error": native_job.error})
                return
            native_training_result = {
                "native_job_id": native_job.job_id,
                "native_backend": "dpo",
                "adapter_path": native_job.adapter_path,
                "metrics": native_job.metrics,
            }
        else:
            native_job = _create_native_rlhf_job(
                base_model=base_model,
                dataset_path=str(tmp_ds),
                adapter_name=f"{base_model.replace('/', '-')}-rlhf-native",
                config=config,
            )
            _run_native_rlhf_training(native_job)
            if native_job.status != "completed":
                update_rlhf_dpo_job(
                    job_id,
                    status="failed",
                    error={"message": f"native RLHF backend failed: {native_job.error or native_job.status}"},
                )
                append_rlhf_dpo_job_event(job_id, "Native RLHF backend failed", level="error", data={"job_id": native_job.job_id, "error": native_job.error})
                return
            native_training_result = {
                "native_job_id": native_job.job_id,
                "native_backend": "rlhf",
                "adapter_path": native_job.adapter_path,
                "metrics": native_job.metrics,
                "rounds_completed": native_job.rounds_completed,
            }

    if training_backend == "orchestration":

        child = _create_finetune_job_from_rows(
        model=base_model,
        rows=[
            {
                "prompt": pair["prompt"],
                "response": pair["chosen"],
                "source": "rlhf_dpo",
                "preference_rejected": pair["rejected"],
                "synthetic_negative": pair["synthetic_negative"],
            }
            for pair in preference_pairs
        ],
        source="rlhf_dpo",
        provenance_extra={
            "method": method,
            "dataset_version_id": dataset_version_id,
            "rlhf_dpo_job_id": job_id,
            "config": latest.get("config") if isinstance(latest.get("config"), dict) else {},
            "pair_count": preference_metrics["pair_count"],
            "synthetic_negative_count": preference_metrics["synthetic_negative_count"],
        },
    )
        fine_tune_job_id = str(child.get("job", {}).get("id") or "")
        append_rlhf_dpo_job_event(job_id, "Fine-tune child job created", data={"fine_tune_job_id": fine_tune_job_id})

        timeout_seconds = 900
        started_at = time.time()
        while time.time() - started_at <= timeout_seconds:
            current = get_rlhf_dpo_job(job_id)
            if not current or str(current.get("status") or "") == "cancelled":
                append_rlhf_dpo_job_event(job_id, "RLHF/DPO job cancelled")
                return
            ft_job = db_get_fine_tuning_job(fine_tune_job_id)
            if ft_job and str(ft_job.get("status") or "") in {"succeeded", "failed", "cancelled"}:
                break
            time.sleep(0.4)

        ft_job = db_get_fine_tuning_job(fine_tune_job_id)
        if not ft_job:
            update_rlhf_dpo_job(job_id, status="failed", error={"message": "fine-tune child job missing"})
            append_rlhf_dpo_job_event(job_id, "Child fine-tune job missing", level="error")
            return
        ft_status = str(ft_job.get("status") or "")
        if ft_status != "succeeded":
            update_rlhf_dpo_job(
                job_id,
                status="failed",
                error={"message": f"fine-tune child job ended with status={ft_status}"},
                result={"fine_tune_job_id": fine_tune_job_id, "dataset_version_id": dataset_version_id},
            )
            append_rlhf_dpo_job_event(job_id, "Child fine-tune failed", level="error", data={"status": ft_status})
            return
    else:
        fine_tune_job_id = ""
        ft_job = {"trained_tokens": 0}

    adapter_id = f"{base_model.replace('/', '-')}-{method}"
    adapter_version = f"v{int(time.time())}"
    adapter = save_lora_adapter_version(
        adapter_id=adapter_id,
        version=adapter_version,
        base_model=base_model,
        checkpoint_uri=f"inline://adapters/{adapter_id}/{adapter_version}",
        metrics={
            "preference_alignment_score": float(preference_metrics["preference_alignment_score"]),
            "preference_pair_count": int(preference_metrics["pair_count"]),
            "chosen_avg": float(preference_metrics["chosen_avg"]),
            "rejected_avg": float(preference_metrics["rejected_avg"]),
            "trained_tokens": int(ft_job.get("trained_tokens") or 0),
        },
        provenance={
            "method": method,
            "dataset_version_id": dataset_version_id,
            "rlhf_dpo_job_id": job_id,
            "fine_tune_job_id": fine_tune_job_id,
        },
        tags=["rlhf", method],
        status="ready",
    )
    result = {
        "preference_alignment_score": float(preference_metrics["preference_alignment_score"]),
        "pair_count": int(preference_metrics["pair_count"]),
        "chosen_avg": float(preference_metrics["chosen_avg"]),
        "rejected_avg": float(preference_metrics["rejected_avg"]),
        "synthetic_negative_count": int(preference_metrics["synthetic_negative_count"]),
        "preference_rows": preference_metrics["rows"][:25],
        "method": method,
        "training_backend": training_backend,
        "telemetry_gates": telemetry_gates,
        "synthetic_negative_ratio": round(synthetic_ratio, 6),
        "fine_tune_job_id": fine_tune_job_id,
        "dataset_version_id": dataset_version_id,
        "adapter": adapter,
    }
    if native_training_result is not None:
        result["native_training"] = native_training_result
    update_rlhf_dpo_job(
        job_id,
        status="succeeded",
        result=result,
        error=None,
    )
    append_rlhf_dpo_job_event(job_id, "RLHF/DPO job completed", data=result)


def _continual_policies_key() -> str:
    return "finetune.continual.policies.v1"


def _update_continual_policy(policy_id: str, updates: dict) -> dict | None:
    rows = _load_continual_policies()
    updated = None
    for idx, row in enumerate(rows):
        if str(row.get("id") or "") != str(policy_id):
            continue
        merged = dict(row)
        merged.update(dict(updates or {}))
        merged["updated_at"] = datetime.now(timezone.utc).isoformat().replace("+00:00", "Z")
        rows[idx] = merged
        updated = merged
        break
    if updated is not None:
        _save_continual_policies(rows)
    return updated


def _execute_continual_re_tune_task(payload: dict) -> str:
    from ..eval_pipeline import run_eval_suite

    policy_id = str(payload.get("policy_id") or "").strip()
    model = str(payload.get("model") or "nexus-prime-base").strip() or "nexus-prime-base"
    threshold = float(payload.get("threshold") or 0.05)
    suites = payload.get("suites") if isinstance(payload.get("suites"), list) else ["code", "autonomy", "rag"]
    n_samples = max(2, min(int(payload.get("n_samples") or 8), 64))
    provider = str(payload.get("provider") or "ollama")
    include_trace = bool(payload.get("include_trace", True))

    batch = run_eval_suite(model=model, provider=provider, suites=suites, n_samples=n_samples)
    avg_score = float(batch.get("average_score") or 0.0)
    has_regression = bool(batch.get("has_regression"))

    current_policy = None
    if policy_id:
        for row in _load_continual_policies():
            if str(row.get("id") or "") == policy_id:
                current_policy = row
                break
    prev_score = float(current_policy.get("last_average_score") or 0.0) if current_policy else 0.0
    delta = avg_score - prev_score

    updates = {
        "last_run_at": datetime.now(timezone.utc).isoformat().replace("+00:00", "Z"),
        "last_average_score": round(avg_score, 6),
        "last_delta": round(delta, 6),
        "run_count": int((current_policy or {}).get("run_count") or 0) + 1,
        "last_eval": {
            "suites": suites,
            "average_score": round(avg_score, 6),
            "has_regression": has_regression,
        },
    }
    if policy_id:
        _update_continual_policy(policy_id, updates)

    should_retune = (delta >= threshold) and (not has_regression)
    if not should_retune:
        return f"continual policy checked: avg={avg_score:.3f}, delta={delta:.3f}, threshold={threshold:.3f}, retune=no"

    rows = _training_rows_from_feedback(include_trace=include_trace, limit=5000)
    if not rows:
        return f"continual policy checked: avg={avg_score:.3f}, delta={delta:.3f}, retune_skipped=no_rows"

    try:
        bundle = _create_finetune_job_from_rows(
            model=model,
            rows=rows,
            source="continual_feedback_trace",
            provenance_extra={
                "policy_id": policy_id,
                "trigger_average_score": avg_score,
                "trigger_delta": delta,
                "include_trace": include_trace,
            },
        )
    except Exception as exc:
        err = str(exc)
        lowered = err.lower()
        transient_markers = ("cuda", "gpu", "out of memory", "nvidia", "resource", "capacity")
        if any(m in lowered for m in transient_markers):
            if policy_id:
                _update_continual_policy(
                    policy_id,
                    {
                        "last_trigger_error": err[:500],
                        "last_trigger_status": "deferred_capacity",
                    },
                )
            return (
                f"continual policy deferred: avg={avg_score:.3f}, delta={delta:.3f}, "
                f"reason={err[:160]}"
            )
        raise
    if policy_id:
        _update_continual_policy(
            policy_id,
            {
                "last_triggered_finetune_job_id": str(bundle.get("job", {}).get("id") or ""),
                "last_triggered_dataset_id": str(bundle.get("dataset_version", {}).get("dataset_id") or ""),
                "retune_count": int((current_policy or {}).get("retune_count") or 0) + 1,
                "last_trigger_status": "triggered",
            },
        )
    return (
        f"continual policy triggered: avg={avg_score:.3f}, delta={delta:.3f}, "
        f"fine_tune_job_id={bundle.get('job', {}).get('id', '')}"
    )


def _load_continual_policies() -> list[dict]:
    raw = db_load_pref(_continual_policies_key(), "[]")
    try:
        data = json.loads(raw)
        return data if isinstance(data, list) else []
    except Exception:
        return []


def _save_continual_policies(rows: list[dict]) -> None:
    db_save_pref(_continual_policies_key(), json.dumps(rows[-5000:]))



# ─────────────────────────────────────────────────────────────────────────────
# Audio ingestion  — podcast / meeting transcript pipeline
# ─────────────────────────────────────────────────────────────────────────────


# Ollama management endpoints
# ─────────────────────────────────────────────────────────────────────────────

def _ollama_base() -> str:
    return os.getenv("OLLAMA_BASE_URL", "http://localhost:11434/v1").rstrip("/").removesuffix("/v1")


@router.post("/ollama/pull")
async def ollama_pull_model(request: Request):
    """Pull (download) an Ollama model on demand.
    POST body: {"model": "llama3.2:3b", "stream": false}
    """
    import requests as _req
    try:
        body = await request.json()
    except Exception:
        return _api_error("invalid JSON body", "validation_error", 400)

    model = str(body.get("model", "")).strip()
    if not model:
        return _api_error("model is required", "validation_error", 422)

    stream = bool(body.get("stream", False))
    base = _ollama_base()

    if stream:
        def _pull_stream():
            try:
                r = _req.post(f"{base}/api/pull",
                              json={"name": model, "stream": True},
                              stream=True, timeout=600)
                r.raise_for_status()
                for line in r.iter_lines():
                    if line:
                        yield line.decode() + "\n"
            except Exception as exc:
                yield json.dumps({"error": str(exc)}) + "\n"

        return StreamingResponse(_pull_stream(), media_type="application/x-ndjson")

    try:
        r = _req.post(f"{base}/api/pull", json={"name": model, "stream": False}, timeout=600)
        r.raise_for_status()
        return {"ok": True, "model": model, "status": r.json().get("status", "success")}
    except Exception as exc:
        return _api_error(f"Ollama pull failed: {exc}", "model_error", 502)


@router.post("/ollama/benchmark")
async def ollama_benchmark(request: Request):
    """Run a latency benchmark against one or more local Ollama models.
    POST body: {"models": ["llama3.2:3b"], "prompt": "Say hello.", "runs": 3}
    """
    from ..benchmark import run_ollama_benchmark
    try:
        body = await request.json()
    except Exception:
        return _api_error("invalid JSON body", "validation_error", 400)
    result = run_ollama_benchmark(
        models=body.get("models") or None,
        prompt=str(body.get("prompt", "Say hello in one sentence.")),
        runs=int(body.get("runs", 3)),
        ollama_base=_ollama_base(),
    )
    if "error" in result and "models" not in result:
        return _api_error(result["error"], "validation_error", 422)
    return result


# ── GGUF model file management ────────────────────────────────────────────────

_GGUF_DIR = os.path.join(os.getenv("DATA_DIR", "/data"), "models")


def _ensure_gguf_dir():
    os.makedirs(_GGUF_DIR, exist_ok=True)


@router.get("/ollama/gguf")
def ollama_list_gguf():
    """List GGUF model files stored in the data/models directory."""
    _ensure_gguf_dir()
    files = []
    for name in os.listdir(_GGUF_DIR):
        if name.lower().endswith(".gguf"):
            path = os.path.join(_GGUF_DIR, name)
            stat = os.stat(path)
            files.append({
                "filename": name,
                "size_bytes": stat.st_size,
                "modified_at": int(stat.st_mtime),
                "path": path,
            })
    return {"object": "list", "data": sorted(files, key=lambda f: f["filename"])}


@router.post("/ollama/gguf/import")
async def ollama_import_gguf(request: Request):
    """Import a GGUF file into Ollama via `ollama create`.
    POST body: {"filename": "model.gguf", "name": "my-model:latest"}
    The file must already be in the data/models directory.
    """
    import subprocess as _sp
    try:
        body = await request.json()
    except Exception:
        return _api_error("invalid JSON body", "validation_error", 400)

    filename = str(body.get("filename", "")).strip()
    model_name = str(body.get("name", "")).strip()
    if not filename or not model_name:
        return _api_error("filename and name are required", "validation_error", 422)

    gguf_path = os.path.join(_GGUF_DIR, filename)
    if not os.path.exists(gguf_path) or not filename.lower().endswith(".gguf"):
        return _api_error(f"GGUF file '{filename}' not found in data/models/", "not_found_error", 404)

    # Write a minimal Modelfile
    modelfile_content = f"FROM {gguf_path}\n"
    modelfile_path = gguf_path + ".Modelfile"
    try:
        with open(modelfile_path, "w") as mf:
            mf.write(modelfile_content)
        result = _sp.run(
            ["ollama", "create", model_name, "-f", modelfile_path],
            capture_output=True, text=True, timeout=300,
        )
        if result.returncode != 0:
            return _api_error(f"ollama create failed: {result.stderr}", "model_error", 502)
        return {"ok": True, "model": model_name, "gguf": filename, "output": result.stdout.strip()}
    except FileNotFoundError:
        return _api_error("ollama CLI not found in PATH", "model_error", 503)
    except Exception as exc:
        return _api_error(str(exc), "model_error", 500)
    finally:
        try:
            os.unlink(modelfile_path)
        except Exception:
            pass


@router.delete("/ollama/gguf/{filename}")
def ollama_delete_gguf(filename: str):
    """Delete a GGUF file from the data/models directory."""
    if ".." in filename or "/" in filename:
        return _api_error("invalid filename", "invalid_request_error", 400)
    path = os.path.join(_GGUF_DIR, filename)
    if not os.path.exists(path):
        return _api_error(f"GGUF file '{filename}' not found", "not_found_error", 404)
    os.unlink(path)
    return {"ok": True, "deleted": filename}


# ─────────────────────────────────────────────────────────────────────────────
# HuggingFace download + Ollama import pipeline
# ─────────────────────────────────────────────────────────────────────────────

@router.post("/huggingface/download")
async def huggingface_download(request: Request):
    """Download a GGUF model from HuggingFace Hub and optionally import it into Ollama.

    POST body:
      {
        "repo_id":   "bartowski/Llama-3.2-3B-Instruct-GGUF",
        "filename":  "Llama-3.2-3B-Instruct-Q4_K_M.gguf",
        "ollama_name": "llama3.2-3b:q4"   // optional — if set, auto-imports after download
      }

    Requires huggingface_hub (`pip install huggingface_hub`) or wget fallback.
    """
    import subprocess as _sp
    try:
        body = await request.json()
    except Exception:
        return _api_error("invalid JSON body", "validation_error", 400)

    repo_id = str(body.get("repo_id", "")).strip()
    filename = str(body.get("filename", "")).strip()
    ollama_name = str(body.get("ollama_name", "")).strip()

    if not repo_id or not filename:
        return _api_error("repo_id and filename are required", "validation_error", 422)
    if not filename.lower().endswith(".gguf"):
        return _api_error("filename must be a .gguf file", "validation_error", 422)
    if ".." in filename or "/" in filename:
        return _api_error("invalid filename", "invalid_request_error", 400)

    _ensure_gguf_dir()
    dest_path = os.path.join(_GGUF_DIR, filename)

    # Try huggingface_hub first
    try:
        from huggingface_hub import hf_hub_download  # type: ignore
        hf_token = os.getenv("HF_TOKEN", "") or None
        hf_hub_download(
            repo_id=repo_id,
            filename=filename,
            local_dir=_GGUF_DIR,
            token=hf_token,
        )
    except ImportError:
        # Fallback: direct HTTPS download via requests
        import requests as _req
        hf_token = os.getenv("HF_TOKEN", "")
        url = f"https://huggingface.co/{repo_id}/resolve/main/{filename}"
        headers = {"Authorization": f"Bearer {hf_token}"} if hf_token else {}
        try:
            with _req.get(url, headers=headers, stream=True, timeout=30) as r:
                r.raise_for_status()
                total = int(r.headers.get("content-length", 0))
                downloaded = 0
                with open(dest_path, "wb") as fh:
                    for chunk in r.iter_content(chunk_size=8192):
                        fh.write(chunk)
                        downloaded += len(chunk)
        except Exception as exc:
            return _api_error(f"Download failed: {exc}", "model_error", 502)
    except Exception as exc:
        return _api_error(f"HuggingFace download failed: {exc}", "model_error", 502)

    result: dict = {
        "ok": True,
        "repo_id": repo_id,
        "filename": filename,
        "path": dest_path,
        "size_bytes": os.path.getsize(dest_path),
    }

    # Auto-import into Ollama if ollama_name is provided
    if ollama_name:
        modelfile_path = dest_path + ".Modelfile"
        try:
            with open(modelfile_path, "w") as mf:
                mf.write(f"FROM {dest_path}\n")
            proc = _sp.run(
                ["ollama", "create", ollama_name, "-f", modelfile_path],
                capture_output=True, text=True, timeout=300,
            )
            if proc.returncode == 0:
                result["ollama_import"] = {"ok": True, "model": ollama_name}
            else:
                result["ollama_import"] = {"ok": False, "error": proc.stderr.strip()}
        except FileNotFoundError:
            result["ollama_import"] = {"ok": False, "error": "ollama CLI not found in PATH"}
        except Exception as exc:
            result["ollama_import"] = {"ok": False, "error": str(exc)}
        finally:
            try:
                os.unlink(modelfile_path)
            except Exception:
                pass

    return result


# ── Autonomy SSE streaming ─────────────────────────────────────────────────────

@router.post("/autonomy/execute/stream")
async def autonomy_execute_stream(request: Request):
    """Stream autonomy execution events via SSE."""
    body = await request.json()
    goal         = str(body.get("goal", "")).strip()
    strategy     = str(body.get("strategy", "parallel"))
    max_subtasks = int(body.get("max_subtasks", 6))
    sid          = str(body.get("sid", ""))
    trace_id     = secrets.token_hex(8)

    if not goal:
        return _api_error("'goal' is required", "invalid_request_error", 400)

    import queue as _queue
    event_queue: "_queue.Queue[dict | None]" = _queue.Queue()
    checkpoint_events: List[Dict[str, Any]] = []

    def _run_autonomy():
        try:
            orchestrator = Orchestrator(llm=_orchestrator_llm)
            start_evt = {"type": "autonomy_start", "goal": goal, "trace_id": trace_id}
            checkpoint_events.append(start_evt)
            _save_autonomy_checkpoint(trace_id, 0, goal, checkpoint_events)
            event_queue.put(start_evt)
            result = orchestrator.execute(goal, context={"strategy": strategy, "max_subtasks": max_subtasks, "sid": sid})
            result["trace_id"] = trace_id
            for idx, subtask in enumerate(result.get("subtasks", []), 1):
                evt = {"type": "subtask_done", "trace_id": trace_id, "subtask": subtask}
                checkpoint_events.append(evt)
                _save_autonomy_checkpoint(trace_id, idx, goal, checkpoint_events)
                event_queue.put(evt)
            done_evt = {
                "type": "autonomy_done",
                "trace_id": trace_id,
                "result": result.get("result", ""),
                "plan_summary": result.get("plan_summary", ""),
                "execution_time": result.get("execution_time", 0),
            }
            checkpoint_events.append(done_evt)
            _save_autonomy_checkpoint(trace_id, len(checkpoint_events), goal, checkpoint_events)
            autonomy_traces[trace_id] = {"type": "execution_stream", "goal": goal, "status": "done", **result}
            db_save_autonomy_trace(trace_id, autonomy_traces[trace_id])
            event_queue.put(done_evt)
        except Exception as exc:
            err_evt = {"type": "error", "trace_id": trace_id, "error": str(exc)}
            checkpoint_events.append(err_evt)
            _save_autonomy_checkpoint(trace_id, len(checkpoint_events), goal, checkpoint_events)
            event_queue.put(err_evt)
        finally:
            event_queue.put(None)  # sentinel

    threading.Thread(target=_run_autonomy, daemon=True).start()

    async def _gen():
        loop = asyncio.get_event_loop()
        while True:
            event = await loop.run_in_executor(None, event_queue.get)
            if event is None:
                yield "data: [DONE]\n\n"
                break
            yield f"data: {json.dumps(event)}\n\n"

    return StreamingResponse(
        _gen(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no", "X-Trace-Id": trace_id},
    )


# ── Task queue routes ─────────────────────────────────────────────────────────

@router.post("/tasks/queue")
async def task_queue_submit(request: Request):
    """Submit a task to the priority queue."""
    from ..task_queue import submit_task, start_worker
    body          = await request.json()
    description   = str(body.get("description", "")).strip()
    priority      = int(body.get("priority", 5))
    deps          = body.get("deps", [])
    metadata      = body.get("metadata", {})
    schedule_cron = str(body.get("schedule_cron", ""))
    dedupe        = bool(body.get("dedupe", True))

    if not description:
        return _api_error("'description' is required", "invalid_request_error", 400)

    start_worker()
    task_id = submit_task(
        description=description,
        priority=priority,
        deps=list(deps) if isinstance(deps, list) else [],
        metadata=dict(metadata) if isinstance(metadata, dict) else {},
        schedule_cron=schedule_cron,
        dedupe=dedupe,
    )
    return {"task_id": task_id, "status": "pending", "dedupe": dedupe}


@router.delete("/tasks/queue/{task_id}")
async def task_queue_cancel(task_id: str):
    """Cancel a pending or running task."""
    from ..task_queue import cancel_task
    cancelled = cancel_task(task_id)
    if not cancelled:
        return _api_error(f"Task '{task_id}' not found or already terminal", "not_found_error", 404)
    return {"task_id": task_id, "cancelled": True}


@router.get("/tasks/queue")
async def task_queue_list(request: Request):
    """List tasks (optionally filtered by status)."""
    from ..task_queue import list_tasks
    status = request.query_params.get("status", "")
    limit  = int(request.query_params.get("limit", 50))
    return {"tasks": list_tasks(status=status, limit=limit)}


@router.get("/tasks/queue/dag")
async def task_queue_dag():
    """Return current DAG state (blocked/ready/running/done)."""
    from ..task_queue import get_dag_status
    return get_dag_status()


@router.get("/tasks/queue/{task_id}")
async def task_queue_get(task_id: str):
    """Get a single task by ID."""
    from ..task_queue import get_task
    task = get_task(task_id)
    if task is None:
        return _api_error(f"Task '{task_id}' not found", "not_found_error", 404)
    return task


# ── Task queue shared memory ──────────────────────────────────────────────────

@router.get("/tasks/memory")
async def task_memory_list():
    """List all cross-task shared memory entries."""
    from ..task_queue import list_shared_memory
    return {"memory": list_shared_memory()}


@router.get("/tasks/memory/{key}")
async def task_memory_get(key: str):
    """Get a single shared memory entry."""
    from ..task_queue import get_shared_memory
    value = get_shared_memory(key)
    if value is None:
        return _api_error(f"Key '{key}' not found", "not_found_error", 404)
    return {"key": key, "value": value}


@router.put("/tasks/memory/{key}")
async def task_memory_set(key: str, request: Request):
    """Set a shared memory entry."""
    from ..task_queue import set_shared_memory
    body  = await request.json()
    value = body.get("value")
    set_shared_memory(key, value)
    return {"key": key, "value": value}


@router.delete("/tasks/memory/{key}")
async def task_memory_delete(key: str):
    """Delete a shared memory entry."""
    from ..task_queue import delete_shared_memory
    deleted = delete_shared_memory(key)
    if not deleted:
        return _api_error(f"Key '{key}' not found", "not_found_error", 404)
    return {"key": key, "deleted": True}


# ── Task worker control ───────────────────────────────────────────────────────

@router.get("/tasks/worker/status")
async def task_worker_status():
    """Return task worker status (running, queue depth, active tasks)."""
    from ..task_queue import worker_status
    return worker_status()


@router.post("/tasks/worker/start")
async def task_worker_start():
    """Start the background task worker."""
    from ..task_queue import start_worker
    start_worker()
    return {"started": True}


@router.post("/tasks/worker/stop")
async def task_worker_stop():
    """Stop the background task worker."""
    from ..task_queue import stop_worker
    stop_worker()
    return {"stopped": True}


# ── Simulation scenario library & comparison ─────────────────────────────────

@router.get("/simulation/scenarios")
async def simulation_scenarios():
    """List pre-built simulation scenario templates."""
    from ..simulation import SCENARIO_LIBRARY
    return {
        "scenarios": [
            {"id": k, "topic": v["topic"], "n_personas": v.get("n_personas", 5),
             "n_rounds": v.get("n_rounds", 3)}
            for k, v in SCENARIO_LIBRARY.items()
        ]
    }


@router.post("/simulation/compare")
async def simulation_compare(request: Request):
    """A/B diff two simulation results."""
    from ..simulation import compare_simulations
    body  = await request.json()
    sim_a = body.get("sim_a")
    sim_b = body.get("sim_b")
    if not sim_a or not sim_b:
        return _api_error("Both 'sim_a' and 'sim_b' are required", "invalid_request_error", 400)
    return compare_simulations(sim_a, sim_b)


@router.post("/simulation/export-training")
async def simulation_export_training(request: Request):
    """Export simulation results as a fine-tuning training dataset."""
    from ..simulation import export_training_dataset
    body    = await request.json()
    results = body.get("results", [])
    if not isinstance(results, list) or not results:
        return _api_error("'results' must be a non-empty array", "invalid_request_error", 400)
    dataset = export_training_dataset(results)
    return {"count": len(dataset), "dataset": dataset}



# =============================================================================
# Section 1 additions — Health probes, metrics, feature flags,
# MFA, org management, audit log, circuit breaker admin, cache admin
# =============================================================================

import time as _time


# ─────────────────────────────────────────────────────────────────────────────
# Enhanced health endpoints
# ─────────────────────────────────────────────────────────────────────────────

@router.get("/health/live")
def health_live():
    """Kubernetes liveness probe — always 200 while the process is alive."""
    return {"status": "alive", "ts": _time.time()}


@router.get("/health/ready")
def health_ready():
    """Kubernetes readiness probe — 200 when DB is reachable."""
    try:
        from ..db import _sql_fetchall
        _sql_fetchall("SELECT 1")
        return {"status": "ready", "ts": _time.time()}
    except Exception as exc:
        return JSONResponse({"status": "not_ready", "error": str(exc)}, status_code=503)


@router.get("/health/deep")
def health_deep():
    """Deep health check — database, Redis, vector store, and provider connectivity."""
    import urllib.request
    results: dict = {"ts": _time.time()}

    # DB check
    try:
        from ..db import _sql_fetchall
        _sql_fetchall("SELECT 1")
        results["db"] = "ok"
    except Exception as exc:
        results["db"] = f"error: {exc}"

    # Redis check
    try:
        from ..redis_state import redis_health
        results["redis"] = redis_health()
    except Exception as exc:
        results["redis"] = f"error: {exc}"

    # Vector store (ChromaDB) check
    try:
        from ..memory import _get_chroma_client
        client = _get_chroma_client()
        client.heartbeat()
        results["vector_store"] = "ok"
    except Exception:
        try:
            # Fallback: try instantiating the RAG vector store directly
            from ..rag.rag_system import get_rag_system
            rag = get_rag_system()
            vs = getattr(rag, "vector_store", None)
            if vs is not None:
                results["vector_store"] = "ok"
            else:
                results["vector_store"] = "unavailable"
        except Exception as exc2:
            results["vector_store"] = f"error: {exc2}"

    # Provider reachability check (attempt a lightweight HTTP probe on the active provider)
    try:
        from ..agent import PROVIDERS, _has_key, _provider_api_key
        reachable: dict = {}
        probe_providers = {
            "ollama": ("http", str(__import__("os").getenv("OLLAMA_BASE_URL", "http://localhost:11434")) + "/api/tags"),
            "groq": ("https", "https://api.groq.com/openai/v1/models"),
            "openai": ("https", "https://api.openai.com/v1/models"),
            "gemini": ("https", "https://generativelanguage.googleapis.com/v1beta/models"),
        }
        for pid, (scheme, url) in probe_providers.items():
            cfg = PROVIDERS.get(pid)
            if not cfg or not _has_key(cfg):
                continue
            try:
                api_key = _provider_api_key(cfg)
                req = urllib.request.Request(url, method="GET")
                if api_key and pid != "ollama":
                    req.add_header("Authorization", f"Bearer {api_key}")
                req.add_header("User-Agent", "NexusAI-HealthProbe/1.0")
                with urllib.request.urlopen(req, timeout=3) as resp:
                    reachable[pid] = "ok" if resp.status < 500 else f"http_{resp.status}"
            except Exception as e:
                reachable[pid] = f"error: {type(e).__name__}"
        results["providers"] = reachable if reachable else {"note": "no configured providers probed"}
    except Exception as exc:
        results["providers"] = f"error: {exc}"

    overall = "healthy" if results.get("db") == "ok" else "degraded"
    results["status"] = overall
    status_code = 200 if overall == "healthy" else 503
    return JSONResponse(results, status_code=status_code)


# ─────────────────────────────────────────────────────────────────────────────
# Prometheus metrics
# ─────────────────────────────────────────────────────────────────────────────

@router.get("/metrics")
def prometheus_metrics(request: Request):
    """Expose Prometheus metrics text."""
    try:
        from ..observability import get_prometheus_metrics_text
        from fastapi.responses import PlainTextResponse
        return PlainTextResponse(
            get_prometheus_metrics_text(),
            media_type="text/plain; version=0.0.4; charset=utf-8",
        )
    except Exception as exc:
        return JSONResponse({"error": str(exc)}, status_code=503)


# ─────────────────────────────────────────────────────────────────────────────
# API changelog
# ─────────────────────────────────────────────────────────────────────────────

@router.get("/api/changelog")
def api_changelog():
    """Return API version changelog."""
    return {
        "versions": [
            {
                "version": "1.0.0",
                "date": "2025-01",
                "changes": ["Initial OpenAI-compatible API layer"],
            },
            {
                "version": "1.1.0",
                "date": "2025-06",
                "changes": [
                    "Added /health/live, /health/ready, /health/deep probes",
                    "Added /metrics Prometheus endpoint",
                    "Added MFA routes",
                    "Added org management routes",
                    "Added feature flag admin routes",
                    "Added circuit breaker admin routes",
                    "Added audit log routes",
                ],
            },
        ]
    }


def _get_current_user(request: Request) -> str:
    """Extract username from request auth, returns '' on failure."""
    try:
        user = require_auth(request)
        return user.get("username", "")
    except Exception:
        return ""


# ─────────────────────────────────────────────────────────────────────────────
# SAML 2.0 enterprise SSO routes
# ─────────────────────────────────────────────────────────────────────────────

def _get_saml_client(provider: str):
    """Build a pysaml2 Saml2Client for the given provider slug."""
    from saml2 import BINDING_HTTP_POST
    from saml2.client import Saml2Client
    from saml2.config import Config as Saml2Config

    idp_metadata_url = os.getenv(f"SAML_{provider.upper()}_IDP_METADATA_URL", "")
    sp_entity_id = os.getenv(f"SAML_{provider.upper()}_SP_ENTITY_ID", f"nexus-ai-{provider}")
    acs_url = os.getenv(f"SAML_{provider.upper()}_ACS_URL", f"https://localhost/auth/saml/{provider}/acs")

    if not idp_metadata_url:
        raise ValueError(f"SAML provider '{provider}' not configured (missing IDP_METADATA_URL)")

    settings = {
        "entityid": sp_entity_id,
        "service": {
            "sp": {
                "endpoints": {
                    "assertion_consumer_service": [
                        (acs_url, BINDING_HTTP_POST),
                    ],
                },
                "allow_unsolicited": True,
                "authn_requests_signed": False,
                "want_assertions_signed": True,
            }
        },
        "metadata": {"remote": [{"url": idp_metadata_url}]},
    }
    cfg = Saml2Config()
    cfg.load(settings)
    return Saml2Client(config=cfg)


# ── Swarm health and control ──────────────────────────────────────────────────

# Module-level swarm pause flag
_swarm_paused: bool = False


@router.get("/swarm/health")
def swarm_health():
    """Return a health summary of all known agents in the swarm.

    Uses the activity log to classify each agent as idle, busy, or errored
    based on recent events.
    """
    from ..agent_bus import all_agents, get_bus
    bus = get_bus()

    # Derive agent states from the last N activity events
    recent = activity_log[-200:]
    agent_state: dict[str, str] = {}
    for event in reversed(recent):
        agent_id = str(event.get("agent") or event.get("agent_id") or "")
        if not agent_id or agent_id in agent_state:
            continue
        event_type = str(event.get("type") or "").lower()
        if "error" in event_type or "fail" in event_type:
            agent_state[agent_id] = "errored"
        elif "start" in event_type or "run" in event_type:
            agent_state[agent_id] = "busy"
        else:
            agent_state[agent_id] = "idle"

    # Also include agents with messages in bus inbox
    for aid in all_agents():
        if aid not in agent_state:
            unread = bus.unread_count(aid)
            agent_state[aid] = "busy" if unread > 0 else "idle"

    counts: dict[str, int] = {"idle": 0, "busy": 0, "errored": 0, "unknown": 0}
    for state in agent_state.values():
        counts[state if state in counts else "unknown"] += 1

    return {
        "paused":  _swarm_paused,
        "agents":  [{"id": aid, "state": state} for aid, state in agent_state.items()],
        "summary": counts,
        "total":   len(agent_state),
    }


@router.post("/swarm/pause")
def swarm_pause():
    """Pause the swarm — signals agents to stop accepting new tasks."""
    global _swarm_paused
    _swarm_paused = True
    return {"paused": True, "message": "Swarm paused. Existing tasks will complete."}


@router.post("/swarm/resume")
def swarm_resume():
    """Resume the swarm after a pause."""
    global _swarm_paused
    _swarm_paused = False
    return {"paused": False, "message": "Swarm resumed."}


# ── Blueprint execution, export, and import ───────────────────────────────────

@router.post("/architecture/blueprints/{name}/execute")
async def execute_architecture_blueprint(name: str, request: Request):
    """Execute a blueprint — spawn the specialist agents it defines.

    Publishes a task message on the agent bus for each agent role in the
    blueprint's agent_layer.  Returns a list of dispatched task message IDs.

    POST body (optional):
        task     — override task description sent to each agent
        version  — specific blueprint version to execute (default: latest)
    """
    body: dict = {}
    try:
        body = await request.json()
    except Exception:
        body = {}

    version = int(body.get("version", 0)) or None
    task    = str(body.get("task") or "Execute blueprint workflow").strip()

    from ..db import load_architecture_blueprint
    blueprint = load_architecture_blueprint(name, version=version)
    if not blueprint:
        return _api_error(f"Blueprint '{name}' not found", "not_found", 404)

    if _swarm_paused:
        return _api_error("Swarm is paused — resume before executing blueprints", "swarm_paused", 409)

    from ..agent_bus import post_message
    snapshot  = blueprint.get("snapshot", {})
    agent_ids = [a.get("id") for a in snapshot.get("agent_layer", []) if a.get("id")]

    if not agent_ids:
        return _api_error("Blueprint has no agents to execute", "empty_blueprint", 422)

    dispatched = []
    for agent_id in agent_ids:
        msg = post_message(
            from_id = "blueprint_executor",
            to_id   = agent_id,
            content = task,
            topic   = f"blueprint:{name}",
        )
        dispatched.append({"agent_id": agent_id, "msg_id": msg.msg_id})

    return {
        "blueprint":  name,
        "version":    blueprint.get("version"),
        "task":       task,
        "dispatched": dispatched,
        "count":      len(dispatched),
    }


@router.get("/architecture/blueprints/{name}/export")
def export_architecture_blueprint(name: str, version: int = 0):
    """Export a blueprint as a portable JSON file (downloadable).

    Query params:
        version — specific version to export (0 = latest)
    """
    from ..db import load_architecture_blueprint
    blueprint = load_architecture_blueprint(name, version=version if version > 0 else None)
    if not blueprint:
        return _api_error(f"Blueprint '{name}' not found", "not_found", 404)

    payload   = json.dumps(blueprint, indent=2).encode()
    safe_name = "".join(c if c.isalnum() or c in "-_" else "_" for c in name)
    filename  = f"blueprint_{safe_name}_v{blueprint.get('version', 1)}.json"
    return Response(
        content    = payload,
        media_type = "application/json",
        headers    = {"Content-Disposition": f'attachment; filename="{filename}"'},
    )


# NOTE: /architecture/blueprints/import must be registered before
# /architecture/blueprints/{name} to avoid {name} capturing "import".
# (FastAPI registers routes in order, so this is handled in app.py include_router.)
@router.post("/architecture/blueprints/import", status_code=201)
async def import_architecture_blueprint(request: Request):
    """Import a previously exported blueprint JSON.

    POST body: the full blueprint JSON object (as exported by GET …/export).
    Optional override keys:
        name  — rename the blueprint on import
        notes — override notes field
    """
    body: dict = {}
    try:
        body = await request.json()
    except Exception:
        return _api_error("Invalid JSON body", "validation_error", 422)

    # Accept either the full export envelope or a raw snapshot
    if "snapshot" in body:
        name     = str(body.get("name") or "imported").strip() or "imported"
        notes    = str(body.get("notes") or "").strip()
        snapshot = body["snapshot"]
    else:
        # Treat the whole body as the snapshot
        name     = str(body.get("name") or "imported").strip() or "imported"
        notes    = ""
        snapshot = body

    if not isinstance(snapshot, dict):
        return _api_error("'snapshot' must be a JSON object", "validation_error", 422)

    from ..db import save_architecture_blueprint
    created = save_architecture_blueprint(name=name, snapshot=snapshot, notes=notes)
    return {"blueprint": created, "status": "imported"}


# ── Swarm health and control ──────────────────────────────────────────────────

# Module-level swarm pause flag
_swarm_paused: bool = False


@router.get("/swarm/health")
def swarm_health():
    """Return a health summary of all known agents in the swarm.

    Uses the activity log to classify each agent as idle, busy, or errored
    based on recent events.
    """
    from ..agent_bus import all_agents, get_bus
    bus = get_bus()

    # Derive agent states from the last N activity events
    recent = activity_log[-200:]
    agent_state: dict[str, str] = {}
    for event in reversed(recent):
        agent_id = str(event.get("agent") or event.get("agent_id") or "")
        if not agent_id or agent_id in agent_state:
            continue
        event_type = str(event.get("type") or "").lower()
        if "error" in event_type or "fail" in event_type:
            agent_state[agent_id] = "errored"
        elif "start" in event_type or "run" in event_type:
            agent_state[agent_id] = "busy"
        else:
            agent_state[agent_id] = "idle"

    # Also include agents with messages in bus inbox
    for aid in all_agents():
        if aid not in agent_state:
            unread = bus.unread_count(aid)
            agent_state[aid] = "busy" if unread > 0 else "idle"

    counts: dict[str, int] = {"idle": 0, "busy": 0, "errored": 0, "unknown": 0}
    for state in agent_state.values():
        counts[state if state in counts else "unknown"] += 1

    return {
        "paused":  _swarm_paused,
        "agents":  [{"id": aid, "state": state} for aid, state in agent_state.items()],
        "summary": counts,
        "total":   len(agent_state),
    }


@router.post("/swarm/pause")
def swarm_pause():
    """Pause the swarm — signals agents to stop accepting new tasks."""
    global _swarm_paused
    _swarm_paused = True
    return {"paused": True, "message": "Swarm paused. Existing tasks will complete."}


@router.post("/swarm/resume")
def swarm_resume():
    """Resume the swarm after a pause."""
    global _swarm_paused
    _swarm_paused = False
    return {"paused": False, "message": "Swarm resumed."}


# ── Blueprint execution, export, and import ───────────────────────────────────

@router.post("/architecture/blueprints/{name}/execute")
async def execute_architecture_blueprint(name: str, request: Request):
    """Execute a blueprint — spawn the specialist agents it defines.

    Publishes a task message on the agent bus for each agent role in the
    blueprint's agent_layer.  Returns a list of dispatched task message IDs.

    POST body (optional):
        task     — override task description sent to each agent
        version  — specific blueprint version to execute (default: latest)
    """
    body: dict = {}
    try:
        body = await request.json()
    except Exception:
        body = {}

    version = int(body.get("version", 0)) or None
    task    = str(body.get("task") or "Execute blueprint workflow").strip()

    from ..db import load_architecture_blueprint
    blueprint = load_architecture_blueprint(name, version=version)
    if not blueprint:
        return _api_error(f"Blueprint '{name}' not found", "not_found", 404)

    if _swarm_paused:
        return _api_error("Swarm is paused — resume before executing blueprints", "swarm_paused", 409)

    from ..agent_bus import post_message
    snapshot  = blueprint.get("snapshot", {})
    agent_ids = [a.get("id") for a in snapshot.get("agent_layer", []) if a.get("id")]

    if not agent_ids:
        return _api_error("Blueprint has no agents to execute", "empty_blueprint", 422)

    dispatched = []
    for agent_id in agent_ids:
        msg = post_message(
            from_id = "blueprint_executor",
            to_id   = agent_id,
            content = task,
            topic   = f"blueprint:{name}",
        )
        dispatched.append({"agent_id": agent_id, "msg_id": msg.msg_id})

    return {
        "blueprint":  name,
        "version":    blueprint.get("version"),
        "task":       task,
        "dispatched": dispatched,
        "count":      len(dispatched),
    }


@router.get("/architecture/blueprints/{name}/export")
def export_architecture_blueprint(name: str, version: int = 0):
    """Export a blueprint as a portable JSON file (downloadable).

    Query params:
        version — specific version to export (0 = latest)
    """
    from ..db import load_architecture_blueprint
    blueprint = load_architecture_blueprint(name, version=version if version > 0 else None)
    if not blueprint:
        return _api_error(f"Blueprint '{name}' not found", "not_found", 404)

    payload   = json.dumps(blueprint, indent=2).encode()
    safe_name = "".join(c if c.isalnum() or c in "-_" else "_" for c in name)
    filename  = f"blueprint_{safe_name}_v{blueprint.get('version', 1)}.json"
    return Response(
        content    = payload,
        media_type = "application/json",
        headers    = {"Content-Disposition": f'attachment; filename="{filename}"'},
    )


# NOTE: /architecture/blueprints/import must be registered before
# /architecture/blueprints/{name} to avoid {name} capturing "import".
# (FastAPI registers routes in order, so this is handled in app.py include_router.)
@router.post("/architecture/blueprints/import", status_code=201)
async def import_architecture_blueprint(request: Request):
    """Import a previously exported blueprint JSON.

    POST body: the full blueprint JSON object (as exported by GET …/export).
    Optional override keys:
        name  — rename the blueprint on import
        notes — override notes field
    """
    body: dict = {}
    try:
        body = await request.json()
    except Exception:
        return _api_error("Invalid JSON body", "validation_error", 422)

    # Accept either the full export envelope or a raw snapshot
    if "snapshot" in body:
        name     = str(body.get("name") or "imported").strip() or "imported"
        notes    = str(body.get("notes") or "").strip()
        snapshot = body["snapshot"]
    else:
        # Treat the whole body as the snapshot
        name     = str(body.get("name") or "imported").strip() or "imported"
        notes    = ""
        snapshot = body

    if not isinstance(snapshot, dict):
        return _api_error("'snapshot' must be a JSON object", "validation_error", 422)

    from ..db import save_architecture_blueprint
    created = save_architecture_blueprint(name=name, snapshot=snapshot, notes=notes)
    return {"blueprint": created, "status": "imported"}




# ═══════════════════════════════════════════════════════════════════════════════
#  NEW ROUTES — Sections 17-25 feature implementations
# ═══════════════════════════════════════════════════════════════════════════════

# ── S3/R2 Backup (Section 17.2) ───────────────────────────────────────────────

@router.post("/backup/s3/configure")
async def api_backup_s3_configure(request: Request):
    from ..object_storage import configure_s3
    body = await request.json()
    result = configure_s3(body)
    if not result["ok"]:
        return _api_error(result.get("error", "S3 config failed"), status_code=400)
    return result


@router.post("/backup/s3/push")
async def api_backup_s3_push():
    from ..object_storage import push_db_to_s3
    result = push_db_to_s3()
    if not result["ok"]:
        return _api_error(result.get("error", "S3 push failed"), status_code=500)
    return result


@router.post("/backup/s3/restore")
async def api_backup_s3_restore(request: Request):
    from ..object_storage import restore_db_from_s3
    body: dict = {}
    try:
        body = await request.json()
    except Exception:
        pass
    result = restore_db_from_s3(key=body.get("key"))
    if not result["ok"]:
        return _api_error(result.get("error", "S3 restore failed"), status_code=500)
    return result


@router.get("/backup/s3/backups")
async def api_backup_s3_list(limit: int = 20):
    from ..object_storage import list_s3_backups
    return {"backups": list_s3_backups(limit=limit)}


# ── Notion / Obsidian Export (Section 17.3) ───────────────────────────────────

@router.post("/export/notion/{conversation_id}")
async def api_export_notion(conversation_id: str, request: Request):
    from ..object_storage import export_chat_to_notion
    body: dict = {}
    try:
        body = await request.json()
    except Exception:
        pass
    chat = body.get("chat") or db_load_chat(conversation_id) or {}
    result = export_chat_to_notion(chat)
    if not result.get("ok"):
        return _api_error(result.get("error", "Notion export failed"), status_code=500)
    return result


@router.post("/export/obsidian/{conversation_id}")
async def api_export_obsidian(conversation_id: str, request: Request):
    from ..object_storage import export_chat_to_obsidian
    body: dict = {}
    try:
        body = await request.json()
    except Exception:
        pass
    chat = body.get("chat") or db_load_chat(conversation_id) or {}
    result = export_chat_to_obsidian(chat)
    if not result.get("ok"):
        return _api_error(result.get("error", "Obsidian export failed"), status_code=500)
    return result


@router.post("/export/obsidian/workspace")
async def api_export_obsidian_workspace(request: Request):
    from ..object_storage import export_workspace_to_obsidian
    body: dict = {}
    try:
        body = await request.json()
    except Exception:
        pass
    chats = body.get("chats") or db_load_chats()
    result = export_workspace_to_obsidian(chats)
    return result


# ── Nexus Hub — multi-instance management (Section 18.1) ─────────────────────

@router.post("/hub/instances")
async def api_hub_register(request: Request):
    from ..nexus_hub import register_instance
    body = await request.json()
    result = register_instance(**body)
    return result


@router.get("/hub/instances")
async def api_hub_list(include_offline: bool = True):
    from ..nexus_hub import list_instances
    return {"instances": list_instances(include_offline=include_offline)}


@router.get("/hub/instances/{instance_id}")
async def api_hub_get(instance_id: str):
    from ..nexus_hub import get_instance
    inst = get_instance(instance_id)
    if not inst:
        return _api_error("Instance not found", status_code=404)
    return inst


@router.delete("/hub/instances/{instance_id}")
async def api_hub_deregister(instance_id: str):
    from ..nexus_hub import deregister_instance
    ok = deregister_instance(instance_id)
    return {"ok": ok}


@router.post("/hub/instances/{instance_id}/ping")
async def api_hub_ping(instance_id: str):
    from ..nexus_hub import ping_instance
    return await ping_instance(instance_id)


@router.post("/hub/instances/{instance_id}/passthrough")
async def api_hub_passthrough(instance_id: str, request: Request):
    from ..nexus_hub import passthrough_request
    body: dict = {}
    try:
        body = await request.json()
    except Exception:
        pass
    method  = body.get("method", "GET")
    path    = body.get("path", "/v1/health")
    payload = body.get("body")
    headers = body.get("headers", {})
    result  = await passthrough_request(instance_id, method, path, payload, headers)
    if not result.get("ok"):
        return _api_error(result.get("error", "Passthrough failed"), status_code=502)
    return result


# ── Nexus Systems passthrough (Section 25.3) ──────────────────────────────────

@router.post("/nexus/passthrough")
async def api_nexus_systems_passthrough(request: Request):
    from ..nexus_hub import nexus_systems_passthrough
    body = await request.json()
    method   = body.get("method", "GET")
    endpoint = body.get("endpoint", "/")
    payload  = body.get("body")
    result   = await nexus_systems_passthrough(method, endpoint, payload)
    if not result.get("ok"):
        return _api_error(result.get("error", "Nexus Systems passthrough failed"), status_code=502)
    return result


# ── Federated learning (Section 18.3) ─────────────────────────────────────────

@router.post("/federated/round")
async def api_federated_round(request: Request):
    from ..federated import compute_and_submit_update
    body = await request.json()
    samples = body.get("samples", [])
    try:
        global_round = int(body.get("global_round", 0))
    except Exception:
        return _api_error("global_round must be an integer", "validation_error", 422)

    result = compute_and_submit_update(samples, global_round)
    payload = result.to_dict()
    status = str(payload.get("status") or "")
    if status == "disabled":
        return JSONResponse(payload, status_code=503)
    if status == "failed":
        return JSONResponse(payload, status_code=400)
    return payload


@router.get("/federated/status")
async def api_federated_status():
    from ..federated import get_federation_status
    return get_federation_status()


@router.get("/federated/rounds")
async def api_federated_rounds(limit: int = 20):
    from ..federated import list_rounds
    return {"rounds": list_rounds(limit=limit)}


# ── Creative async jobs (Phase 2 contract) ───────────────────────────────────

@router.post("/creative/music/jobs")
async def api_create_music_job(request: Request):
    require_auth(request)
    from ..creative_jobs import create_creative_job

    body = await _read_json_body(request, "invalid JSON body")
    prompt = str(body.get("prompt") or "").strip()
    if not prompt:
        return _api_error("prompt is required", "validation_error", 422)

    created = create_creative_job(
        kind="music",
        prompt=prompt,
        params={
            "duration": int(body.get("duration") or 15),
            "style": str(body.get("style") or "ambient"),
        },
    )
    if not bool(created.get("ok", False)):
        code = str(created.get("error") or "creative_job_failed")
        status = 503 if code == "creative_tools_disabled" else 400
        return JSONResponse(created, status_code=status)
    return JSONResponse({"job": created}, status_code=202)


@router.post("/creative/3d/jobs")
async def api_create_3d_job(request: Request):
    require_auth(request)
    from ..creative_jobs import create_creative_job

    body = await _read_json_body(request, "invalid JSON body")
    prompt = str(body.get("prompt") or "").strip()
    if not prompt:
        return _api_error("prompt is required", "validation_error", 422)

    created = create_creative_job(
        kind="3d",
        prompt=prompt,
        params={"format": str(body.get("format") or "glb")},
    )
    if not bool(created.get("ok", False)):
        code = str(created.get("error") or "creative_job_failed")
        status = 503 if code == "creative_tools_disabled" else 400
        return JSONResponse(created, status_code=status)
    return JSONResponse({"job": created}, status_code=202)


@router.get("/creative/jobs/{job_id}")
async def api_get_creative_job(job_id: str, request: Request):
    require_auth(request)
    from ..creative_jobs import get_creative_job_status

    row = get_creative_job_status(job_id)
    if row is None:
        return _api_error("creative job not found", "not_found", 404)
    return {"job": row}


@router.get("/creative/jobs")
async def api_list_creative_jobs(request: Request, kind: str = "", limit: int = 50):
    require_auth(request)
    from ..creative_jobs import list_creative_jobs_status

    rows = list_creative_jobs_status(kind=kind, limit=limit)
    return {"jobs": rows, "count": len(rows)}


# ── MoE routing + reasoning modes (Section 20.1) ─────────────────────────────

@router.post("/routing/moe")
async def api_moe_route(request: Request):
    from ..moe_router import route_to_expert
    body = await request.json()
    prompt     = str(body.get("prompt", ""))
    persona    = body.get("persona", "nexus")
    complexity = int(body.get("complexity", 5))
    return route_to_expert(prompt, persona=persona, complexity=complexity)


@router.get("/reasoning/sessions")
async def api_reasoning_sessions(limit: int = 50):
    from ..moe_router import list_reasoning_sessions
    return {"sessions": list_reasoning_sessions(limit=limit)}


# ── SLO dashboard (Section 24.1) ──────────────────────────────────────────────

@router.get("/slo/dashboard")
async def api_slo_dashboard():
    from ..slo import get_slo_dashboard
    return get_slo_dashboard()


@router.post("/slo/latency")
async def api_slo_record_latency(request: Request):
    from ..slo import record_latency
    body     = await request.json()
    endpoint = str(body.get("endpoint", "unknown"))
    latency  = float(body.get("latency_ms", 0))
    success  = bool(body.get("success", True))
    record_latency(endpoint, latency, success)
    return {"ok": True}


@router.get("/slo/hotspots")
async def api_slo_hotspots(threshold_ms: float = 2000.0):
    from ..slo import detect_latency_hotspots
    return {"hotspots": detect_latency_hotspots(threshold_ms=threshold_ms)}


@router.get("/slo/hardware-hint")
async def api_slo_hardware_hint():
    from ..slo import get_hardware_routing_hint
    return get_hardware_routing_hint()


# ── Team budgets + cost attribution (Sections 24.2, 24.3) ────────────────────

@router.post("/billing/teams/{team_id}/budget")
async def api_set_team_budget(team_id: str, request: Request):
    from ..slo import set_team_budget
    body = await request.json()
    return set_team_budget(
        team_id=team_id,
        monthly_limit_usd=float(body.get("monthly_limit_usd", 100.0)),
        daily_limit_usd=float(body.get("daily_limit_usd", 0.0)),
        alert_threshold_pct=float(body.get("alert_threshold_pct", body.get("alert_at_pct", 80.0))),
    )


@router.get("/billing/teams/{team_id}/budget")
async def api_get_team_budget(team_id: str):
    from ..slo import get_team_budget
    b = get_team_budget(team_id)
    if not b:
        return _api_error("Team budget not found", status_code=404)
    return b


@router.get("/billing/teams")
async def api_list_team_budgets():
    from ..slo import list_team_budgets
    return {"budgets": list_team_budgets()}


@router.get("/billing/alerts")
async def api_list_budget_alerts():
    from ..slo import list_budget_alerts
    return {"alerts": list_budget_alerts()}


@router.post("/billing/attribution")
async def api_record_attribution(request: Request):
    from ..slo import record_attribution
    body = await request.json()
    record_attribution(
        team_id=str(body.get("team_id", "")),
        department=str(body.get("department", "")),
        user=str(body.get("user", "")),
        cost_usd=float(body.get("cost_usd", 0.0)),
        tokens=int(body.get("tokens", 0)),
        model=str(body.get("model", "")),
        endpoint=str(body.get("endpoint", "")),
    )
    return {"ok": True}


@router.get("/billing/attribution/report")
async def api_attribution_report(team_id: str | None = None,
                                   department: str | None = None):
    from ..slo import get_attribution_report
    return get_attribution_report(team_id=team_id, department=department)


@router.put("/billing/teams/{team_id}/capacity")
async def api_set_reserved_capacity(team_id: str, request: Request):
    from ..slo import set_reserved_capacity
    body = await request.json()
    return set_reserved_capacity(
        team_id=team_id,
        reserved_rps=float(body.get("reserved_rps", 0.0)),
        max_concurrency=int(body.get("max_concurrency", 0)),
        guarantee_tier=str(body.get("guarantee_tier", "standard")),
        expires_at=str(body.get("expires_at", "")),
    )


@router.get("/billing/teams/{team_id}/capacity")
async def api_get_reserved_capacity(team_id: str):
    from ..slo import get_reserved_capacity
    row = get_reserved_capacity(team_id)
    if not row:
        return _api_error("Reserved capacity not found", status_code=404)
    return row


@router.get("/billing/capacity")
async def api_list_reserved_capacity():
    from ..slo import list_reserved_capacity
    return {"capacity": list_reserved_capacity()}


@router.post("/billing/teams/{team_id}/capacity/check")
async def api_check_rate_guarantee(team_id: str, request: Request):
    from ..slo import check_rate_guarantee
    body = await request.json()
    return check_rate_guarantee(
        team_id=team_id,
        requested_rps=float(body.get("requested_rps", 0.0)),
        requested_concurrency=int(body.get("requested_concurrency", 0)),
    )


@router.put("/billing/teams/{team_id}/spot")
async def api_set_spot_policy(team_id: str, request: Request):
    from ..slo import set_spot_policy
    body = await request.json()
    return set_spot_policy(
        team_id=team_id,
        enabled=bool(body.get("enabled", False)),
        max_discount_pct=float(body.get("max_discount_pct", 50.0)),
        fallback_on_preempt=bool(body.get("fallback_on_preempt", True)),
        max_preemptions_per_hour=int(body.get("max_preemptions_per_hour", 2)),
    )


@router.get("/billing/teams/{team_id}/spot")
async def api_get_spot_policy(team_id: str):
    from ..slo import get_spot_policy
    row = get_spot_policy(team_id)
    if not row:
        return _api_error("Spot policy not found", status_code=404)
    return row


@router.post("/billing/teams/{team_id}/spot/events")
async def api_record_spot_event(team_id: str, request: Request):
    from ..slo import record_spot_event
    body = await request.json()
    return record_spot_event(
        team_id=team_id,
        event_type=str(body.get("event_type", "preempted")),
        details=body.get("details") if isinstance(body.get("details"), dict) else {},
    )


# ── Custom tool registry (Section 25.1) ───────────────────────────────────────

@router.post("/tools/registry")
async def api_tools_register(request: Request):
    from ..marketplace_registry import register_tool
    body = await request.json()
    result = register_tool(**body)
    return result


@router.get("/tools/registry")
async def api_tools_list(category: str | None = None):
    from ..marketplace_registry import list_tools
    return {"tools": list_tools(tag=category)}


@router.get("/tools/registry/{tool_id}")
async def api_tools_get(tool_id: str):
    from ..marketplace_registry import get_tool
    t = get_tool(tool_id)
    if not t:
        return _api_error("Tool not found", status_code=404)
    return t


@router.post("/tools/registry/{tool_id}/invoke")
async def api_tools_invoke(tool_id: str, request: Request):
    from ..marketplace_registry import invoke_custom_tool
    body   = await request.json()
    params = body.get("params", {})
    result = await invoke_custom_tool(tool_id, params)
    if not result.get("ok"):
        return _api_error(result.get("error", "Tool invocation failed"), status_code=400)
    return result


@router.delete("/tools/registry/{tool_id}")
async def api_tools_delete(tool_id: str):
    from ..marketplace_registry import deactivate_tool
    ok = deactivate_tool(tool_id)
    return {"ok": ok}


# ── Marketplace connectors (Section 25.2) ─────────────────────────────────────

@router.get("/marketplace/connectors")
async def api_connectors_list(installed_only: bool = False):
    from ..marketplace_registry import list_connectors
    return {"connectors": list_connectors(installed_only=installed_only)}


@router.get("/marketplace/connectors/{connector_id}")
async def api_connectors_get(connector_id: str):
    from ..marketplace_registry import get_connector
    c = get_connector(connector_id)
    if not c:
        return _api_error("Connector not found", status_code=404)
    return c


@router.post("/marketplace/connectors/{connector_id}/install")
async def api_connectors_install(connector_id: str, request: Request):
    from ..marketplace_registry import install_connector
    body   = await request.json()
    config = body.get("config", {})
    result = install_connector(connector_id, config=config)
    if not result.get("ok"):
        return _api_error(result.get("error", "Install failed"), status_code=400)
    return result


@router.delete("/marketplace/connectors/{connector_id}/install")
async def api_connectors_uninstall(connector_id: str):
    from ..marketplace_registry import uninstall_connector
    result = uninstall_connector(connector_id)
    return result


# ── Agent templates (Section 25.2) ───────────────────────────────────────────

@router.post("/marketplace/templates")
async def api_templates_publish(request: Request):
    from ..marketplace_registry import publish_template
    body   = await request.json()
    result = publish_template(**body)
    return result


@router.get("/marketplace/templates")
async def api_templates_list(category: str | None = None, tags: str | None = None):
    from ..marketplace_registry import list_templates
    tags_list = tags.split(",") if tags else None
    return {"templates": list_templates(tag=category)}


@router.get("/marketplace/templates/{template_id}")
async def api_templates_get(template_id: str):
    from ..marketplace_registry import get_template
    t = get_template(template_id)
    if not t:
        return _api_error("Template not found", status_code=404)
    return t


@router.post("/marketplace/templates/{template_id}/deploy")
async def api_templates_deploy(template_id: str, request: Request):
    from ..marketplace_registry import deploy_template
    body      = await request.json()
    overrides = body.get("overrides", {})
    result    = deploy_template(template_id)
    if not result.get("ok"):
        return _api_error(result.get("error", "Deploy failed"), status_code=400)
    return result


# ── Community personas (Section 25.2) ────────────────────────────────────────

@router.post("/marketplace/personas")
async def api_community_persona_publish(request: Request):
    from ..marketplace_registry import publish_persona
    body   = await request.json()
    result = publish_persona(**body)
    return result


@router.get("/marketplace/personas")
async def api_community_personas_list(tags: str | None = None):
    from ..marketplace_registry import list_community_personas
    tags_list = tags.split(",") if tags else None
    return {"personas": list_community_personas(tag=tags_list[0] if tags_list else None)}


@router.get("/marketplace/personas/{persona_id}")
async def api_community_persona_get(persona_id: str):
    from ..marketplace_registry import get_community_persona
    p = get_community_persona(persona_id)
    if not p:
        return _api_error("Persona not found", status_code=404)
    return p


# ── Provider plugins (Section 25.2) ──────────────────────────────────────────

@router.post("/providers/plugins")
async def api_providers_plugin_register(request: Request):
    from ..marketplace_registry import register_provider_plugin
    body   = await request.json()
    result = register_provider_plugin(**body)
    return result


@router.get("/providers/plugins")
async def api_providers_plugins_list(active_only: bool = False):
    from ..marketplace_registry import list_provider_plugins
    return {"plugins": list_provider_plugins(active_only=active_only)}


@router.post("/providers/plugins/{plugin_id}/activate")
async def api_providers_plugin_activate(plugin_id: str, request: Request):
    from ..marketplace_registry import activate_provider_plugin
    body   = await request.json()
    active = bool(body.get("active", True))
    if not active:
        # deactivation not yet implemented, return current state
        return {"ok": False, "note": "Deactivation not yet implemented"}
    result = activate_provider_plugin(plugin_id)
    if not result:
        return _api_error("Plugin not found", status_code=404)
    return result


# ── Vision extras — chart extraction + document understanding (Section 22) ────

@router.post("/vision/extract-charts")
async def api_vision_extract_charts(request: Request):
    from ..vision import extract_charts_and_tables
    import base64
    body      = await request.json()
    image_b64 = str(body.get("image_base64", ""))
    mime_type = str(body.get("mime_type", "image/png"))
    if not image_b64:
        return _api_error("image_base64 required", status_code=422)
    try:
        image_bytes = base64.b64decode(image_b64)
    except Exception:
        return _api_error("Invalid base64 image", status_code=422)
    return extract_charts_and_tables(image_bytes, mime_type=mime_type)


@router.post("/documents/understand")
async def api_documents_understand(request: Request):
    from ..vision import understand_pdf, understand_office_doc
    import base64
    body      = await request.json()
    file_b64  = str(body.get("file_base64", ""))
    filename  = str(body.get("filename", "document.pdf"))
    mime_type = str(body.get("mime_type", "application/pdf"))
    if not file_b64:
        return _api_error("file_base64 required", status_code=422)
    try:
        file_bytes = base64.b64decode(file_b64)
    except Exception:
        return _api_error("Invalid base64 file", status_code=422)
    if "pdf" in mime_type or filename.lower().endswith(".pdf"):
        return understand_pdf(file_bytes)
    return understand_office_doc(file_bytes, filename=filename)


@router.post("/documents/diff")
async def api_documents_diff(request: Request):
    from ..vision import diff_documents
    body     = await request.json()
    text_a   = str(body.get("text_a", ""))
    text_b   = str(body.get("text_b", ""))
    ctx      = int(body.get("context_lines", 3))
    return diff_documents(text_a, text_b, context_lines=ctx)


# ── Audio extras — analysis + diarization + streaming (Section 22) ────────────
















# ── Video extras — transcription + chapters + editing (Section 22) ───────────

@router.post("/video/transcribe")
async def api_video_transcribe(request: Request):
    from ..generation import video_to_text
    import base64
    body           = await request.json()
    video_b64      = str(body.get("video_base64", ""))
    frame_interval = float(body.get("frame_interval_s", 5.0))
    max_frames     = int(body.get("max_frames", 20))
    if not video_b64:
        return _api_error("video_base64 required", status_code=422)
    video_bytes = base64.b64decode(video_b64)
    return video_to_text(video_bytes, frame_interval_s=frame_interval, max_frames=max_frames)


@router.post("/video/chapters")
async def api_video_chapters(request: Request):
    from ..generation import detect_video_chapters_from_transcript
    import base64
    body      = await request.json()
    video_b64 = body.get("video_base64")
    video_url = body.get("video_url")
    if video_b64:
        video_bytes = base64.b64decode(video_b64)
        chapters    = detect_video_chapters_from_transcript(video_bytes)
    elif video_url:
        from ..generation import detect_video_chapters
        chapters = detect_video_chapters(video_url)
    else:
        return _api_error("video_base64 or video_url required", status_code=422)
    return {"chapters": chapters}


@router.post("/video/edit")
async def api_video_edit(request: Request):
    from ..generation import edit_video
    import base64
    body        = await request.json()
    video_b64   = str(body.get("video_base64", ""))
    operations  = body.get("operations", [])
    if not video_b64:
        return _api_error("video_base64 required", status_code=422)
    video_bytes = base64.b64decode(video_b64)
    result      = edit_video(video_bytes, operations)
    if result.get("output_bytes"):
        output_b64         = base64.b64encode(result["output_bytes"]).decode()
        result["output_bytes"] = None
        result["output_base64"] = output_b64
    return result





# ── Developer tooling — code review + bug fix + migration (Section 18) ────────

_BUG_FIX_CHECKPOINTS_KEY = "agent_bug_fix_checkpoints_v1"


def _load_bug_fix_checkpoints() -> list[dict[str, Any]]:
    raw = db_load_pref(_BUG_FIX_CHECKPOINTS_KEY, "[]")
    if isinstance(raw, list):
        return [row for row in raw if isinstance(row, dict)]
    if isinstance(raw, str):
        try:
            parsed = json.loads(raw) if raw.strip() else []
            if isinstance(parsed, list):
                return [row for row in parsed if isinstance(row, dict)]
        except Exception:
            return []
    return []


def _save_bug_fix_checkpoints(rows: list[dict[str, Any]]) -> None:
    db_save_pref(_BUG_FIX_CHECKPOINTS_KEY, json.dumps(rows[-1000:], separators=(",", ":")))


@router.post("/benchmark/safety")
async def api_benchmark_safety(request: Request):
    from ..benchmark import run_safety_benchmark
    body = await request.json()
    return run_safety_benchmark(test_cases=body.get("test_cases") or None)


@router.get("/benchmark/safety/results")
async def api_benchmark_safety_results():
    from ..benchmark import load_safety_benchmark_results, evaluate_safety_release_gate
    results = load_safety_benchmark_results(limit=500)
    return {"results": results, "release_gate": evaluate_safety_release_gate(results)}


@router.get("/benchmark/safety/gate")
async def api_benchmark_safety_gate():
    from ..benchmark import load_safety_gate_config, load_safety_benchmark_results, evaluate_safety_release_gate
    cfg = load_safety_gate_config()
    return {"config": cfg, "status": evaluate_safety_release_gate(load_safety_benchmark_results(limit=500), cfg)}


@router.post("/benchmark/safety/gate")
async def api_benchmark_safety_gate_config(request: Request):
    from ..benchmark import update_safety_gate_config
    body = await _read_json_body(request, "invalid JSON body")
    return update_safety_gate_config(body)


# ── Benchmark regression (Section 19.3) ──────────────────────────────────────

@router.get("/benchmark/regression")
async def api_benchmark_regression():
    from ..benchmark import get_regression_report
    return get_regression_report()


@router.post("/benchmark/regression/baseline")
async def api_benchmark_regression_baseline(request: Request):
    from ..benchmark import set_regression_baseline
    body = await request.json()
    return set_regression_baseline(body)


# ── Concurrent task management extras (Section 23) ───────────────────────────

@router.post("/tasks/queue/{task_id}/cancel")
async def api_task_cancel(task_id: str):
    from ..task_queue import cancel_task, get_task_runtime_status
    ok = cancel_task(task_id)
    if not ok:
        return _api_error("Task not found or already terminal", status_code=404)
    return {
        "ok": True,
        "task_id": task_id,
        "status": get_task_runtime_status(task_id),
    }


@router.get("/tasks/queue/{task_id}/status")
async def api_task_status(task_id: str):
    from ..task_queue import get_task_runtime_status
    status = get_task_runtime_status(task_id)
    if status is None:
        return _api_error("Task not found", status_code=404)
    return status


@router.post("/tasks/sessions")
@router.post("/tasks/queue/sessions")
@router.post("/task-sessions")
async def api_task_session_create(request: Request):
    body: dict = {}
    try:
        body = await request.json()
    except Exception:
        body = {}
    session_id = str(body.get("session_id") or f"ts_{secrets.token_hex(6)}")
    label = str(body.get("label") or "")
    created_by = str(body.get("created_by") or "")
    return {
        "session_id": session_id,
        "label": label,
        "created_by": created_by,
        "created_at": datetime.now(timezone.utc).isoformat(),
    }


@router.get("/tasks/sessions")
@router.get("/tasks/queue/sessions")
@router.get("/task-sessions")
async def api_task_sessions(limit: int = 200):
    from ..task_queue import list_task_sessions
    return {"sessions": list_task_sessions(limit=max(1, min(int(limit), 1000)))}


@router.get("/tasks/sessions/{session_id}")
@router.get("/tasks/queue/sessions/{session_id}")
@router.get("/task-sessions/{session_id}")
async def api_task_session_detail(session_id: str, status: str = "", limit: int = 200):
    from ..task_queue import get_session_tasks
    tasks = get_session_tasks(session_id=session_id, status=status, limit=max(1, min(int(limit), 1000)))
    if not tasks:
        return _api_error("Task session not found", status_code=404)
    return {"session_id": session_id, "tasks": tasks, "total": len(tasks)}


@router.post("/tasks/sessions/{session_id}/cancel")
@router.post("/tasks/queue/sessions/{session_id}/cancel")
@router.post("/task-sessions/{session_id}/cancel")
async def api_task_session_cancel(session_id: str, request: Request):
    from ..task_queue import cancel_session_tasks
    body: dict = {}
    try:
        body = await request.json()
    except Exception:
        body = {}
    include_running = bool(body.get("include_running", True))
    result = cancel_session_tasks(session_id=session_id, include_running=include_running)
    if result.get("cancelled", 0) == 0:
        return _api_error("No cancellable tasks found for session", status_code=404)
    return result


# ── 26.4 Evaluation: A/B testing, human eval ─────────────────────────────────

@router.get("/evals/experiments")
async def api_list_experiments(request: Request):
    try:
        from ..evals.ab_testing import list_experiments
        return {"experiments": list_experiments()}
    except Exception as exc:
        return _api_error(str(exc))


@router.post("/evals/experiments")
async def api_create_experiment(request: Request):
    body = await request.json()
    name = str(body.get("name", "")).strip()
    variants = body.get("variants", [])
    if not name or not variants:
        return _api_error("name and variants are required", status_code=400)
    try:
        from ..evals.ab_testing import create_experiment
        exp = create_experiment(
            name=name, variants=variants,
            metric=str(body.get("metric", "quality_score")),
            description=str(body.get("description", "")),
        )
        return {"experiment_id": exp.experiment_id, "name": exp.name, "status": exp.status}
    except Exception as exc:
        return _api_error(str(exc))


@router.post("/evals/experiments/{experiment_id}/start")
async def api_start_experiment(request: Request, experiment_id: str):
    try:
        from ..evals.ab_testing import start_experiment
        ok = start_experiment(experiment_id)
        return {"started": ok, "experiment_id": experiment_id}
    except Exception as exc:
        return _api_error(str(exc))


@router.post("/evals/experiments/{experiment_id}/pause")
async def api_pause_experiment(request: Request, experiment_id: str):
    try:
        from ..evals.ab_testing import pause_experiment
        ok = pause_experiment(experiment_id)
        return {"paused": ok, "experiment_id": experiment_id}
    except Exception as exc:
        return _api_error(str(exc))


@router.get("/evals/experiments/{experiment_id}/analysis")
async def api_analyze_experiment(request: Request, experiment_id: str):
    try:
        from ..evals.ab_testing import analyze_experiment
        return analyze_experiment(experiment_id)
    except Exception as exc:
        return _api_error(str(exc))


@router.post("/webhooks/outbound")
async def api_enqueue_webhook(request: Request):
    body = await request.json()
    url = str(body.get("url", "")).strip()
    event_type = str(body.get("event_type", "generic")).strip()
    payload = body.get("payload", {})
    if not url:
        return _api_error("url is required", status_code=400)
    try:
        from ..webhooks_delivery import enqueue_webhook
        delivery = enqueue_webhook(url=url, event_type=event_type, payload=payload)
        return {"delivery_id": delivery.delivery_id, "status": delivery.status, "url": url}
    except Exception as exc:
        return _api_error(str(exc))


# ── 26.7 Agent Capabilities: persistent state, tool policies, structured output

@router.post("/agents/state")
async def api_create_agent_state(request: Request):
    body = await request.json()
    objective = str(body.get("objective", "")).strip()
    if not objective:
        return _api_error("objective is required", status_code=400)
    user = body.get("username", "") or ""
    try:
        from ..agent_state import create_agent_state
        state = create_agent_state(
            objective=objective,
            session_id=str(body.get("session_id", "")),
            username=user,
            persona_id=str(body.get("persona_id", "")),
            metadata=body.get("metadata", {}),
        )
        return {"agent_id": state.agent_id, "objective": state.objective, "status": state.status}
    except Exception as exc:
        return _api_error(str(exc))


@router.get("/agents/state/{agent_id}")
async def api_get_agent_state(request: Request, agent_id: str):
    try:
        from ..agent_state import get_agent_state, get_plan_summary
        state = get_agent_state(agent_id)
        if not state:
            return _api_error("Agent state not found", status_code=404)
        summary = get_plan_summary(state)
        return {
            "agent_id": state.agent_id, "status": state.status,
            "objective": state.objective, "plan_summary": summary,
            "step_count": state.step_count, "updated_at": state.updated_at,
            "working_memory": state.working_memory,
        }
    except Exception as exc:
        return _api_error(str(exc))


@router.get("/agents/active")
async def api_list_active_agents(request: Request, username: str = ""):
    try:
        from ..agent_state import list_active_agents
        return {"agents": list_active_agents(username=username or None)}
    except Exception as exc:
        return _api_error(str(exc))


@router.post("/structured-output/generate")
async def api_structured_generate(request: Request):
    body = await request.json()
    prompt = str(body.get("prompt", "")).strip()
    schema = body.get("schema", {})
    if not prompt or not schema:
        return _api_error("prompt and schema are required", status_code=400)
    try:
        from ..structured_output import generate_structured
        result = generate_structured(
            prompt=prompt, schema=schema,
            model_name=str(body.get("model", "")),
            max_tokens=int(body.get("max_tokens", 2048)),
            temperature=float(body.get("temperature", 0.1)),
        )
        return result
    except Exception as exc:
        return _api_error(str(exc))


@router.post("/structured-output/validate")
async def api_structured_validate(request: Request):
    body = await request.json()
    output = str(body.get("output", ""))
    schema = body.get("schema", {})
    if not output or not schema:
        return _api_error("output and schema are required", status_code=400)
    try:
        from ..structured_output import validate_output
        return validate_output(output, schema)
    except Exception as exc:
        return _api_error(str(exc))


# ── Dataset-backed benchmark runners ─────────────────────────────────────────

@router.post("/benchmark/dataset/run")
async def api_benchmark_dataset_run(request: Request):
    """Run a publishable dataset benchmark (gsm8k, truthfulqa, humaneval, mmlu, hellaswag)."""
    body = await request.json()
    dataset = str(body.get("dataset", "gsm8k")).strip()
    provider = str(body.get("provider", "")).strip()
    model = str(body.get("model", "")).strip()
    max_samples = min(int(body.get("max_samples", 10)), 50)
    try:
        from ..benchmark import run_dataset_benchmark
        return run_dataset_benchmark(dataset=dataset, provider=provider, model=model, max_samples=max_samples)
    except Exception as exc:
        return _api_error(str(exc))


@router.post("/benchmark/dataset/suite")
async def api_benchmark_dataset_suite(request: Request):
    """Run a full suite of dataset benchmarks and return aggregated results."""
    body = await request.json()
    datasets = body.get("datasets") or None
    provider = str(body.get("provider", "")).strip()
    model = str(body.get("model", "")).strip()
    max_samples = min(int(body.get("max_samples_per_dataset", 10)), 50)
    try:
        from ..benchmark import run_dataset_suite_benchmark
        return run_dataset_suite_benchmark(datasets=datasets, provider=provider, model=model, max_samples_per_dataset=max_samples)
    except Exception as exc:
        return _api_error(str(exc))


@router.get("/benchmark/dataset/history")
async def api_benchmark_dataset_history(request: Request):
    """Return persisted dataset benchmark history."""
    dataset = str(request.query_params.get("dataset", "")).strip()
    limit = min(int(request.query_params.get("limit", 50)), 500)
    try:
        from ..benchmark import get_dataset_benchmark_history
        return get_dataset_benchmark_history(dataset=dataset, limit=limit)
    except Exception as exc:
        return _api_error(str(exc))


@router.get("/benchmark/dataset/datasets")
async def api_benchmark_dataset_list(request: Request):
    """List available dataset benchmark runners."""
    from ..evals.dataset_runners import DATASET_RUNNERS
    return {
        "datasets": list(DATASET_RUNNERS.keys()),
        "hf_enabled": __import__("os").getenv("BENCHMARK_USE_HF_DATASETS", "false").lower() == "true",
        "note": "Set BENCHMARK_USE_HF_DATASETS=true to load live samples from HuggingFace",
    }


# ── Benchmark artifact export ─────────────────────────────────────────────────

@router.get("/benchmark/export/{run_id}")
async def api_benchmark_export(request: Request, run_id: str):
    """Export a benchmark run as publishable artifacts (JSONL, CSV, HTML, leaderboard JSON)."""
    formats_raw = str(request.query_params.get("formats", "")).strip()
    formats = [f.strip() for f in formats_raw.split(",") if f.strip()] or None
    try:
        from ..benchmark import export_benchmark_run
        result = export_benchmark_run(run_id=run_id, formats=formats)
        if "error" in result:
            return _api_error(result["error"], status_code=404)
        return result
    except Exception as exc:
        return _api_error(str(exc))


@router.post("/benchmark/export/suite")
async def api_benchmark_export_suite(request: Request):
    """Export a complete benchmark suite as publishable artifacts."""
    body = await request.json()
    suite_data = body.get("suite_data", {})
    full_results = body.get("full_results", [])
    formats_raw = str(body.get("formats", "")).strip()
    formats = [f.strip() for f in formats_raw.split(",") if f.strip()] or None
    if not suite_data or not full_results:
        return _api_error("suite_data and full_results are required", status_code=400)
    try:
        from ..benchmark import export_dataset_suite_artifacts
        return export_dataset_suite_artifacts(suite_data=suite_data, full_results=full_results, formats=formats)
    except Exception as exc:
        return _api_error(str(exc))


@router.get("/benchmark/export/{run_id}/html")
async def api_benchmark_export_html(request: Request, run_id: str):
    """Return the HTML benchmark report directly (text/html content type)."""
    from fastapi.responses import HTMLResponse
    try:
        from ..benchmark import export_benchmark_run
        result = export_benchmark_run(run_id=run_id, formats=["html"])
        if "error" in result:
            return HTMLResponse(f"<h1>Not found</h1><p>{result['error']}</p>", status_code=404)
        return HTMLResponse(content=result.get("html", "<p>No HTML output.</p>"))
    except Exception as exc:
        return HTMLResponse(f"<h1>Error</h1><p>{exc}</p>", status_code=500)


# ── Domain-route sub-routers ────────────────────────────
from ..routes.finetune import router as _finetune_router
router.include_router(_finetune_router)

from ..routes.auth import router as _auth_router
router.include_router(_auth_router)

from ..routes.safety import router as _safety_router
router.include_router(_safety_router)

from ..routes.admin import router as _admin_router
router.include_router(_admin_router)

from ..routes.agent import router as _agent_router
router.include_router(_agent_router)

from ..routes.audio import router as _audio_router
router.include_router(_audio_router)

from ..routes.orgs import router as _orgs_router
router.include_router(_orgs_router)

from ..routes.workspace import router as _workspace_router
router.include_router(_workspace_router)

from ..routes.rag import router as _rag_router
router.include_router(_rag_router)

from ..routes.reasoning import router as _reasoning_router
router.include_router(_reasoning_router)

from ..routes.browser import router as _browser_router
router.include_router(_browser_router)

from ..routes.collab import router as _collab_router
router.include_router(_collab_router)

from ..routes.mcp import router as _mcp_router
router.include_router(_mcp_router)

from ..routes.v1 import router as _v1_router
router.include_router(_v1_router)

from ..routes.rlhf import router as _rlhf_router
router.include_router(_rlhf_router)
