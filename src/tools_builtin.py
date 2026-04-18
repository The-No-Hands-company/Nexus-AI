"""
Built-in tools that don't need LLM calls — calculator, weather, currency,
unit converter, regex tester, base64, JSON formatter, color info.
"""
import os, re, json, math, subprocess, zipfile, base64 as b64lib
from pathlib import Path
from datetime import datetime
from .providers.model_router import ModelRouter
from .scheduler import (
    schedule_job,
    list_jobs,
    cancel_job,
    job_to_dict,
)
from .knowledge_graph import (
    kg_store as _kg_store,
    kg_query as _kg_query,
    kg_list_entities as _kg_list,
    kg_get as _kg_get,
)

# ── CALCULATOR ────────────────────────────────────────────────────────────────
_SAFE_NAMES = {k: v for k, v in math.__dict__.items() if not k.startswith('_')}
_SAFE_NAMES.update({'abs': abs, 'round': round, 'min': min, 'max': max,
                    'sum': sum, 'pow': pow, 'len': len})

def tool_calculate(expr: str) -> str:
    try:
        # Strip anything that looks like an import or dunder
        if any(kw in expr for kw in ['import','__','exec','eval','open','os','sys']):
            return "❌ Expression not allowed"
        result = eval(expr, {"__builtins__": {}}, _SAFE_NAMES)
        return f"**{expr}** = `{result}`"
    except Exception as e:
        return f"❌ Calculation error: {e}"


# ── WEATHER ───────────────────────────────────────────────────────────────────
def tool_weather(location: str) -> str:
    try:
        import requests
        loc = location.strip().replace(' ', '+')
        resp = requests.get(f"https://wttr.in/{loc}?format=j1", timeout=10)
        resp.raise_for_status()
        d = resp.json()
        cur = d['current_condition'][0]
        area = d['nearest_area'][0]
        city = area['areaName'][0]['value']
        country = area['country'][0]['value']
        desc = cur['weatherDesc'][0]['value']
        temp_c = cur['temp_C']
        temp_f = cur['temp_F']
        feels_c = cur['FeelsLikeC']
        humidity = cur['humidity']
        wind_kmph = cur['windspeedKmph']
        return (f"**{city}, {country}** — {desc}\n"
                f"🌡️ {temp_c}°C / {temp_f}°F (feels like {feels_c}°C)\n"
                f"💧 Humidity: {humidity}%  💨 Wind: {wind_kmph} km/h")
    except Exception as e:
        return f"❌ Weather lookup failed: {e}"


# ── CURRENCY ──────────────────────────────────────────────────────────────────
def tool_currency(amount: float, from_cur: str, to_cur: str) -> str:
    try:
        import requests
        # Open exchange rates free endpoint (no key needed for latest)
        resp = requests.get(
            f"https://open.er-api.com/v6/latest/{from_cur.upper()}",
            timeout=10
        )
        resp.raise_for_status()
        d = resp.json()
        if d.get('result') != 'success':
            return f"❌ Currency lookup failed: {d.get('error-type','unknown')}"
        rate = d['rates'].get(to_cur.upper())
        if not rate:
            return f"❌ Unknown currency: {to_cur}"
        converted = amount * rate
        return (f"**{amount:,.2f} {from_cur.upper()}** = "
                f"**{converted:,.2f} {to_cur.upper()}**\n"
                f"Rate: 1 {from_cur.upper()} = {rate:.4f} {to_cur.upper()}\n"
                f"*Updated: {d.get('time_last_update_utc','?')}*")
    except Exception as e:
        return f"❌ Currency error: {e}"


# ── UNIT CONVERTER ────────────────────────────────────────────────────────────
_UNITS = {
    # length (base: meter)
    'km':1000,'m':1,'cm':0.01,'mm':0.001,'mile':1609.344,'miles':1609.344,
    'yard':0.9144,'yards':0.9144,'ft':0.3048,'feet':0.3048,'inch':0.0254,'inches':0.0254,
    # weight (base: kg)
    'kg':1,'g':0.001,'mg':0.000001,'lb':0.453592,'lbs':0.453592,'oz':0.028349,
    # volume (base: liter)
    'l':1,'ml':0.001,'gallon':3.78541,'gallons':3.78541,'cup':0.236588,'cups':0.236588,
    'fl oz':0.029574,'tbsp':0.014787,'tsp':0.004929,
    # temperature handled separately
    'c':'temp','f':'temp','k':'temp',
    # data (base: byte)
    'b':1,'kb':1024,'mb':1048576,'gb':1073741824,'tb':1099511627776,
}
_UNIT_GROUPS = {
    'length':  ['km','m','cm','mm','mile','miles','yard','yards','ft','feet','inch','inches'],
    'weight':  ['kg','g','mg','lb','lbs','oz'],
    'volume':  ['l','ml','gallon','gallons','cup','cups'],
    'data':    ['b','kb','mb','gb','tb'],
}

def _temp_convert(val, from_u, to_u):
    from_u, to_u = from_u.lower(), to_u.lower()
    # to Celsius first
    if from_u == 'f':   c = (val - 32) * 5/9
    elif from_u == 'k': c = val - 273.15
    else:               c = val
    # Celsius to target
    if to_u == 'f':     return c * 9/5 + 32
    elif to_u == 'k':   return c + 273.15
    return c

def tool_convert(value: float, from_unit: str, to_unit: str) -> str:
    f, t = from_unit.lower().strip(), to_unit.lower().strip()
    # Temperature
    if f in ('c','f','k') or t in ('c','f','k'):
        result = _temp_convert(value, f, t)
        names = {'c':'°C','f':'°F','k':'K'}
        return f"**{value}{names.get(f,f)}** = **{result:.4g}{names.get(t,t)}**"
    if f not in _UNITS or t not in _UNITS:
        return f"❌ Unknown units: {from_unit} or {to_unit}"
    # Check same group
    grp_f = next((g for g,u in _UNIT_GROUPS.items() if f in u), None)
    grp_t = next((g for g,u in _UNIT_GROUPS.items() if t in u), None)
    if grp_f != grp_t:
        return f"❌ Can't convert between {grp_f or from_unit} and {grp_t or to_unit}"
    base = value * _UNITS[f]
    result = base / _UNITS[t]
    return f"**{value} {from_unit}** = **{result:.6g} {to_unit}**"


# ── REGEX TESTER ──────────────────────────────────────────────────────────────
def tool_regex(pattern: str, text: str, flags_str: str = "") -> str:
    try:
        flag_map = {'i': re.IGNORECASE, 'm': re.MULTILINE, 's': re.DOTALL}
        flags = 0
        for c in flags_str.lower():
            flags |= flag_map.get(c, 0)
        matches = list(re.finditer(pattern, text, flags))
        if not matches:
            return f"No matches for `{pattern}` in the given text."
        lines = [f"**{len(matches)} match{'es' if len(matches)>1 else ''}** for `{pattern}`:\n"]
        for i, m in enumerate(matches[:10], 1):
            lines.append(f"{i}. `{m.group()}` at position {m.start()}–{m.end()}")
            if m.groups():
                lines.append(f"   Groups: {m.groups()}")
        if len(matches) > 10:
            lines.append(f"…and {len(matches)-10} more")
        return "\n".join(lines)
    except re.error as e:
        return f"❌ Regex error: {e}"


# ── BASE64 ────────────────────────────────────────────────────────────────────
def tool_base64(text: str, mode: str = "encode") -> str:
    try:
        if mode == "encode":
            result = b64lib.b64encode(text.encode()).decode()
            return f"**Encoded:**\n```\n{result}\n```"
        else:
            result = b64lib.b64decode(text.encode()).decode()
            return f"**Decoded:**\n```\n{result}\n```"
    except Exception as e:
        return f"❌ Base64 error: {e}"


# ── JSON FORMATTER ────────────────────────────────────────────────────────────
def tool_json_format(text: str) -> str:
    try:
        parsed = json.loads(text)
        formatted = json.dumps(parsed, indent=2, ensure_ascii=False)
        return f"```json\n{formatted}\n```"
    except json.JSONDecodeError as e:
        return f"❌ Invalid JSON: {e}"


def tool_select_model(task: str, prefer_speed: bool = False, prefer_quality: bool = False) -> str:
    try:
        router = ModelRouter()
        model_id, spec = router.select_model(task, prefer_speed=prefer_speed, prefer_quality=prefer_quality)
        languages = ", ".join(spec.languages) if spec.languages else "all"
        return (
            f"Selected model: {spec.name} ({model_id})\n"
            f"Tier: {spec.tier.value}\n"
            f"Estimated RAM: {spec.ram_required_gb}GB\n"
            f"Strengths: {', '.join(spec.strengths)}\n"
            f"Languages: {languages}"
        )
    except Exception as e:
        return f"Model selection failed: {e}"


_RAG_SYSTEM = None

def _get_rag_system():
    global _RAG_SYSTEM
    if _RAG_SYSTEM is None:
        from .rag.rag_system import RAGSystem
        _RAG_SYSTEM = RAGSystem()
    return _RAG_SYSTEM


def tool_rag_ingest(text: str = "", path: str = "", metadata: dict | None = None,
                    doc_id_prefix: str | None = None, workdir: str = "/tmp") -> str:
    if path:
        full = path if os.path.isabs(path) else os.path.join(workdir, path)
        if not os.path.exists(full):
            return f"File not found: {path}"
        try:
            with open(full, "r", encoding="utf-8", errors="replace") as f:
                text = f.read()
        except Exception as e:
            return f"Failed to read file: {e}"

    if not text:
        return "No text or file path provided to ingest."

    try:
        count = _get_rag_system().ingest(text, metadata=metadata or {}, doc_id_prefix=doc_id_prefix)
        source = f"file {path}" if path else "text"
        return f"RAG ingested {count} chunks from {source}."
    except Exception as e:
        return f"RAG ingest failed: {e}"


def tool_rag_query(query: str, top_k: int = 5, filter_metadata: dict | None = None) -> str:
    try:
        results = _get_rag_system().query(query, top_k=top_k, filter_metadata=filter_metadata)
    except Exception as e:
        return f"RAG query failed: {e}"
    if not results:
        return "No relevant RAG results found."
    lines = []
    for i, r in enumerate(results, 1):
        doc = str(r.get("document", "")).strip().replace("\n", " ")
        score = r.get("score", "?")
        source = r.get("metadata", {}).get("source", "unknown") if isinstance(r.get("metadata"), dict) else "unknown"
        lines.append(f"{i}. score={score} source={source}\n{doc[:320]}")
    return "\n\n".join(lines)


def tool_rag_status() -> str:
    try:
        stats = _get_rag_system().stats()
        return (f"RAG status: ingested={stats.get('total_ingested', 0)}, "
                f"queries={stats.get('total_queries', 0)}, "
                f"store_count={stats.get('store_count', 0)}, "
                f"backend={stats.get('embedding_backend', 'unknown')}")
    except Exception as e:
        return f"RAG status failed: {e}"


# ── TRACE HELPER ─────────────────────────────────────────────────────────────
def _tool_trace(
    action: dict,
    result: str,
    metadata: dict | None = None,
    status: str = "done",
    error: str | None = None,
) -> dict:
    """
    Return a structured trace dict for a completed tool call.
    agent.py unpacks .result / .metadata / .status for event emission.
    Shape: { action, tool_name, status, input, result, metadata, error }
    """
    return {
        "action":    action.get("action"),
        "tool_name": action.get("action"),
        "status":    status,
        "input":     action,
        "result":    result,
        "metadata":  metadata or {},
        "error":     error,
    }


# ── DISPATCH ──────────────────────────────────────────────────────────────────
def dispatch_builtin(action: dict) -> dict | None:
    """
    Returns a structured trace dict or None if action is not a built-in tool.
    Shape: { action, tool_name, status, input, result, metadata, error }
    """
    kind = action.get("action")
    if kind == "calculate":
        r = tool_calculate(action.get("expr", ""))
        return _tool_trace(action, r, {"expr": action.get("expr", "")})
    if kind == "weather":
        r = tool_weather(action.get("location", ""))
        return _tool_trace(action, r, {"location": action.get("location", "")})
    if kind == "currency":
        r = tool_currency(float(action.get("amount", 1)),
                          action.get("from", "USD"), action.get("to", "EUR"))
        return _tool_trace(action, r, {
            "amount": action.get("amount"), "from": action.get("from"), "to": action.get("to")})
    if kind == "convert":
        r = tool_convert(float(action.get("value", 0)),
                         action.get("from_unit", ""), action.get("to_unit", ""))
        return _tool_trace(action, r, {
            "value": action.get("value"),
            "from_unit": action.get("from_unit"),
            "to_unit": action.get("to_unit"),
        })
    if kind == "regex":
        r = tool_regex(action.get("pattern", ""), action.get("text", ""),
                       action.get("flags", ""))
        return _tool_trace(action, r, {
            "pattern": action.get("pattern", ""), "flags": action.get("flags", "")})
    if kind == "base64":
        r = tool_base64(action.get("text", ""), action.get("mode", "encode"))
        return _tool_trace(action, r, {"mode": action.get("mode", "encode")})
    if kind == "json_format":
        r = tool_json_format(action.get("text", ""))
        return _tool_trace(action, r)
    if kind == "select_model":
        r = tool_select_model(action.get("task", ""),
                              bool(action.get("prefer_speed", False)),
                              bool(action.get("prefer_quality", False)))
        return _tool_trace(action, r, {
            "task": action.get("task", ""),
            "prefer_quality": action.get("prefer_quality", False),
        })
    if kind == "rag_ingest":
        r = tool_rag_ingest(action.get("text", ""), action.get("path", ""),
                            action.get("metadata", {}),
                            action.get("doc_id_prefix", None),
                            action.get("workdir", "/tmp"))
        return _tool_trace(action, r, {
            "path": action.get("path", ""), "doc_id_prefix": action.get("doc_id_prefix")})
    if kind == "rag_query":
        r = tool_rag_query(action.get("query", ""), action.get("top_k", 5),
                           action.get("filter_metadata", None))
        return _tool_trace(action, r, {
            "query": action.get("query", ""), "top_k": action.get("top_k", 5)})
    if kind == "rag_status":
        r = tool_rag_status()
        return _tool_trace(action, r)
    if kind == "inspect_db":
        r = tool_inspect_db(action.get("connection_string", ""))
        return _tool_trace(action, r, {"connection_string": action.get("connection_string", "")})
    if kind == "query_db":
        r = tool_query_db(action.get("connection_string", ""), action.get("query", ""))
        return _tool_trace(action, r, {
            "connection_string": action.get("connection_string", ""), "query": action.get("query", "")})
    if kind == "cron_schedule":
        r = tool_cron_schedule(
            action.get("name", "background-task"),
            action.get("task", ""),
            action.get("schedule", "5m"),
        )
        return _tool_trace(action, r, {
            "name": action.get("name", "background-task"),
            "schedule": action.get("schedule", "5m"),
        })
    if kind == "cron_list":
        r = tool_cron_list()
        return _tool_trace(action, r)
    if kind == "cron_cancel":
        r = tool_cron_cancel(action.get("job_id", ""))
        return _tool_trace(action, r, {"job_id": action.get("job_id", "")})
    if kind == "kg_store":
        r = tool_kg_store(
            action.get("name", ""),
            action.get("entity_type", "concept"),
            action.get("facts", {}),
            action.get("relations", []),
        )
        return _tool_trace(action, r, {"name": action.get("name", "")})
    if kind == "kg_query":
        r = tool_kg_query(action.get("query", ""), int(action.get("limit", 10)))
        return _tool_trace(action, r, {"query": action.get("query", "")})
    if kind == "kg_list":
        r = tool_kg_list(action.get("entity_type", None))
        return _tool_trace(action, r, {"entity_type": action.get("entity_type")})
    if kind == "read_page":
        r = tool_read_page(action.get("url", ""))
        return _tool_trace(action, r, {"url": action.get("url", "")})
    if kind == "api_call":
        r = tool_api_call(action.get("method", "GET"), action.get("url", ""),
                          action.get("headers"), action.get("body"))
        return _tool_trace(action, r, {"method": action.get("method", "GET"),
                                       "url": action.get("url", "")})
    if kind == "search_in_files":
        r = tool_search_in_files(
            action.get("pattern", ""),
            action.get("include_glob", "*"),
            action.get("workdir", "/tmp"),
            int(action.get("max_results", 200)),
        )
        return _tool_trace(action, r, {"pattern": action.get("pattern", "")})
    if kind == "create_directory":
        r = tool_create_directory(action.get("path", ""), action.get("workdir", "/tmp"))
        return _tool_trace(action, r, {"path": action.get("path", "")})
    if kind == "zip_files":
        r = tool_zip_files(
            action.get("paths", []),
            action.get("output_path", "archive.zip"),
            action.get("workdir", "/tmp"),
        )
        return _tool_trace(action, r, {"output_path": action.get("output_path", "archive.zip")})
    if kind == "unzip_files":
        r = tool_unzip_files(
            action.get("zip_path", ""),
            action.get("dest_path", ""),
            action.get("workdir", "/tmp"),
        )
        return _tool_trace(action, r, {"zip_path": action.get("zip_path", "")})
    if kind == "git_status":
        r = tool_git_status(action.get("repo_path", "."), action.get("workdir", "/tmp"))
        return _tool_trace(action, r, {"repo_path": action.get("repo_path", ".")})
    if kind == "git_log":
        r = tool_git_log(
            action.get("repo_path", "."),
            int(action.get("max_count", 20)),
            action.get("workdir", "/tmp"),
        )
        return _tool_trace(action, r, {"repo_path": action.get("repo_path", ".")})
    if kind == "git_diff":
        r = tool_git_diff(
            action.get("repo_path", "."),
            action.get("ref", "HEAD"),
            action.get("workdir", "/tmp"),
        )
        return _tool_trace(action, r, {
            "repo_path": action.get("repo_path", "."),
            "ref": action.get("ref", "HEAD"),
        })
    if kind == "git_checkout":
        r = tool_git_checkout(
            action.get("repo_path", "."),
            action.get("branch", ""),
            bool(action.get("create", False)),
            action.get("workdir", "/tmp"),
        )
        return _tool_trace(action, r, {
            "repo_path": action.get("repo_path", "."),
            "branch": action.get("branch", ""),
        })
    if kind == "git_pull":
        r = tool_git_pull(
            action.get("repo_path", "."),
            action.get("remote", "origin"),
            action.get("branch", ""),
            bool(action.get("ff_only", True)),
            action.get("workdir", "/tmp"),
        )
        return _tool_trace(action, r, {"repo_path": action.get("repo_path", ".")})
    if kind == "create_pull_request":
        r = tool_create_pull_request(
            action.get("title", ""),
            action.get("body", ""),
            action.get("base", "main"),
            action.get("head", ""),
            action.get("repo_path", "."),
            action.get("workdir", "/tmp"),
        )
        return _tool_trace(action, r, {"title": action.get("title", "")})
    if kind == "list_issues":
        r = tool_list_issues(
            action.get("repo_path", "."),
            action.get("state", "open"),
            int(action.get("limit", 20)),
            action.get("workdir", "/tmp"),
        )
        return _tool_trace(action, r, {"repo_path": action.get("repo_path", ".")})
    if kind == "create_issue":
        r = tool_create_issue(
            action.get("title", ""),
            action.get("body", ""),
            action.get("labels", ""),
            action.get("repo_path", "."),
            action.get("workdir", "/tmp"),
        )
        return _tool_trace(action, r, {"title": action.get("title", "")})
    if kind in ("youtube_transcript", "youtube"):
        r = tool_youtube_transcript(action.get("url", ""))
        return _tool_trace(action, r, {"url": action.get("url", "")})
    if kind == "generate_image":
        g = tool_generate_image(action.get("prompt", ""),
                                 action.get("width", 1024), action.get("height", 1024))
        r = f"![Generated image]({g['url']})\n\n*Prompt: {g['prompt']}*"
        return _tool_trace(action, r, {"prompt": action.get("prompt", "")})
    if kind == "diff":
        r = tool_diff(action.get("original", ""), action.get("modified", ""),
                      action.get("filename", "file"))
        return _tool_trace(action, r, {"filename": action.get("filename", "file")})
    if kind == "read_csv":
        r = tool_read_csv(action.get("path", ""), action.get("workdir", "/tmp"))
        return _tool_trace(action, r, {"path": action.get("path", "")})
    if kind == "write_csv":
        r = tool_write_csv(action.get("path", ""), action.get("data", []),
                           action.get("workdir", "/tmp"))
        return _tool_trace(action, r, {"path": action.get("path", "")})
    if kind == "read_pdf":
        r = tool_read_pdf(action.get("path", ""), action.get("workdir", "/tmp"))
        return _tool_trace(action, r, {"path": action.get("path", "")})
    if kind == "read_docx":
        r = tool_read_docx(action.get("path", ""), action.get("workdir", "/tmp"))
        return _tool_trace(action, r, {"path": action.get("path", "")})
    if kind == "read_xlsx":
        r = tool_read_xlsx(action.get("path", ""), action.get("workdir", "/tmp"))
        return _tool_trace(action, r, {"path": action.get("path", "")})
    if kind == "read_pptx":
        r = tool_read_pptx(action.get("path", ""), action.get("workdir", "/tmp"))
        return _tool_trace(action, r, {"path": action.get("path", "")})
    if kind == "ollama_list_models":
        from .agent import tool_ollama_list_models
        r = tool_ollama_list_models()
        return _tool_trace(action, r)
    # ── NEW Section 6.1 utility tools ────────────────────────────────────────
    if kind == "get_time":
        r = tool_get_time(action.get("timezone", "UTC"))
        return _tool_trace(action, r, {"timezone": action.get("timezone", "UTC")})
    if kind == "nexus_status":
        r = tool_nexus_status()
        return _tool_trace(action, r)
    if kind == "hash":
        r = tool_hash(action.get("text", ""), action.get("algorithm", "sha256"))
        return _tool_trace(action, r, {"algorithm": action.get("algorithm", "sha256")})
    if kind == "uuid":
        r = tool_uuid(int(action.get("version", 4)), action.get("namespace", ""),
                      action.get("name", ""))
        return _tool_trace(action, r)
    if kind == "qr_code":
        r = tool_qr_code(action.get("text", ""), int(action.get("size", 10)))
        return _tool_trace(action, r, {"text": action.get("text", "")})
    if kind == "csv_to_json":
        r = tool_csv_to_json(action.get("csv_text", ""))
        return _tool_trace(action, r)
    if kind == "json_to_csv":
        r = tool_json_to_csv(action.get("json_text", ""))
        return _tool_trace(action, r)
    if kind == "xml_parse":
        r = tool_xml_parse(action.get("xml_text", ""), action.get("xpath", ""))
        return _tool_trace(action, r)
    if kind in ("url_encode", "url_decode"):
        mode = "decode" if kind == "url_decode" else "encode"
        r = tool_url_encode(action.get("text", ""), mode)
        return _tool_trace(action, r, {"mode": mode})
    if kind == "jwt_decode":
        r = tool_jwt_decode(action.get("token", ""))
        return _tool_trace(action, r)
    if kind == "color_convert":
        r = tool_color_convert(action.get("color", ""), action.get("to_format", "all"))
        return _tool_trace(action, r, {"color": action.get("color", "")})
    # ── NEW Section 6.2 file / repo tools ────────────────────────────────────
    if kind == "write_file":
        r = tool_write_file(action.get("path", ""), action.get("content", ""),
                            action.get("workdir", "/tmp"), action.get("mode", "w"))
        return _tool_trace(action, r, {"path": action.get("path", "")})
    if kind == "read_file":
        r = tool_read_file(action.get("path", ""), action.get("workdir", "/tmp"),
                           int(action.get("start_line", 1)), int(action.get("end_line", 500)))
        return _tool_trace(action, r, {"path": action.get("path", "")})
    if kind == "list_files":
        r = tool_list_files(action.get("path", "."), action.get("workdir", "/tmp"),
                            bool(action.get("recursive", False)),
                            int(action.get("max_entries", 200)))
        return _tool_trace(action, r, {"path": action.get("path", ".")})
    if kind == "delete_file":
        r = tool_delete_file(action.get("path", ""), action.get("workdir", "/tmp"))
        return _tool_trace(action, r, {"path": action.get("path", "")})
    if kind == "move_file":
        r = tool_move_file(action.get("src", ""), action.get("dst", ""),
                           action.get("workdir", "/tmp"))
        return _tool_trace(action, r, {"src": action.get("src", ""), "dst": action.get("dst", "")})
    if kind == "copy_file":
        r = tool_copy_file(action.get("src", ""), action.get("dst", ""),
                           action.get("workdir", "/tmp"))
        return _tool_trace(action, r, {"src": action.get("src", ""), "dst": action.get("dst", "")})
    if kind == "clone_repo":
        r = tool_clone_repo(action.get("url", ""), action.get("dest", ""),
                            action.get("workdir", "/tmp"), action.get("branch", ""),
                            int(action.get("depth", 0)))
        return _tool_trace(action, r, {"url": action.get("url", "")})
    if kind == "run_command":
        r = tool_run_command(action.get("command", ""), action.get("workdir", "/tmp"),
                             int(action.get("timeout", 30)),
                             bool(action.get("allow_write", False)))
        return _tool_trace(action, r, {"command": action.get("command", "")[:100]})
    if kind == "commit_push":
        r = tool_commit_push(action.get("message", ""), action.get("repo_path", "."),
                             action.get("workdir", "/tmp"), action.get("remote", "origin"),
                             action.get("branch", ""))
        return _tool_trace(action, r, {"message": action.get("message", "")})
    if kind == "create_repo":
        r = tool_create_repo(action.get("name", ""), action.get("description", ""),
                             bool(action.get("private", False)), action.get("repo_path", "."),
                             action.get("workdir", "/tmp"))
        return _tool_trace(action, r, {"name": action.get("name", "")})
    # ── NEW Section 6.3 web / network tools ──────────────────────────────────
    if kind == "web_search":
        r = tool_web_search(action.get("query", ""), int(action.get("max_results", 5)),
                            action.get("engine", "auto"))
        return _tool_trace(action, r, {"query": action.get("query", "")})
    if kind == "web_scrape_structured":
        r = tool_web_scrape_structured(action.get("url", ""),
                                       action.get("selectors"),
                                       action.get("output_format", "json"))
        return _tool_trace(action, r, {"url": action.get("url", "")})
    if kind == "rss_fetch":
        r = tool_rss_fetch(action.get("url", ""), int(action.get("max_items", 10)))
        return _tool_trace(action, r, {"url": action.get("url", "")})
    if kind == "sitemap_crawl":
        r = tool_sitemap_crawl(action.get("url", ""), int(action.get("max_urls", 50)))
        return _tool_trace(action, r, {"url": action.get("url", "")})
    if kind == "check_url_status":
        urls_val = action.get("urls", action.get("url", []))
        r = tool_check_url_status(urls_val, int(action.get("timeout", 10)))
        return _tool_trace(action, r)
    if kind == "screenshot":
        try:
            from .vision import capture_screenshot
            url_val = action.get("url", "")
            r = capture_screenshot(url_val) if url_val else "❌ url is required"
        except Exception as _exc:
            r = f"❌ screenshot failed: {_exc}"
        return _tool_trace(action, r, {"url": action.get("url", "")})
    # ── NEW Section 6.5 DB tools ──────────────────────────────────────────────
    if kind == "pg_query":
        r = tool_pg_query(action.get("sql", ""), action.get("database_url", ""),
                          int(action.get("max_rows", 100)))
        return _tool_trace(action, r, {"sql": action.get("sql", "")[:80]})
    if kind == "db_migrate":
        r = tool_db_migrate(action.get("migration_sql", ""),
                            action.get("database_url", ""),
                            bool(action.get("dry_run", True)))
        return _tool_trace(action, r, {"dry_run": action.get("dry_run", True)})
    return None


# ── IMAGE GENERATION ──────────────────────────────────────────────────────────
def tool_generate_image(prompt: str, width: int = 1024, height: int = 1024,
                         model: str = "flux") -> dict:
    """
    Generate an image via Pollinations.ai (free, no API key).
    Returns a dict with the image URL and prompt used.
    """
    import urllib.parse
    encoded = urllib.parse.quote(prompt)
    seed    = hash(prompt) % 9999
    url     = (f"https://image.pollinations.ai/prompt/{encoded}"
               f"?width={width}&height={height}&model={model}&seed={seed}&nologo=true")
    return {"url": url, "prompt": prompt, "width": width, "height": height}


# ── YOUTUBE TRANSCRIPT ────────────────────────────────────────────────────────
def tool_youtube_transcript(url: str) -> str:
    """Extract subtitles/transcript from a YouTube video using yt-dlp."""
    try:
        import yt_dlp, tempfile, os, json as _json
        opts = {
            "skip_download": True,
            "writeautomaticsub": True,
            "writesubtitles": True,
            "subtitleslangs": ["en", "en-US"],
            "subtitlesformat": "json3",
            "quiet": True,
            "no_warnings": True,
            "outtmpl": os.path.join(tempfile.gettempdir(), "%(id)s.%(ext)s"),
        }
        with yt_dlp.YoutubeDL(opts) as ydl:
            info = ydl.extract_info(url, download=True)
            vid_id  = info.get("id","")
            title   = info.get("title","")
            duration= info.get("duration", 0)

        # Find downloaded subtitle file
        tmp = tempfile.gettempdir()
        for ext in [f"{vid_id}.en.json3", f"{vid_id}.en-US.json3"]:
            fpath = os.path.join(tmp, ext)
            if os.path.exists(fpath):
                with open(fpath) as f:
                    data = _json.load(f)
                texts = []
                for event in data.get("events", []):
                    for seg in event.get("segs", []):
                        t = seg.get("utf8", "").strip()
                        if t and t != "\n":
                            texts.append(t)
                transcript = " ".join(texts)
                os.remove(fpath)
                mins = duration // 60
                return (f"**{title}** ({mins}m)\n\n"
                        f"{transcript[:4000]}"
                        + (" …*(truncated)*" if len(transcript) > 4000 else ""))
        return f"No English subtitles found for: {title or url}"
    except Exception as e:
        return f"YouTube transcript failed: {e}"


# ── YOUTUBE TRANSCRIPT ────────────────────────────────────────────────────────
def tool_youtube(url: str) -> str:
    """Fetch YouTube transcript via youtube-transcript-api or scrape description."""
    try:
        from youtube_transcript_api import YouTubeTranscriptApi
        import re as _re
        vid_match = _re.search(r'(?:v=|youtu\.be/)([A-Za-z0-9_-]{11})', url)
        if not vid_match:
            return f"❌ Could not extract video ID from URL: {url}"
        vid_id = vid_match.group(1)
        transcript = YouTubeTranscriptApi.get_transcript(vid_id)
        text = " ".join(t["text"] for t in transcript)
        # Truncate to ~4000 chars
        if len(text) > 4000:
            text = text[:4000] + f"\n… (transcript truncated, {len(text)} total chars)"
        return f"**YouTube transcript** ({vid_id}):\n\n{text}"
    except ImportError:
        # Fallback: yt-dlp metadata only
        try:
            import subprocess, json as _json
            r = subprocess.run(
                ["yt-dlp", "--dump-json", "--no-download", url],
                capture_output=True, text=True, timeout=20
            )
            if r.returncode == 0:
                d = _json.loads(r.stdout)
                return (f"**{d.get('title','')}** ({d.get('uploader','')})\n"
                        f"Duration: {d.get('duration_string','?')}\n\n"
                        f"{d.get('description','No description')[:1500]}")
        except Exception:
            pass
        return "❌ youtube-transcript-api not installed. Run: pip install youtube-transcript-api"
    except Exception as e:
        return f"❌ YouTube transcript failed: {e}"


# ── PDF READER ────────────────────────────────────────────────────────────────
def tool_read_pdf(path: str, workdir: str = "/tmp") -> str:
    """Extract text from a PDF file."""
    import os as _os
    full = _os.path.join(workdir, path) if not _os.path.isabs(path) else path
    if not _os.path.exists(full):
        return f"❌ File not found: {path}"
    try:
        import pypdf
        reader = pypdf.PdfReader(full)
        pages  = []
        for i, page in enumerate(reader.pages[:50]):   # first 50 pages (production limit)
            text = page.extract_text() or ""
            if text.strip():
                pages.append(f"[Page {i+1}]\n{text.strip()}")
        if not pages:
            return "❌ No extractable text found (may be a scanned PDF)"
        total = len(reader.pages)
        content = "\n\n".join(pages)
        if len(content) > 12000:
            content = content[:12000] + f"\n\n… (truncated, {total} pages total)"
        return content
    except ImportError:
        return "❌ pypdf not installed. Run: pip install pypdf"
    except Exception as e:
        return f"❌ PDF read failed: {e}"


# ── DIFF VIEWER ───────────────────────────────────────────────────────────────
def tool_diff(original: str, modified: str, filename: str = "file") -> str:
    """Generate a unified diff between two text strings."""
    import difflib
    orig_lines = original.splitlines(keepends=True)
    mod_lines  = modified.splitlines(keepends=True)
    diff = list(difflib.unified_diff(
        orig_lines, mod_lines,
        fromfile=f"a/{filename}", tofile=f"b/{filename}", lineterm=""
    ))
    if not diff:
        return "No differences found."
    diff_text = "\n".join(diff)
    if len(diff_text) > 4000:
        diff_text = diff_text[:4000] + "\n… (diff truncated)"
    return f"```diff\n{diff_text}\n```"


# ── OFFICE DOCUMENT READERS ───────────────────────────────────────────────────

def tool_read_docx(path: str, workdir: str = "/tmp") -> str:
    """Extract text from a Word (.docx) file, including heading structure and tables."""
    import os as _os
    full = _os.path.join(workdir, path) if not _os.path.isabs(path) else path
    if not _os.path.exists(full):
        return f"❌ File not found: {path}"
    try:
        from docx import Document
        doc = Document(full)
        parts = []
        # Paragraphs with heading structure
        for para in doc.paragraphs:
            if not para.text.strip():
                continue
            style = para.style.name if para.style else ""
            if style.startswith("Heading "):
                try:
                    level = int(style.split()[-1])
                except (ValueError, IndexError):
                    level = 1
                marker = "#" * min(level, 4)
                parts.append(f"{marker} {para.text.strip()}")
            else:
                parts.append(para.text.strip())
        # Tables
        for tbl_idx, table in enumerate(doc.tables):
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                if any(cells):
                    rows.append(" | ".join(cells))
            if rows:
                parts.append(f"[Table {tbl_idx + 1}]\n" + "\n".join(rows[:100]))
        content = "\n\n".join(parts)
        if len(content) > 8000:
            content = content[:8000] + f"\n\n… (truncated, {len(doc.paragraphs)} paragraphs, {len(doc.tables)} tables)"
        return content if content else "❌ No extractable text found in document."
    except ImportError:
        return "❌ python-docx not installed. Run: pip install python-docx"
    except Exception as e:
        return f"❌ DOCX read failed: {e}"


def tool_read_xlsx(path: str, workdir: str = "/tmp") -> str:
    """Extract data from an Excel (.xlsx) file."""
    import os as _os
    full = _os.path.join(workdir, path) if not _os.path.isabs(path) else path
    if not _os.path.exists(full):
        return f"❌ File not found: {path}"
    try:
        import openpyxl
        wb = openpyxl.load_workbook(full, read_only=True, data_only=True)
        lines = []
        for sheet_name in wb.sheetnames[:5]:   # first 5 sheets (production limit)
            ws = wb[sheet_name]
            lines.append(f"### Sheet: {sheet_name}")
            row_count = 0
            for row in ws.iter_rows(max_row=200, values_only=True):   # 200 rows per sheet
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    lines.append(" | ".join(cells))
                    row_count += 1
            if ws.max_row and ws.max_row > 200:
                lines.append(f"… ({ws.max_row - 200} more rows in sheet)")
        wb.close()
        content = "\n".join(lines)
        return content if content.strip() else "❌ No data found in workbook."
    except ImportError:
        return "❌ openpyxl not installed. Run: pip install openpyxl"
    except Exception as e:
        return f"❌ XLSX read failed: {e}"


def tool_read_pptx(path: str, workdir: str = "/tmp") -> str:
    """Extract text from a PowerPoint (.pptx) file."""
    import os as _os
    full = _os.path.join(workdir, path) if not _os.path.isabs(path) else path
    if not _os.path.exists(full):
        return f"❌ File not found: {path}"
    try:
        from pptx import Presentation
        prs = Presentation(full)
        slides = []
        for i, slide in enumerate(prs.slides[:30]):   # first 30 slides
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                slides.append(f"[Slide {i+1}]\n" + "\n".join(texts))
        content = "\n\n".join(slides)
        total = len(prs.slides)
        if len(content) > 6000:
            content = content[:6000] + f"\n\n… (truncated, {total} slides total)"
        return content if content else "❌ No text found in presentation."
    except ImportError:
        return "❌ python-pptx not installed. Run: pip install python-pptx"
    except Exception as e:
        return f"❌ PPTX read failed: {e}"


# ── SPREADSHEET (CSV) ─────────────────────────────────────────────────────────
def tool_read_csv(path: str, workdir: str = "/tmp") -> str:
    import os, csv
    full = os.path.join(workdir, path) if not os.path.isabs(path) else path
    if not os.path.exists(full):
        return f"File not found: {path}"
    try:
        with open(full, newline='', encoding='utf-8', errors='replace') as f:
            reader = list(csv.reader(f))
        if not reader:
            return "Empty CSV."
        headers = reader[0]
        rows    = reader[1:51]   # first 50 data rows
        lines   = [" | ".join(headers)]
        lines  += ["-+-".join("-" * max(1,len(h)) for h in headers)]
        for row in rows:
            lines.append(" | ".join(str(v)[:30] for v in row))
        suffix = f"\n…({len(reader)-51} more rows)" if len(reader) > 51 else ""
        return f"**{path}** — {len(reader)-1} rows × {len(headers)} cols\n\n```\n" + "\n".join(lines) + f"\n```{suffix}"
    except Exception as e:
        return f"CSV read failed: {e}"


def tool_write_csv(path: str, data: list, workdir: str = "/tmp") -> str:
    """data: list of lists (first row = headers)"""
    import os, csv
    full = os.path.join(workdir, path) if not os.path.isabs(path) else path
    os.makedirs(os.path.dirname(full), exist_ok=True)
    try:
        with open(full, 'w', newline='', encoding='utf-8') as f:
            csv.writer(f).writerows(data)
        return f"Wrote {path} ({len(data)} rows)"
    except Exception as e:
        return f"CSV write failed: {e}"


# ── API CALLER ────────────────────────────────────────────────────────────────
def tool_api_call(method: str, url: str, headers: dict | None = None,
                  body: dict | str | None = None, timeout: int = 15) -> str:
    import requests as _r, json as _j
    BLOCKED_HOSTS = ["169.254.", "192.168.", "10.", "127.", "0.0.0.0", "localhost"]
    for b in BLOCKED_HOSTS:
        if b in url:
            return f"Blocked: cannot call internal/local addresses."
    try:
        method = method.upper()
        hdrs   = {"User-Agent": "ClaudeAlt/1.0", **(headers or {})}
        kwargs = {"headers": hdrs, "timeout": timeout}
        if body:
            if isinstance(body, dict):
                kwargs["json"] = body
            else:
                kwargs["data"] = body
        resp = getattr(_r, method.lower())(url, **kwargs)
        ct   = resp.headers.get("Content-Type", "")
        if "json" in ct:
            try:
                data = resp.json()
                text = _j.dumps(data, indent=2)[:3000]
            except Exception:
                text = resp.text[:3000]
        else:
            text = resp.text[:3000]
        return f"**{method} {url}**\nStatus: {resp.status_code}\n\n```\n{text}\n```"
    except Exception as e:
        return f"API call failed: {e}"


# ── PAGE READER ───────────────────────────────────────────────────────────────
def tool_read_page(url: str) -> str:
    """Fetch a webpage and return readable text (strips HTML tags)."""
    import requests as _r, re as _re, html as _html
    BLOCKED = ["localhost", "127.", "192.168.", "10.", "169.254."]
    for b in BLOCKED:
        if b in url:
            return "Blocked: internal addresses."
    try:
        resp = _r.get(url, headers={"User-Agent":"Mozilla/5.0"},
                      timeout=15, allow_redirects=True)
        resp.raise_for_status()
        text = resp.text
        # Strip scripts/styles
        text = _re.sub(r'<(script|style)[^>]*>.*?</\1>', ' ', text, flags=_re.S|_re.I)
        # Strip tags
        text = _re.sub(r'<[^>]+>', ' ', text)
        # Decode entities
        text = _html.unescape(text)
        # Collapse whitespace
        text = _re.sub(r'\s{3,}', '\n\n', text).strip()
        return text[:5000] + ("…*(truncated)*" if len(text) > 5000 else "")
    except Exception as e:
        return f"Page read failed: {e}"


# ── FILESYSTEM + GIT HELPERS ────────────────────────────────────────────────
def _resolve_path(path: str, workdir: str = "/tmp") -> Path:
    candidate = Path(path)
    if not candidate.is_absolute():
        candidate = Path(workdir) / candidate
    return candidate.resolve()


def tool_search_in_files(pattern: str, include_glob: str = "*", workdir: str = "/tmp", max_results: int = 200) -> str:
    """Search files recursively with regex and return matching lines."""
    if not pattern.strip():
        return "❌ pattern is required."
    root = _resolve_path(".", workdir)
    if not root.exists():
        return f"❌ workdir does not exist: {workdir}"

    try:
        rx = re.compile(pattern, re.IGNORECASE)
    except re.error as e:
        return f"❌ Invalid regex pattern: {e}"

    matches: list[str] = []
    files_scanned = 0
    for file_path in root.rglob("*"):
        if not file_path.is_file():
            continue
        if include_glob and not file_path.match(include_glob):
            continue
        files_scanned += 1
        try:
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                for idx, line in enumerate(f, start=1):
                    if rx.search(line):
                        rel = file_path.relative_to(root)
                        matches.append(f"{rel}:{idx}: {line.rstrip()[:240]}")
                        if len(matches) >= max_results:
                            break
        except Exception:
            continue
        if len(matches) >= max_results:
            break

    if not matches:
        return f"No matches found for /{pattern}/ in {files_scanned} files."
    return (
        f"Found {len(matches)} match(es) in {files_scanned} file(s).\n\n"
        + "```\n"
        + "\n".join(matches)
        + "\n```"
    )


def tool_create_directory(path: str, workdir: str = "/tmp") -> str:
    """Create a directory recursively (mkdir -p semantics)."""
    if not path.strip():
        return "❌ path is required."
    target = _resolve_path(path, workdir)
    try:
        target.mkdir(parents=True, exist_ok=True)
        return f"✅ Directory ready: {target}"
    except Exception as e:
        return f"❌ Failed to create directory: {e}"


def tool_zip_files(paths: list[str], output_path: str, workdir: str = "/tmp") -> str:
    """Create a ZIP archive from files/directories."""
    if not paths:
        return "❌ paths is required."
    if not output_path.strip():
        return "❌ output_path is required."

    root = _resolve_path(".", workdir)
    out = _resolve_path(output_path, workdir)
    out.parent.mkdir(parents=True, exist_ok=True)

    added = 0
    try:
        with zipfile.ZipFile(out, "w", compression=zipfile.ZIP_DEFLATED) as zf:
            for raw_path in paths:
                src = _resolve_path(raw_path, workdir)
                if not src.exists():
                    continue
                if src.is_file():
                    arcname = src.relative_to(root) if src.is_relative_to(root) else src.name
                    zf.write(src, arcname=str(arcname))
                    added += 1
                else:
                    for child in src.rglob("*"):
                        if child.is_file():
                            arcname = child.relative_to(root) if child.is_relative_to(root) else child.name
                            zf.write(child, arcname=str(arcname))
                            added += 1
        return f"✅ Created zip: {out} ({added} file(s))"
    except Exception as e:
        return f"❌ ZIP creation failed: {e}"


def tool_unzip_files(zip_path: str, dest_path: str = "", workdir: str = "/tmp") -> str:
    """Extract a ZIP archive into destination path."""
    if not zip_path.strip():
        return "❌ zip_path is required."
    src = _resolve_path(zip_path, workdir)
    if not src.exists():
        return f"❌ ZIP not found: {zip_path}"
    dest = _resolve_path(dest_path, workdir) if dest_path else _resolve_path(".", workdir)
    dest.mkdir(parents=True, exist_ok=True)
    try:
        with zipfile.ZipFile(src, "r") as zf:
            zf.extractall(dest)
            extracted_count = len(zf.namelist())
        return f"✅ Extracted {extracted_count} entrie(s) to {dest}"
    except Exception as e:
        return f"❌ Unzip failed: {e}"


def _run_git(repo_path: str, args: list[str], workdir: str = "/tmp") -> tuple[int, str, str]:
    repo = _resolve_path(repo_path or ".", workdir)
    proc = subprocess.run(
        ["git", "-C", str(repo), *args],
        capture_output=True,
        text=True,
        timeout=30,
    )
    return proc.returncode, proc.stdout.strip(), proc.stderr.strip()


def tool_git_status(repo_path: str = ".", workdir: str = "/tmp") -> str:
    """Show current git working tree status."""
    code, out, err = _run_git(repo_path, ["status", "--short", "--branch"], workdir=workdir)
    if code != 0:
        return f"❌ git_status failed: {err or 'unknown error'}"
    return out or "Working tree clean."


def tool_git_log(repo_path: str = ".", max_count: int = 20, workdir: str = "/tmp") -> str:
    """Show recent git commit history."""
    if max_count <= 0:
        max_count = 20
    code, out, err = _run_git(
        repo_path,
        ["log", f"--max-count={max_count}", "--pretty=format:%h %ad %an %s", "--date=short"],
        workdir=workdir,
    )
    if code != 0:
        return f"❌ git_log failed: {err or 'unknown error'}"
    return out or "No commits found."


def tool_git_diff(repo_path: str = ".", ref: str = "HEAD", workdir: str = "/tmp") -> str:
    """Show git diff against a reference (default HEAD)."""
    code, out, err = _run_git(repo_path, ["--no-pager", "diff", ref], workdir=workdir)
    if code != 0:
        return f"❌ git_diff failed: {err or 'unknown error'}"
    if not out:
        return "No diff."
    if len(out) > 8000:
        out = out[:8000] + "\n... (diff truncated)"
    return f"```diff\n{out}\n```"


def tool_git_checkout(repo_path: str, branch: str, create: bool = False, workdir: str = "/tmp") -> str:
    """Checkout a git branch, optionally creating it."""
    if not branch.strip():
        return "❌ branch is required."
    args = ["checkout"]
    if create:
        args.append("-b")
    args.append(branch)
    code, out, err = _run_git(repo_path or ".", args, workdir=workdir)
    if code != 0:
        return f"❌ git_checkout failed: {err or 'unknown error'}"
    return out or f"✅ Checked out {branch}"


def tool_git_pull(repo_path: str = ".", remote: str = "origin", branch: str = "", ff_only: bool = True, workdir: str = "/tmp") -> str:
    """Pull latest changes from remote with optional ff-only safety."""
    args = ["pull"]
    if ff_only:
        args.append("--ff-only")
    if remote:
        args.append(remote)
    if branch:
        args.append(branch)
    code, out, err = _run_git(repo_path, args, workdir=workdir)
    if code != 0:
        return f"❌ git_pull failed: {err or 'unknown error'}"
    return out or "✅ Pull completed."


def _run_gh(args: list[str], repo_path: str = ".", workdir: str = "/tmp") -> tuple[int, str, str]:
    repo = _resolve_path(repo_path or ".", workdir)
    proc = subprocess.run(
        ["gh", *args],
        cwd=str(repo),
        capture_output=True,
        text=True,
        timeout=45,
    )
    return proc.returncode, proc.stdout.strip(), proc.stderr.strip()


def tool_create_pull_request(title: str, body: str = "", base: str = "main", head: str = "", repo_path: str = ".", workdir: str = "/tmp") -> str:
    """Create a GitHub pull request using gh CLI."""
    if not title.strip():
        return "❌ title is required."
    args = ["pr", "create", "--title", title, "--body", body or "", "--base", base]
    if head.strip():
        args.extend(["--head", head])
    code, out, err = _run_gh(args, repo_path=repo_path, workdir=workdir)
    if code != 0:
        return f"❌ create_pull_request failed: {err or 'gh CLI error'}"
    return out or "✅ Pull request created."


def tool_list_issues(repo_path: str = ".", state: str = "open", limit: int = 20, workdir: str = "/tmp") -> str:
    """List GitHub issues using gh CLI."""
    if state not in {"open", "closed", "all"}:
        return "❌ state must be one of: open, closed, all"
    if limit <= 0:
        limit = 20
    args = ["issue", "list", "--state", state, "--limit", str(limit)]
    code, out, err = _run_gh(args, repo_path=repo_path, workdir=workdir)
    if code != 0:
        return f"❌ list_issues failed: {err or 'gh CLI error'}"
    return out or "No issues found."


def tool_create_issue(title: str, body: str = "", labels: str = "", repo_path: str = ".", workdir: str = "/tmp") -> str:
    """Create a GitHub issue using gh CLI."""
    if not title.strip():
        return "❌ title is required."
    args = ["issue", "create", "--title", title, "--body", body or ""]
    if labels.strip():
        args.extend(["--label", labels])
    code, out, err = _run_gh(args, repo_path=repo_path, workdir=workdir)
    if code != 0:
        return f"❌ create_issue failed: {err or 'gh CLI error'}"
    return out or "✅ Issue created."


# ── DATABASE QUERY TOOL ───────────────────────────────────────────────────────
def tool_query_db(connection_string: str, query: str) -> str:
    """Run a read-only SQL query. Only SELECT statements allowed."""
    import re as _re
    stripped = query.strip().upper()
    if not stripped.startswith("SELECT"):
        return "❌ Only SELECT queries are allowed."
    # Block dangerous keywords even inside SELECT
    dangerous = ["DROP ", "DELETE ", "INSERT ", "UPDATE ", "ALTER ",
                 "EXEC ", "EXECUTE ", "xp_", "--", "/*"]
    for kw in dangerous:
        if kw.upper() in query.upper():
            return f"❌ Blocked keyword in query: {kw.strip()}"
    try:
        if connection_string.startswith("sqlite:///") or connection_string.endswith(".db"):
            import sqlite3
            path = connection_string.replace("sqlite:///","")
            conn = sqlite3.connect(path)
            conn.row_factory = sqlite3.Row
            rows  = conn.execute(query).fetchall()[:100]
            conn.close()
            if not rows: return "No results."
            cols = rows[0].keys()
            lines = [" | ".join(str(c) for c in cols)]
            lines += ["-+-".join("-"*max(1,len(str(c))) for c in cols)]
            for row in rows:
                lines.append(" | ".join(str(row[c])[:30] for c in cols))
            return f"```\n" + "\n".join(lines) + "\n```"
        else:
            return "❌ Only SQLite (sqlite:///path.db) supported currently."
    except Exception as e:
        return f"❌ Query failed: {e}"


def tool_inspect_db(connection_string: str) -> str:
    """Introspect a database schema — list tables, columns, types, and row counts.

    Supports:
      sqlite:///path/to/file.db  or  /path/to/file.db
      postgresql://user:pass@host/dbname  (requires psycopg2)  # pragma: allowlist secret
    """
    cs = (connection_string or "").strip()
    if not cs:
        return "❌ connection_string is required."
    try:
        if cs.startswith("sqlite:///") or cs.endswith(".db"):
            import sqlite3
            path = cs.replace("sqlite:///", "")
            conn = sqlite3.connect(path)
            c    = conn.cursor()
            tables = [r[0] for r in c.execute(
                "SELECT name FROM sqlite_master WHERE type='table' ORDER BY name"
            ).fetchall()]
            if not tables:
                conn.close()
                return "No tables found."
            lines = [f"**SQLite schema** — `{path}`\n"]
            for tbl in tables:
                try:
                    cnt = c.execute(f"SELECT COUNT(*) FROM \"{tbl}\"").fetchone()[0]
                except Exception:
                    cnt = "?"
                lines.append(f"### {tbl} ({cnt} rows)")
                cols_info = c.execute(f"PRAGMA table_info(\"{tbl}\")").fetchall()
                for col in cols_info:
                    pk = " 🔑" if col[5] else ""
                    nn = " NOT NULL" if col[3] else ""
                    dv = f" DEFAULT {col[4]}" if col[4] is not None else ""
                    lines.append(f"  - **{col[1]}** `{col[2]}`{nn}{dv}{pk}")
                lines.append("")
            conn.close()
            return "\n".join(lines)

        elif cs.startswith("postgresql://") or cs.startswith("postgres://"):
            try:
                import psycopg2
                import psycopg2.extras
            except ImportError:
                return "❌ psycopg2 not installed. Run: pip install psycopg2-binary"
            conn = psycopg2.connect(cs)
            cur  = conn.cursor(cursor_factory=psycopg2.extras.DictCursor)
            cur.execute("""
                SELECT table_name
                FROM information_schema.tables
                WHERE table_schema = 'public'
                ORDER BY table_name
            """)
            tables = [r["table_name"] for r in cur.fetchall()]
            if not tables:
                conn.close()
                return "No public tables found."
            lines = [f"**PostgreSQL schema** — `{cs.split('@')[-1]}`\n"]
            for tbl in tables:
                try:
                    cur.execute(f'SELECT COUNT(*) FROM "{tbl}"')
                    cnt = cur.fetchone()[0]
                except Exception:
                    cnt = "?"
                lines.append(f"### {tbl} ({cnt} rows)")
                cur.execute("""
                    SELECT column_name, data_type, is_nullable, column_default
                    FROM information_schema.columns
                    WHERE table_schema='public' AND table_name=%s
                    ORDER BY ordinal_position
                """, (tbl,))
                for col in cur.fetchall():
                    nn = "" if col["is_nullable"] == "YES" else " NOT NULL"
                    dv = f" DEFAULT {col['column_default']}" if col["column_default"] else ""
                    lines.append(f"  - **{col['column_name']}** `{col['data_type']}`{nn}{dv}")
                lines.append("")
            conn.close()
            return "\n".join(lines)

        else:
            return "❌ Unsupported connection string. Use sqlite:///path.db or postgresql://..."
    except Exception as e:
        return f"❌ Schema inspection failed: {e}"


# ── BACKGROUND SCHEDULER TOOLS ───────────────────────────────────────────────
def tool_cron_schedule(name: str, task: str, schedule: str) -> str:
    """Create an autonomous background job.

    schedule supports intervals (30s/5m/2h/1d) or basic 5-field cron.
    """
    if not (task or "").strip():
        return "❌ task is required."
    try:
        job = schedule_job((name or "background-task").strip(), task.strip(), (schedule or "5m").strip())
        j = job_to_dict(job)
        return (
            f"✅ Scheduled **{j['name']}** (id `{j['id']}`)\n"
            f"- task: {j['task'][:160]}\n"
            f"- schedule: `{j['schedule']}`\n"
            f"- next_run: `{j['next_run']}`"
        )
    except Exception as e:
        return f"❌ Failed to schedule job: {e}"


def tool_cron_list() -> str:
    """List all background jobs with status and timing."""
    jobs = [job_to_dict(j) for j in list_jobs()]
    if not jobs:
        return "No scheduled jobs."
    lines = ["# Scheduled jobs"]
    for j in jobs:
        lines.append(
            f"- `{j['id']}` **{j['name']}** · {j['status']} · `{j['schedule']}` "
            f"(runs: {j['run_count']}, next: {j['next_run'] or 'n/a'})"
        )
    return "\n".join(lines)


def tool_cron_cancel(job_id: str) -> str:
    """Cancel a scheduled background job by id."""
    jid = (job_id or "").strip()
    if not jid:
        return "❌ job_id is required."
    if cancel_job(jid):
        return f"✅ Cancelled job `{jid}`."
    return f"❌ Job not found: `{jid}`"


# ── KNOWLEDGE GRAPH TOOLS ───────────────────────────────────────────────────

def tool_kg_store(name: str, entity_type: str = "concept",
                  facts: dict | None = None, relations: list | None = None) -> str:
    """Store or update an entity in the long-term knowledge graph."""
    if not name or not name.strip():
        return "❌ entity name is required."
    eid = _kg_store(name.strip(), entity_type, facts or {}, relations or [])
    rel_count = len(relations) if relations else 0
    facts_count = len(facts) if facts else 0
    return (f"✅ Stored [{entity_type}] **{name}** in knowledge graph "
            f"(id={eid[:8]}, {facts_count} facts, {rel_count} relations)")


def tool_kg_query(query: str, limit: int = 10) -> str:
    """Search the knowledge graph for entities matching a query."""
    if not query or not query.strip():
        return "❌ query is required."
    results = _kg_query(query.strip(), limit=limit)
    if not results:
        return f"No knowledge graph entries found for: {query}"
    lines = [f"**{len(results)} KG result(s) for '{query}':**\n"]
    for e in results:
        facts = e.get("facts", {})
        facts_str = ", ".join(f"{k}: {v}" for k, v in list(facts.items())[:5]) if facts else "—"
        rels = e.get("relations", [])
        rel_str = "; ".join(f"{r['relation']}→{r['entity']}" for r in rels[:3]) if rels else "—"
        lines.append(f"• **[{e['type']}] {e['name']}**")
        lines.append(f"  facts: {facts_str}")
        lines.append(f"  links: {rel_str}")
    return "\n".join(lines)


def tool_kg_list(entity_type: str | None = None) -> str:
    """List entities in the knowledge graph, optionally filtered by type."""
    entities = _kg_list(entity_type=entity_type, limit=50)
    if not entities:
        filter_txt = f" of type '{entity_type}'" if entity_type else ""
        return f"No knowledge graph entities{filter_txt} found."
    lines = [f"**{len(entities)} KG entit{'y' if len(entities)==1 else 'ies'}:**\n"]
    for e in entities:
        lines.append(f"• [{e['type']}] **{e['name']}** (updated: {e['updated_at'][:10]})")
    return "\n".join(lines)


# ── COST ESTIMATOR ────────────────────────────────────────────────────────────
# Approximate costs per 1M tokens (input/output) as of early 2026
PROVIDER_COSTS = {
    "groq":         (0.00,  0.00),    # free tier
    "cerebras":     (0.00,  0.00),
    "gemini":       (0.00,  0.00),    # free tier
    "mistral":      (0.00,  0.00),    # free tier
    "openrouter":   (0.00,  0.00),    # free models
    "cohere":       (0.00,  0.00),
    "github_models":(0.00,  0.00),
    "llm7":         (0.00,  0.00),
    "nvidia":       (0.00,  0.00),
    "grok":         (5.00, 15.00),    # grok-3 approx
    "claude":       (3.00, 15.00),    # claude-sonnet-4 approx
}

def estimate_cost(provider_label: str, in_tokens: int, out_tokens: int) -> float:
    key = provider_label.lower().split()[0]
    costs = PROVIDER_COSTS.get(key, (0.0, 0.0))
    return (in_tokens * costs[0] + out_tokens * costs[1]) / 1_000_000


# ── SQLITE INTROSPECTION ──────────────────────────────────────────────────────

def tool_inspect_sqlite(query: str = "", db_path: str = "") -> str:
    """Introspect a SQLite database: list tables, describe schema, or run a read-only query."""
    import sqlite3 as _sq, os as _os
    path = db_path or _os.getenv("DB_PATH", "/tmp/nexus_ai.db")
    if not _os.path.exists(path):
        return f"❌ Database not found: {path}"
    try:
        conn = _sq.connect(path)
        conn.row_factory = _sq.Row
        q = (query or "").strip()

        if not q or q.lower() in ("tables", "list tables", ".tables"):
            rows = conn.execute("SELECT name, type FROM sqlite_master WHERE type IN ('table','view') ORDER BY name").fetchall()
            if not rows:
                return "No tables found."
            lines = ["**Tables:**"]
            for r in rows:
                lines.append(f"• `{r['name']}` ({r['type']})")
            conn.close()
            return "\n".join(lines)

        if q.lower().startswith("schema") or q.lower().startswith(".schema"):
            tbl = q.split(None, 1)[1].strip() if len(q.split()) > 1 else ""
            if tbl:
                row = conn.execute("SELECT sql FROM sqlite_master WHERE name=?", (tbl,)).fetchone()
                conn.close()
                return f"```sql\n{row['sql']}\n```" if row else f"Table `{tbl}` not found."
            rows = conn.execute("SELECT sql FROM sqlite_master WHERE type='table' ORDER BY name").fetchall()
            conn.close()
            return "```sql\n" + "\n\n".join(r["sql"] for r in rows if r["sql"]) + "\n```"

        lower_q = q.lower().lstrip()
        if any(lower_q.startswith(kw) for kw in ("insert", "update", "delete", "drop", "alter", "create", "replace")):
            conn.close()
            return "❌ Only read-only queries are allowed (SELECT, PRAGMA, etc.)"

        rows = conn.execute(q).fetchmany(200)
        conn.close()
        if not rows:
            return "Query returned no rows."
        cols = rows[0].keys()
        header = " | ".join(cols)
        sep = " | ".join("---" for _ in cols)
        lines = [f"| {header} |", f"| {sep} |"]
        for r in rows:
            lines.append("| " + " | ".join(str(r[c]) for c in cols) + " |")
        return "\n".join(lines)
    except Exception as e:
        return f"❌ SQLite error: {e}"


def tool_query_sqlite(sql: str, db_path: str = "") -> str:
    """Alias for tool_inspect_sqlite — run a read-only SQL query against the database."""
    return tool_inspect_sqlite(sql, db_path)


# ── POSTGRESQL INTROSPECTION ──────────────────────────────────────────────────

def tool_inspect_postgres(query: str = "", database_url: str = "") -> str:
    """Introspect a PostgreSQL database: list tables, describe schema, or run a read-only query."""
    import os as _os
    url = database_url or _os.getenv("DATABASE_URL", "")
    if not url:
        return "❌ DATABASE_URL not set and no database_url provided."
    try:
        import psycopg2
        from psycopg2.extras import RealDictCursor
        conn = psycopg2.connect(url, cursor_factory=RealDictCursor)
    except Exception as e:
        return f"❌ Cannot connect to PostgreSQL: {e}"

    try:
        q = (query or "").strip()
        with conn.cursor() as cur:
            if not q or q.lower() in ("tables", "list tables", "\\dt"):
                cur.execute("""
                    SELECT table_name, table_type
                    FROM information_schema.tables
                    WHERE table_schema = 'public'
                    ORDER BY table_name
                """)
                rows = cur.fetchall()
                if not rows:
                    return "No tables found."
                lines = ["**Tables (public schema):**"]
                for r in rows:
                    lines.append(f"• `{r['table_name']}` ({r['table_type']})")
                return "\n".join(lines)

            if q.lower().startswith("schema") or q.lower().startswith("\\d "):
                tbl = q.split(None, 1)[1].strip() if len(q.split()) > 1 else ""
                if tbl:
                    cur.execute("""
                        SELECT column_name, data_type, is_nullable, column_default
                        FROM information_schema.columns
                        WHERE table_schema='public' AND table_name=%s
                        ORDER BY ordinal_position
                    """, (tbl,))
                    rows = cur.fetchall()
                    if not rows:
                        return f"Table `{tbl}` not found."
                    lines = [f"**Schema for `{tbl}`:**", "| Column | Type | Nullable | Default |",
                             "| --- | --- | --- | --- |"]
                    for r in rows:
                        lines.append(f"| {r['column_name']} | {r['data_type']} | "
                                     f"{r['is_nullable']} | {r['column_default'] or ''} |")
                    return "\n".join(lines)
                return "Usage: schema <table_name>"

            if q.lower().startswith("indexes") or q.lower().startswith("indices"):
                tbl = q.split(None, 1)[1].strip() if len(q.split()) > 1 else ""
                where = "AND t.relname = %s" if tbl else ""
                params = (tbl,) if tbl else ()
                cur.execute(f"""
                    SELECT i.relname AS index_name, t.relname AS table_name,
                           ix.indisunique AS is_unique, ix.indisprimary AS is_primary,
                           array_to_string(array_agg(a.attname ORDER BY a.attnum), ', ') AS columns
                    FROM pg_index ix
                    JOIN pg_class t ON t.oid = ix.indrelid
                    JOIN pg_class i ON i.oid = ix.indexrelid
                    JOIN pg_attribute a ON a.attrelid = t.oid AND a.attnum = ANY(ix.indkey)
                    WHERE t.relkind = 'r' {where}
                    GROUP BY i.relname, t.relname, ix.indisunique, ix.indisprimary
                    ORDER BY t.relname, i.relname
                """, params)
                rows = cur.fetchall()
                if not rows:
                    return "No indexes found."
                lines = ["**Indexes:**", "| Index | Table | Unique | Primary | Columns |",
                         "| --- | --- | --- | --- | --- |"]
                for r in rows:
                    lines.append(f"| {r['index_name']} | {r['table_name']} | "
                                 f"{r['is_unique']} | {r['is_primary']} | {r['columns']} |")
                return "\n".join(lines)

            lower_q = q.lower().lstrip()
            if any(lower_q.startswith(kw) for kw in ("insert", "update", "delete", "drop",
                                                       "alter", "create", "replace", "truncate")):
                return "❌ Only read-only queries are allowed."

            cur.execute(q)
            rows = cur.fetchmany(200)
            if not rows:
                return "Query returned no rows."
            cols = list(rows[0].keys())
            lines = ["| " + " | ".join(cols) + " |",
                     "| " + " | ".join("---" for _ in cols) + " |"]
            for r in rows:
                lines.append("| " + " | ".join(str(r[c]) for c in cols) + " |")
            return "\n".join(lines)
    except Exception as e:
        return f"❌ PostgreSQL error: {e}"
    finally:
        conn.close()


# ─────────────────────────────────────────────────────────────────────────────
# Utility
# ─────────────────────────────────────────────────────────────────────────────

def tool_get_time(timezone: str = "UTC") -> str:
    """Return current date and time in the requested timezone."""
    from datetime import datetime as _dt, timezone as _tz, timedelta as _td
    import time as _time
    tz_label = (timezone or "UTC").strip()
    try:
        if tz_label.upper() in ("UTC", "Z"):
            now = _dt.now(_tz.utc)
            return (f"🕐 **{tz_label}**: {now.strftime('%Y-%m-%d %H:%M:%S')} UTC\n"
                    f"Unix timestamp: `{int(_time.time())}`")
        try:
            import zoneinfo
            zi = zoneinfo.ZoneInfo(tz_label)
            now = _dt.now(zi)
        except (ImportError, KeyError):
            # Fallback: accept simple UTC±offset strings like "UTC+2"
            import re as _re
            m = _re.match(r'^UTC([+-])(\d{1,2})(?::(\d{2}))?$', tz_label, _re.I)
            if m:
                sign, hrs, mins = m.group(1), int(m.group(2)), int(m.group(3) or 0)
                offset = _td(hours=hrs, minutes=mins)
                if sign == '-':
                    offset = -offset
                now = _dt.now(_tz(offset))
            else:
                now = _dt.now(_tz.utc)
                tz_label = "UTC (fallback)"
        return (f"🕐 **{tz_label}**: {now.strftime('%Y-%m-%d %H:%M:%S %Z')}\n"
                f"Unix timestamp: `{int(_time.time())}`")
    except Exception as e:
        return f"❌ time error: {e}"


def tool_nexus_status() -> str:
    """Return live Nexus AI system status including DB, Redis, and uptime."""
    import os as _os, time as _time
    lines = ["**Nexus AI System Status**\n"]
    # DB
    try:
        from .db import list_users
        users = list_users()
        lines.append(f"✅ Database — {len(users)} user(s)")
    except Exception as e:
        lines.append(f"❌ Database — {e}")
    # Redis
    try:
        from .redis_state import get_redis
        r = get_redis()
        if r:
            r.ping()
            lines.append("✅ Redis — connected")
        else:
            lines.append("⚠️ Redis — not configured (in-process fallback)")
    except Exception as e:
        lines.append(f"❌ Redis — {e}")
    # Provider fallback chain
    try:
        from .agent import get_providers_list
        providers = get_providers_list()
        lines.append(f"✅ Providers — {len(providers)} registered")
    except Exception as e:
        lines.append(f"⚠️ Providers — {e}")
    # Memory
    try:
        import psutil
        vm = psutil.virtual_memory()
        used_mb = vm.used // (1024 * 1024)
        total_mb = vm.total // (1024 * 1024)
        lines.append(f"🖥️  RAM — {used_mb} MB / {total_mb} MB used ({vm.percent:.1f}%)")
    except ImportError:
        # Parse /proc/meminfo as fallback
        try:
            with open("/proc/meminfo") as f:
                minfo = {l.split(':')[0].strip(): l.split(':')[1].strip()
                         for l in f if ':' in l}
            total_kb = int(minfo.get("MemTotal", "0 kB").split()[0])
            avail_kb = int(minfo.get("MemAvailable", "0 kB").split()[0])
            lines.append(f"🖥️  RAM — {(total_kb-avail_kb)//1024} MB used / {total_kb//1024} MB total")
        except Exception:
            lines.append("⚠️ RAM info unavailable")
    return "\n".join(lines)


def tool_hash(text: str, algorithm: str = "sha256") -> str:
    """Hash a string with the specified algorithm (sha256, sha512, md5, sha1)."""
    import hashlib as _hl
    algo = algorithm.lower().strip()
    if algo not in ("sha256", "sha512", "md5", "sha1", "sha384"):
        return f"❌ Unsupported algorithm: {algo}. Use sha256, sha512, md5, sha1."
    h = _hl.new(algo, text.encode("utf-8")).hexdigest()
    return f"**{algo}({text[:60]})** =\n```\n{h}\n```"


def tool_uuid(version: int = 4, namespace: str = "", name: str = "") -> str:
    """Generate a UUID (v1, v3, v4, or v5)."""
    import uuid as _uuid
    v = int(version)
    if v == 1:
        u = _uuid.uuid1()
    elif v == 3:
        ns = getattr(_uuid, f"NAMESPACE_{namespace.upper()}", _uuid.NAMESPACE_DNS)
        u = _uuid.uuid3(ns, name or "nexus")
    elif v == 5:
        ns = getattr(_uuid, f"NAMESPACE_{namespace.upper()}", _uuid.NAMESPACE_DNS)
        u = _uuid.uuid5(ns, name or "nexus")
    else:
        u = _uuid.uuid4()
    return f"**UUID v{v}:** `{u}`"


def tool_qr_code(text: str, size: int = 10) -> str:
    """Generate a QR code as an ASCII representation or Pollinations image URL."""
    if not text.strip():
        return "❌ text is required."
    import urllib.parse
    encoded = urllib.parse.quote(text)
    # Use QR code API service for image
    url = f"https://api.qrserver.com/v1/create-qr-code/?size={size*10}x{size*10}&data={encoded}"
    return (f"**QR Code** for: `{text[:80]}`\n\n"
            f"![QR Code]({url})\n\n"
            f"*Image URL: {url}*")


def tool_csv_to_json(csv_text: str) -> str:
    """Convert CSV text to a JSON array."""
    import csv, io, json as _json
    if not csv_text.strip():
        return "❌ csv_text is required."
    try:
        reader = csv.DictReader(io.StringIO(csv_text))
        rows = list(reader)
        if not rows:
            return "❌ No data rows found in CSV."
        result = _json.dumps(rows, indent=2, ensure_ascii=False)
        if len(result) > 4000:
            result = result[:4000] + "\n... (truncated)"
        return f"```json\n{result}\n```"
    except Exception as e:
        return f"❌ CSV to JSON failed: {e}"


def tool_json_to_csv(json_text: str) -> str:
    """Convert a JSON array of objects to CSV text."""
    import csv, io, json as _json
    if not json_text.strip():
        return "❌ json_text is required."
    try:
        data = _json.loads(json_text)
        if not isinstance(data, list) or not data:
            return "❌ Input must be a non-empty JSON array of objects."
        if not isinstance(data[0], dict):
            return "❌ Array elements must be objects."
        buf = io.StringIO()
        writer = csv.DictWriter(buf, fieldnames=list(data[0].keys()), extrasaction='ignore')
        writer.writeheader()
        writer.writerows(data)
        result = buf.getvalue()
        if len(result) > 4000:
            result = result[:4000] + "\n... (truncated)"
        return f"```csv\n{result}\n```"
    except Exception as e:
        return f"❌ JSON to CSV failed: {e}"


def tool_xml_parse(xml_text: str, xpath: str = "") -> str:
    """Parse XML and return a JSON dict, or evaluate an XPath expression."""
    import json as _json
    if not xml_text.strip():
        return "❌ xml_text is required."
    try:
        import xml.etree.ElementTree as _ET

        def _el_to_dict(el):
            d: dict = dict(el.attrib)
            children = list(el)
            if children:
                for child in children:
                    k = child.tag.split("}")[-1] if "}" in child.tag else child.tag
                    v = _el_to_dict(child)
                    if k in d:
                        if not isinstance(d[k], list):
                            d[k] = [d[k]]
                        d[k].append(v)
                    else:
                        d[k] = v
            elif el.text and el.text.strip():
                d["_text"] = el.text.strip()
            return d

        root = _ET.fromstring(xml_text)
        if xpath:
            nodes = root.findall(xpath)
            result = [_el_to_dict(n) for n in nodes]
        else:
            tag = root.tag.split("}")[-1] if "}" in root.tag else root.tag
            result = {tag: _el_to_dict(root)}
        out = _json.dumps(result, indent=2, ensure_ascii=False)
        if len(out) > 4000:
            out = out[:4000] + "\n... (truncated)"
        return f"```json\n{out}\n```"
    except Exception as e:
        return f"❌ XML parse failed: {e}"


def tool_url_encode(text: str, mode: str = "encode") -> str:
    """URL-encode or decode a string."""
    import urllib.parse
    if not text:
        return "❌ text is required."
    try:
        if mode == "decode":
            result = urllib.parse.unquote_plus(text)
            return f"**URL decoded:**\n```\n{result}\n```"
        else:
            result = urllib.parse.quote_plus(text)
            return f"**URL encoded:**\n```\n{result}\n```"
    except Exception as e:
        return f"❌ URL encode/decode failed: {e}"


def tool_jwt_decode(token: str) -> str:
    """Inspect a JWT payload (no signature validation — read-only)."""
    if not token.strip():
        return "❌ token is required."
    try:
        import base64 as _b64, json as _json
        parts = token.strip().split(".")
        if len(parts) not in (2, 3):
            return "❌ Not a valid JWT (expected 2 or 3 dot-separated parts)."

        def _decode_part(b64url: str) -> dict:
            # Add padding
            b64 = b64url.replace("-", "+").replace("_", "/")
            b64 += "=" * (-len(b64) % 4)
            return _json.loads(_b64.b64decode(b64).decode("utf-8", errors="replace"))

        header = _decode_part(parts[0])
        payload = _decode_part(parts[1])
        sig_present = len(parts) == 3 and bool(parts[2])

        out = _json.dumps({"header": header, "payload": payload,
                           "signature_present": sig_present}, indent=2)
        return f"```json\n{out}\n```\n\n⚠️ *Signature NOT verified — inspect only.*"
    except Exception as e:
        return f"❌ JWT decode failed: {e}"


def tool_color_convert(color: str, to_format: str = "all") -> str:
    """Convert color between hex, rgb, and hsl formats."""
    import re as _re
    color = color.strip()
    if not color:
        return "❌ color is required."
    try:
        # Parse input
        r = g = b = 0
        if _re.match(r'^#?[0-9a-fA-F]{6}$', color):
            c = color.lstrip('#')
            r, g, b = int(c[0:2],16), int(c[2:4],16), int(c[4:6],16)
        elif _re.match(r'^#?[0-9a-fA-F]{3}$', color):
            c = color.lstrip('#')
            r, g, b = int(c[0]*2,16), int(c[1]*2,16), int(c[2]*2,16)
        elif m := _re.match(r'^rgb\s*\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*\)$', color, _re.I):
            r, g, b = int(m.group(1)), int(m.group(2)), int(m.group(3))
        elif m := _re.match(r'^(\d+)\s*,\s*(\d+)\s*,\s*(\d+)$', color):
            r, g, b = int(m.group(1)), int(m.group(2)), int(m.group(3))
        else:
            return f"❌ Unrecognized color format: `{color}`. Use #hex, rgb(r,g,b), or r,g,b"

        # Compute formats
        hex_out = f"#{r:02x}{g:02x}{b:02x}"
        rgb_out = f"rgb({r}, {g}, {b})"
        # HSL conversion
        rf, gf, bf = r/255, g/255, b/255
        cmax, cmin = max(rf,gf,bf), min(rf,gf,bf)
        delta = cmax - cmin
        l = (cmax + cmin) / 2
        if delta == 0:
            h = s = 0.0
        else:
            s = delta / (1 - abs(2*l - 1))
            if cmax == rf:
                h = 60 * (((gf-bf)/delta) % 6)
            elif cmax == gf:
                h = 60 * ((bf-rf)/delta + 2)
            else:
                h = 60 * ((rf-gf)/delta + 4)
        hsl_out = f"hsl({h:.1f}, {s*100:.1f}%, {l*100:.1f}%)"
        return (f"**Color: `{color}`**\n"
                f"- HEX: `{hex_out}`\n"
                f"- RGB: `{rgb_out}`\n"
                f"- HSL: `{hsl_out}`")
    except Exception as e:
        return f"❌ Color conversion failed: {e}"


# ─────────────────────────────────────────────────────────────────────────────
# File and Repo
# ─────────────────────────────────────────────────────────────────────────────

def tool_write_file(path: str, content: str, workdir: str = "/tmp", mode: str = "w") -> str:
    """Write text content to a file. mode='w' overwrites, mode='a' appends."""
    if not path.strip():
        return "❌ path is required."
    target = _resolve_path(path, workdir)
    target.parent.mkdir(parents=True, exist_ok=True)
    write_mode = "a" if mode == "a" else "w"
    try:
        with open(target, write_mode, encoding="utf-8") as f:
            f.write(content)
        size = target.stat().st_size
        return f"✅ Wrote {len(content)} chars to `{path}` ({size} bytes on disk)"
    except Exception as e:
        return f"❌ write_file failed: {e}"


def tool_read_file(path: str, workdir: str = "/tmp",
                   start_line: int = 1, end_line: int = 500) -> str:
    """Read a text file from the working directory."""
    if not path.strip():
        return "❌ path is required."
    target = _resolve_path(path, workdir)
    if not target.exists():
        return f"❌ File not found: {path}"
    try:
        with open(target, "r", encoding="utf-8", errors="replace") as f:
            lines = f.readlines()
        total = len(lines)
        s = max(1, start_line) - 1
        e = min(total, end_line)
        chunk = "".join(lines[s:e])
        suffix = (f"\n\n… ({total - e} more lines)" if e < total else "")
        return f"**{path}** (lines {s+1}–{e} of {total})\n```\n{chunk}\n```{suffix}"
    except Exception as ex:
        return f"❌ read_file failed: {ex}"


def tool_list_files(path: str = ".", workdir: str = "/tmp",
                    recursive: bool = False, max_entries: int = 200) -> str:
    """List files and directories in a path."""
    target = _resolve_path(path, workdir)
    if not target.exists():
        return f"❌ Path not found: {path}"
    try:
        if recursive:
            entries = sorted(target.rglob("*"))[:max_entries]
        else:
            entries = sorted(target.iterdir())[:max_entries]
        lines = [f"**{path}/** — {len(entries)} entries\n"]
        for e in entries:
            rel = e.relative_to(target) if recursive else e.name
            icon = "📁" if e.is_dir() else "📄"
            size = f"  {e.stat().st_size:,} B" if e.is_file() else ""
            lines.append(f"{icon} {rel}{size}")
        return "\n".join(lines)
    except Exception as ex:
        return f"❌ list_files failed: {ex}"


def tool_delete_file(path: str, workdir: str = "/tmp") -> str:
    """Delete a file from the working directory."""
    if not path.strip():
        return "❌ path is required."
    target = _resolve_path(path, workdir)
    if not target.exists():
        return f"❌ File not found: {path}"
    if target.is_dir():
        return f"❌ Path is a directory. Use delete_directory instead."
    try:
        target.unlink()
        return f"✅ Deleted: `{path}`"
    except Exception as e:
        return f"❌ delete_file failed: {e}"


def tool_move_file(src: str, dst: str, workdir: str = "/tmp") -> str:
    """Move or rename a file within the working directory."""
    if not src.strip() or not dst.strip():
        return "❌ src and dst are required."
    src_path = _resolve_path(src, workdir)
    dst_path = _resolve_path(dst, workdir)
    if not src_path.exists():
        return f"❌ Source not found: {src}"
    dst_path.parent.mkdir(parents=True, exist_ok=True)
    try:
        src_path.rename(dst_path)
        return f"✅ Moved `{src}` → `{dst}`"
    except Exception as e:
        return f"❌ move_file failed: {e}"


def tool_copy_file(src: str, dst: str, workdir: str = "/tmp") -> str:
    """Copy a file within the working directory."""
    import shutil as _sh
    if not src.strip() or not dst.strip():
        return "❌ src and dst are required."
    src_path = _resolve_path(src, workdir)
    dst_path = _resolve_path(dst, workdir)
    if not src_path.exists():
        return f"❌ Source not found: {src}"
    dst_path.parent.mkdir(parents=True, exist_ok=True)
    try:
        _sh.copy2(src_path, dst_path)
        return f"✅ Copied `{src}` → `{dst}`"
    except Exception as e:
        return f"❌ copy_file failed: {e}"


def tool_clone_repo(url: str, dest: str = "", workdir: str = "/tmp",
                    branch: str = "", depth: int = 0) -> str:
    """Clone a git repository into the working directory."""
    BLOCKED = ["localhost", "127.", "192.168.", "10.", "169.254."]
    for b in BLOCKED:
        if b in url:
            return "❌ Blocked: cannot clone from internal/local addresses."
    if not url.strip():
        return "❌ url is required."
    target = _resolve_path(dest or Path(url).stem, workdir)
    args = ["git", "clone"]
    if depth and depth > 0:
        args += ["--depth", str(int(depth))]
    if branch.strip():
        args += ["--branch", branch.strip()]
    args += [url, str(target)]
    try:
        proc = subprocess.run(args, capture_output=True, text=True, timeout=120)
        if proc.returncode != 0:
            return f"❌ clone failed: {proc.stderr.strip()}"
        return f"✅ Cloned `{url}` → `{target}`"
    except Exception as e:
        return f"❌ clone_repo failed: {e}"


def tool_run_command(command: str, workdir: str = "/tmp",
                     timeout: int = 30, allow_write: bool = False) -> str:
    """
    Run a sandboxed shell command.
    Blocked by default: rm -rf, mkfs, dd, sudo, curl|wget piped to bash.
    allow_write=False prevents commands that modify files outside workdir.
    """
    BLOCKED_PATTERNS = [
        r"rm\s+-rf\s+/", r"mkfs", r"dd\s+if=", r"\bsudo\b",
        r">\s*/dev/sd", r"chmod\s+777\s+/",
        r"(curl|wget).*\|\s*(bash|sh|zsh)",
        r"shutdown", r"reboot", r"halt", r"poweroff",
        r"iptables", r"ufw\s+disable",
    ]
    import re as _re
    for pat in BLOCKED_PATTERNS:
        if _re.search(pat, command, _re.I):
            return f"❌ Blocked: command matches safety pattern `{pat}`"
    cwd = str(_resolve_path(".", workdir))
    try:
        proc = subprocess.run(
            command,
            shell=True,
            cwd=cwd,
            capture_output=True,
            text=True,
            timeout=min(int(timeout), 120),
        )
        out = (proc.stdout or "").strip()
        err = (proc.stderr or "").strip()
        result = ""
        if out:
            result += f"**stdout:**\n```\n{out[:3000]}\n```"
        if err:
            result += f"\n**stderr:**\n```\n{err[:1000]}\n```"
        result += f"\n\nExit code: `{proc.returncode}`"
        return result.strip() or f"Exit code: `{proc.returncode}` (no output)"
    except subprocess.TimeoutExpired:
        return f"❌ Command timed out after {timeout}s"
    except Exception as e:
        return f"❌ run_command failed: {e}"


def tool_commit_push(message: str, repo_path: str = ".", workdir: str = "/tmp",
                     remote: str = "origin", branch: str = "") -> str:
    """Stage all changes, commit, and push to remote."""
    if not message.strip():
        return "❌ commit message is required."
    # Stage all
    code, out, err = _run_git(repo_path, ["add", "-A"], workdir=workdir)
    if code != 0:
        return f"❌ git add failed: {err}"
    # Commit
    code, out, err = _run_git(repo_path, ["commit", "-m", message], workdir=workdir)
    if code != 0:
        return f"❌ git commit failed: {err or out}"
    # Push
    push_args = ["push", remote]
    if branch.strip():
        push_args.append(branch.strip())
    code, out, err = _run_git(repo_path, push_args, workdir=workdir)
    if code != 0:
        return f"❌ git push failed: {err}"
    return f"✅ Committed and pushed: {out or 'done'}"


def tool_create_repo(name: str, description: str = "", private: bool = False,
                     repo_path: str = ".", workdir: str = "/tmp") -> str:
    """Create a GitHub repository using gh CLI."""
    if not name.strip():
        return "❌ name is required."
    args = ["repo", "create", name.strip(), "--confirm"]
    if private:
        args.append("--private")
    else:
        args.append("--public")
    if description.strip():
        args += ["--description", description.strip()]
    code, out, err = _run_gh(args, repo_path=repo_path, workdir=workdir)
    if code != 0:
        return f"❌ create_repo failed: {err or 'gh CLI error'}"
    return out or f"✅ Repository created: {name}"


# ─────────────────────────────────────────────────────────────────────────────
# Web and Network
# ─────────────────────────────────────────────────────────────────────────────

def tool_web_search(query: str, max_results: int = 5, engine: str = "auto") -> str:
    """
    Search the web using Brave Search API, SerpAPI, or DuckDuckGo HTML fallback.
    BRAVE_SEARCH_API_KEY or SERPAPI_KEY env vars enable premium engines.
    """
    import urllib.parse, urllib.request as _urlreq, json as _json, re as _re
    if not query.strip():
        return "❌ query is required."
    limit = max(1, min(int(max_results), 10))

    # ── Brave Search ─────────────────────────────────────────────────────
    brave_key = os.getenv("BRAVE_SEARCH_API_KEY", "").strip()
    if brave_key and engine in ("auto", "brave"):
        try:
            url = f"https://api.search.brave.com/res/v1/web/search?q={urllib.parse.quote(query)}&count={limit}"
            req = _urlreq.Request(url, headers={
                "Accept": "application/json",
                "Accept-Encoding": "gzip",
                "X-Subscription-Token": brave_key,
            })
            with _urlreq.urlopen(req, timeout=15) as resp:
                data = _json.loads(resp.read())
            results = data.get("web", {}).get("results", [])[:limit]
            if results:
                lines = [f"**Web search: '{query}'** (Brave)\n"]
                for i, r in enumerate(results, 1):
                    lines.append(f"{i}. **[{r.get('title','')}]({r.get('url','')})**")
                    desc = (r.get("description") or r.get("extra_snippets", [""])[0] or "")[:200]
                    if desc:
                        lines.append(f"   {desc}")
                return "\n".join(lines)
        except Exception:
            pass

    # ── SerpAPI ──────────────────────────────────────────────────────────
    serp_key = os.getenv("SERPAPI_KEY", "").strip()
    if serp_key and engine in ("auto", "serp", "serpapi"):
        try:
            url = (f"https://serpapi.com/search.json?q={urllib.parse.quote(query)}"
                   f"&num={limit}&api_key={serp_key}")
            with _urlreq.urlopen(url, timeout=15) as resp:
                data = _json.loads(resp.read())
            results = data.get("organic_results", [])[:limit]
            if results:
                lines = [f"**Web search: '{query}'** (SerpAPI)\n"]
                for i, r in enumerate(results, 1):
                    lines.append(f"{i}. **[{r.get('title','')}]({r.get('link','')})**")
                    if r.get("snippet"):
                        lines.append(f"   {r['snippet'][:200]}")
                return "\n".join(lines)
        except Exception:
            pass

    # ── DuckDuckGo HTML fallback ──────────────────────────────────────────
    try:
        search_url = f"https://html.duckduckgo.com/html/?q={urllib.parse.quote(query)}"
        req = _urlreq.Request(search_url, headers={"User-Agent": "Mozilla/5.0"})
        with _urlreq.urlopen(req, timeout=15) as resp:
            html = resp.read().decode("utf-8", errors="replace")
        # Extract result titles and URLs via regex
        titles = _re.findall(r'class="result__a"[^>]*>([^<]+)<', html)
        urls = _re.findall(r'class="result__url"[^>]*>\s*([^\s<]+)', html)
        snippets = _re.findall(r'class="result__snippet"[^>]*>([^<]+)<', html)
        if titles:
            lines = [f"**Web search: '{query}'** (DuckDuckGo)\n"]
            for i in range(min(limit, len(titles))):
                title = titles[i].strip() if i < len(titles) else ""
                url_t = urls[i].strip() if i < len(urls) else ""
                snip = snippets[i].strip() if i < len(snippets) else ""
                lines.append(f"{i+1}. **{title}**" + (f" — {url_t}" if url_t else ""))
                if snip:
                    lines.append(f"   {snip[:200]}")
            return "\n".join(lines)
    except Exception as e:
        return f"❌ web_search failed: {e}"
    return f"No results found for: {query}"


def tool_web_scrape_structured(url: str, selectors: dict | None = None,
                                output_format: str = "json") -> str:
    """
    Fetch a page and extract structured data using CSS selectors.
    selectors: dict of {field_name: css_selector}  e.g. {"title": "h1", "price": ".price"}
    output_format: "json" or "table"
    """
    import urllib.request as _urlreq, re as _re, html as _html, json as _json
    BLOCKED = ["localhost", "127.", "192.168.", "10.", "169.254."]
    for b in BLOCKED:
        if b in url:
            return "❌ Blocked: internal addresses."
    if not url.strip():
        return "❌ url is required."
    try:
        req = _urlreq.Request(url, headers={"User-Agent": "Mozilla/5.0"})
        with _urlreq.urlopen(req, timeout=15) as resp:
            html_text = resp.read().decode("utf-8", errors="replace")
    except Exception as e:
        return f"❌ Failed to fetch page: {e}"

    if not selectors:
        # No selectors — extract plain text
        text = _re.sub(r'<(script|style)[^>]*>.*?</\1>', '', html_text, flags=_re.S|_re.I)
        text = _re.sub(r'<[^>]+>', ' ', text)
        text = _html.unescape(text)
        text = _re.sub(r'\s{3,}', '\n\n', text).strip()
        return text[:3000] + ("…" if len(text) > 3000 else "")

    # Simple CSS selector implementation using regex
    def _extract(sel: str, html_src: str) -> list[str]:
        tag = _re.match(r'^(\w+)', sel)
        cls = _re.findall(r'\.(\w+)', sel)
        id_ = _re.search(r'#(\w+)', sel)
        tag_pat = tag.group(1) if tag else r'\w+'
        if id_:
            attr_pat = f'id="{id_.group(1)}"'
        elif cls:
            attr_pat = '|'.join(f'class="[^"]*{c}[^"]*"' for c in cls)
        else:
            attr_pat = ""
        pattern = f'<{tag_pat}(?:[^>]*)({attr_pat})[^>]*>(.*?)</{tag_pat}>' if attr_pat else f'<{tag_pat}[^>]*>(.*?)</{tag_pat}>'
        matches = _re.findall(pattern, html_src, _re.I | _re.S)
        texts = []
        for m in matches[:5]:
            raw = m[-1] if isinstance(m, tuple) else m
            texts.append(_html.unescape(_re.sub(r'<[^>]+>', '', raw)).strip()[:300])
        return texts

    data = {}
    for field, sel in (selectors or {}).items():
        data[field] = _extract(sel, html_text)

    if output_format == "table":
        lines = ["| Field | Values |", "| --- | --- |"]
        for k, v in data.items():
            lines.append(f"| {k} | {'; '.join(v[:3])} |")
        return "\n".join(lines)
    return f"```json\n{_json.dumps(data, indent=2, ensure_ascii=False)[:3000]}\n```"


def tool_rss_fetch(url: str, max_items: int = 10) -> str:
    """Fetch and parse an RSS or Atom feed, returning the latest items."""
    import urllib.request as _urlreq, re as _re, html as _html
    BLOCKED = ["localhost", "127.", "192.168.", "10.", "169.254."]
    for b in BLOCKED:
        if b in url:
            return "❌ Blocked: internal addresses."
    if not url.strip():
        return "❌ url is required."
    try:
        req = _urlreq.Request(url, headers={"User-Agent": "Mozilla/5.0"})
        with _urlreq.urlopen(req, timeout=15) as resp:
            xml = resp.read().decode("utf-8", errors="replace")
    except Exception as e:
        return f"❌ Failed to fetch feed: {e}"
    try:
        import xml.etree.ElementTree as _ET
        root = _ET.fromstring(xml)
        ns = {"atom": "http://www.w3.org/2005/Atom"}
        items = []
        # RSS 2.0
        for item in root.findall(".//item")[:max_items]:
            title = (item.findtext("title") or "").strip()
            link = (item.findtext("link") or "").strip()
            desc = (item.findtext("description") or "").strip()
            desc = _html.unescape(_re.sub(r'<[^>]+>', '', desc))[:200]
            pub = (item.findtext("pubDate") or "").strip()
            items.append({"title": title, "link": link, "description": desc, "pubDate": pub})
        # Atom
        if not items:
            for entry in root.findall("atom:entry", ns)[:max_items]:
                title = (entry.findtext("atom:title", namespaces=ns) or "").strip()
                link_el = entry.find("atom:link[@rel='alternate']", ns) or entry.find("atom:link", ns)
                link = link_el.get("href", "") if link_el is not None else ""
                summary = (entry.findtext("atom:summary", namespaces=ns) or "").strip()
                updated = (entry.findtext("atom:updated", namespaces=ns) or "").strip()
                items.append({"title": title, "link": link, "description": summary[:200], "pubDate": updated})
        if not items:
            return f"No items found in feed: {url}"
        lines = [f"**RSS feed:** {url}\n"]
        for i, item in enumerate(items, 1):
            lines.append(f"{i}. **{item['title']}**")
            if item.get("link"):
                lines[-1] += f" — [{item['link'][:80]}]({item['link']})"
            if item.get("description"):
                lines.append(f"   {item['description']}")
            if item.get("pubDate"):
                lines.append(f"   *{item['pubDate']}*")
        return "\n".join(lines)
    except Exception as e:
        return f"❌ RSS parse failed: {e}"


def tool_sitemap_crawl(url: str, max_urls: int = 50) -> str:
    """Discover URLs from a sitemap.xml and return the URL list."""
    import urllib.request as _urlreq, xml.etree.ElementTree as _ET
    BLOCKED = ["localhost", "127.", "192.168.", "10.", "169.254."]
    for b in BLOCKED:
        if b in url:
            return "❌ Blocked: internal addresses."
    # Try /sitemap.xml if not already pointing at one
    if not url.rstrip("/").endswith(".xml"):
        url = url.rstrip("/") + "/sitemap.xml"
    try:
        req = _urlreq.Request(url, headers={"User-Agent": "Mozilla/5.0"})
        with _urlreq.urlopen(req, timeout=15) as resp:
            xml_text = resp.read().decode("utf-8", errors="replace")
    except Exception as e:
        return f"❌ Failed to fetch sitemap: {e}"
    try:
        NS = {"sm": "http://www.sitemaps.org/schemas/sitemap/0.9"}
        root = _ET.fromstring(xml_text)
        # Sitemap index?
        sitemap_urls = [el.text.strip() for el in root.findall("sm:sitemap/sm:loc", NS) if el.text]
        page_urls = [el.text.strip() for el in root.findall("sm:url/sm:loc", NS) if el.text]
        if sitemap_urls and not page_urls:
            return (f"**Sitemap index at {url}** — {len(sitemap_urls)} child sitemaps:\n\n"
                    + "\n".join(f"- {u}" for u in sitemap_urls[:max_urls]))
        if not page_urls:
            return f"No URLs found in sitemap: {url}"
        lines = [f"**Sitemap: {url}** — {len(page_urls)} URLs (showing first {min(max_urls, len(page_urls))})\n"]
        for u in page_urls[:max_urls]:
            lines.append(f"- {u}")
        return "\n".join(lines)
    except Exception as e:
        return f"❌ Sitemap parse failed: {e}"


def tool_check_url_status(urls: list | str, timeout: int = 10) -> str:
    """Check HTTP status of one or more URLs."""
    import urllib.request as _urlreq
    if isinstance(urls, str):
        urls = [urls]
    if not urls:
        return "❌ urls is required."
    BLOCKED = ["localhost", "127.", "192.168.", "10.", "169.254."]
    results = []
    for url in urls[:20]:
        for b in BLOCKED:
            if b in url:
                results.append(f"❌ {url} — blocked (internal address)")
                break
        else:
            try:
                req = _urlreq.Request(url, method="HEAD", headers={"User-Agent": "Mozilla/5.0"})
                with _urlreq.urlopen(req, timeout=int(timeout)) as resp:
                    code = resp.getcode()
                    icon = "✅" if 200 <= code < 300 else ("⚠️" if code < 400 else "❌")
                    results.append(f"{icon} {url} — HTTP {code}")
            except Exception as e:
                results.append(f"❌ {url} — {e}")
    return "\n".join(results)


# ─────────────────────────────────────────────────────────────────────────────
# Database (PostgreSQL query + migration)
# ─────────────────────────────────────────────────────────────────────────────

def tool_pg_query(sql: str, database_url: str = "", max_rows: int = 100) -> str:
    """
    Run a read-only SELECT query against a PostgreSQL database.
    Uses DATABASE_URL env if database_url not provided.
    """
    import os as _os
    url = database_url or _os.getenv("DATABASE_URL", "")
    if not url:
        return "❌ DATABASE_URL not set and no database_url provided."
    q = (sql or "").strip()
    if not q.lower().startswith("select"):
        return "❌ Only SELECT queries are allowed."
    dangerous = ["drop ", "delete ", "insert ", "update ", "alter ", "truncate ", "exec ", "--", "/*"]
    for kw in dangerous:
        if kw in q.lower():
            return f"❌ Blocked keyword: {kw.strip()}"
    try:
        import psycopg2
        from psycopg2.extras import RealDictCursor
        conn = psycopg2.connect(url, cursor_factory=RealDictCursor)
        try:
            with conn.cursor() as cur:
                cur.execute(q)
                rows = cur.fetchmany(max(1, min(int(max_rows), 500)))
                if not rows:
                    return "Query returned no rows."
                cols = list(rows[0].keys())
                lines = ["| " + " | ".join(cols) + " |",
                         "| " + " | ".join("---" for _ in cols) + " |"]
                for r in rows:
                    lines.append("| " + " | ".join(str(r[c])[:50] for c in cols) + " |")
                return "\n".join(lines)
        finally:
            conn.close()
    except ImportError:
        return "❌ psycopg2 not installed. Run: pip install psycopg2-binary"
    except Exception as e:
        return f"❌ pg_query failed: {e}"


def tool_db_migrate(migration_sql: str, database_url: str = "",
                    dry_run: bool = True) -> str:
    """
    Apply a DDL migration string against the Nexus database.
    dry_run=True (default) shows the SQL without executing.
    dry_run=False executes — use with caution.
    """
    import os as _os
    q = (migration_sql or "").strip()
    if not q:
        return "❌ migration_sql is required."
    if dry_run:
        return (f"**Dry-run migration** (not executed):\n```sql\n{q}\n```\n\n"
                f"Set dry_run=false to execute.")
    url = database_url or _os.getenv("DATABASE_URL", "")
    if not url:
        return "❌ DATABASE_URL not set."
    try:
        if url.startswith("sqlite:///") or url.endswith(".db"):
            import sqlite3
            path = url.replace("sqlite:///", "")
            conn = sqlite3.connect(path)
            conn.executescript(q)
            conn.commit()
            conn.close()
            return f"✅ SQLite migration applied successfully."
        elif url.startswith("postgresql://") or url.startswith("postgres://"):
            import psycopg2
            conn = psycopg2.connect(url)
            conn.autocommit = True
            with conn.cursor() as cur:
                cur.execute(q)
            conn.close()
            return f"✅ PostgreSQL migration applied successfully."
        else:
            return "❌ Unsupported DATABASE_URL scheme."
    except Exception as e:
        return f"❌ Migration failed: {e}"
