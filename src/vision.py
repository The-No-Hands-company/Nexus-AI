"""Vision helpers for OCR, screenshot capture, and image understanding."""

from __future__ import annotations

import base64
import os
import shutil
import subprocess
import tempfile


VISION_CAPABLE_PROVIDERS = [
    "ollama",
    "claude",
    "gemini",
    "openrouter",
    "github_models",
]


def is_vision_request(messages: list[dict]) -> bool:
    for msg in messages:
        content = msg.get("content", "")
        if isinstance(content, list):
            for part in content:
                if isinstance(part, dict) and part.get("type") == "image_url":
                    return True
    return False


def route_vision_provider(messages: list[dict], preferred_provider: str | None = None) -> str:
    if preferred_provider and preferred_provider in VISION_CAPABLE_PROVIDERS:
        return preferred_provider
    return VISION_CAPABLE_PROVIDERS[0]


def ocr_image_bytes(image_bytes: bytes, mime_type: str = "image/png") -> str:
    ocr_prompt = (
        "You are an OCR engine. Extract and return ALL visible text from this image, "
        "preserving line breaks and formatting as much as possible. "
        "Include any text in signs, documents, code, captions, or any other visible text. "
        "Return only the extracted text, nothing else."
    )
    return describe_image(image_bytes, mime_type=mime_type, prompt=ocr_prompt)


def ocr_image_path(path: str) -> str:
    with open(path, "rb") as f:
        data = f.read()
    mime_type = "image/png" if path.lower().endswith(".png") else "image/jpeg"
    return ocr_image_bytes(data, mime_type=mime_type)


def _capture_with_playwright(url: str, width: int, height: int) -> bytes | None:
    try:
        from playwright.sync_api import sync_playwright
    except Exception:
        return None

    try:
        with sync_playwright() as p:
            browser = p.chromium.launch(headless=True)
            page = browser.new_page(viewport={"width": width, "height": height})
            page.goto(url, wait_until="networkidle", timeout=45000)
            image_bytes = page.screenshot(type="png", full_page=True)
            browser.close()
            return image_bytes
    except Exception:
        return None


def _capture_with_browser_binary(url: str, width: int, height: int) -> bytes | None:
    browser_candidates = [
        "chromium",
        "chromium-browser",
        "google-chrome",
        "google-chrome-stable",
        "chrome",
        "firefox",
    ]
    binary = next((candidate for candidate in browser_candidates if shutil.which(candidate)), None)
    if not binary:
        return None

    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
        output_path = tmp.name
    try:
        if "firefox" in binary:
            cmd = [binary, "--headless", "--screenshot", output_path, "--window-size", f"{width},{height}", url]
        else:
            cmd = [
                binary,
                "--headless=new",
                "--disable-gpu",
                "--hide-scrollbars",
                f"--window-size={width},{height}",
                f"--screenshot={output_path}",
                url,
            ]
        proc = subprocess.run(cmd, capture_output=True, timeout=60)
        if proc.returncode != 0 or not os.path.exists(output_path):
            return None
        with open(output_path, "rb") as f:
            return f.read()
    except Exception:
        return None
    finally:
        try:
            os.unlink(output_path)
        except Exception:
            pass


def capture_screenshot(url: str, width: int = 1280, height: int = 800) -> bytes:
    if not url or not isinstance(url, str):
        raise ValueError("url must be a non-empty string")
    if width <= 0 or height <= 0:
        raise ValueError("width and height must be positive")
    if width > 4096 or height > 4096:
        raise ValueError("width and height must be <= 4096")

    screenshot = _capture_with_playwright(url, width, height)
    if screenshot is None:
        screenshot = _capture_with_browser_binary(url, width, height)
    if screenshot is None:
        raise RuntimeError(
            "No screenshot backend available. Install Playwright or ensure a headless browser binary is present in PATH"
        )
    return screenshot


def describe_image(
    image_bytes: bytes,
    mime_type: str = "image/png",
    prompt: str = "Describe this image in detail.",
    provider: str | None = None,
) -> str:
    from .agent import AllProvidersExhausted, call_llm_with_fallback

    image_b64 = base64.b64encode(image_bytes).decode("utf-8")
    messages = [
        {
            "role": "user",
            "content": [
                {"type": "text", "text": prompt},
                {"type": "image_url", "image_url": {"url": f"data:{mime_type};base64,{image_b64}"}},
            ],
        }
    ]

    try:
        response, _provider_used = call_llm_with_fallback(messages, task="image_description")
        if isinstance(response, dict):
            return response.get("content", str(response))
        return str(response)
    except AllProvidersExhausted:
        raise
    except Exception as e:
        raise RuntimeError(f"Vision description failed: {e}") from e


def describe_image_path(path: str, prompt: str = "Describe this image in detail.") -> str:
    with open(path, "rb") as f:
        data = f.read()
    mime_type = "image/png" if path.lower().endswith(".png") else "image/jpeg"
    return describe_image(data, mime_type=mime_type, prompt=prompt)


# ── Chart / table extraction ──────────────────────────────────────────────────

def extract_charts_and_tables(image_bytes: bytes, mime_type: str = "image/png") -> dict:
    """Extract structured chart data and tables from an image using a vision model.

    Returns a dict with keys:
        charts: list of {type, title, x_label, y_label, data_points}
        tables: list of {headers, rows}
        raw_description: str
    """
    prompt = (
        "Analyze this image carefully. Extract ALL charts and tables.\n\n"
        "For each CHART respond with JSON:\n"
        "  {\"type\": \"bar|line|pie|scatter\", \"title\": \"...\", \"x_label\": \"...\", "
        "\"y_label\": \"...\", \"data_points\": [{\"label\": \"...\", \"value\": ...}]}\n\n"
        "For each TABLE respond with JSON:\n"
        "  {\"headers\": [\"col1\", \"col2\"], \"rows\": [[\"v1\", \"v2\"], ...]}\n\n"
        "Wrap your entire response in a JSON object:\n"
        "{\"charts\": [...], \"tables\": [...], \"summary\": \"...\"}"
    )
    raw = describe_image(image_bytes, mime_type=mime_type, prompt=prompt)
    import json, re
    # Try to parse JSON from the response
    match = re.search(r"\{.*\}", raw, re.DOTALL)
    if match:
        try:
            parsed = json.loads(match.group(0))
            return {
                "charts": parsed.get("charts", []),
                "tables": parsed.get("tables", []),
                "summary": parsed.get("summary", ""),
                "raw_description": raw,
                "ok": True,
            }
        except json.JSONDecodeError:
            pass
    return {"charts": [], "tables": [], "summary": "", "raw_description": raw, "ok": False}


# ── PDF / Office document understanding ──────────────────────────────────────

def understand_pdf(pdf_bytes: bytes, extract_tables: bool = True,
                   extract_images: bool = False) -> dict:
    """Extract text, tables, and optionally images from a PDF.

    Backend priority: pdfplumber > PyMuPDF > pdfminer > raw text fallback
    """
    # keep parameter for forward compatibility
    _ = extract_images
    # Try pdfplumber (best for tables)
    try:
        import pdfplumber  # type: ignore
        import io
        pages = []
        tables_all = []
        with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
            for i, page in enumerate(pdf.pages):
                text = (page.extract_text() or "").strip()
                page_data: dict = {"page": i + 1, "text": text}
                if extract_tables:
                    raw_tables = page.extract_tables() or []
                    page_tables = []
                    for tbl in raw_tables:
                        if tbl:
                            headers = [str(c or "") for c in (tbl[0] or [])]
                            rows = [[str(c or "") for c in row] for row in tbl[1:]]
                            page_tables.append({"headers": headers, "rows": rows})
                    page_data["tables"] = page_tables
                    tables_all.extend(page_tables)
                pages.append(page_data)
        full_text = "\n\n".join(p["text"] for p in pages if p["text"])
        return {"ok": True, "backend": "pdfplumber", "pages": pages,
                "full_text": full_text, "tables": tables_all, "page_count": len(pages)}
    except ImportError:
        pass
    except Exception:
        pass

    # Try PyMuPDF
    try:
        import fitz  # type: ignore  (PyMuPDF)
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        pages = []
        for i in range(len(doc)):
            page = doc[i]
            text = page.get_text("text").strip()
            pages.append({"page": i + 1, "text": text})
        doc.close()
        full_text = "\n\n".join(p["text"] for p in pages if p["text"])
        return {"ok": True, "backend": "pymupdf", "pages": pages,
                "full_text": full_text, "tables": [], "page_count": len(pages)}
    except ImportError:
        pass
    except Exception:
        pass

    return {"ok": False, "error": "No PDF backend available (install pdfplumber or PyMuPDF)",
            "pages": [], "full_text": "", "tables": []}


def understand_office_doc(file_bytes: bytes, filename: str) -> dict:
    """Extract text from DOCX / XLSX / PPTX files."""
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    if ext in ("docx",):
        try:
            import docx  # type: ignore
            import io
            doc = docx.Document(io.BytesIO(file_bytes))
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            return {"ok": True, "format": "docx", "text": text, "paragraphs": len(doc.paragraphs)}
        except ImportError:
            return {"ok": False, "error": "python-docx not installed", "text": ""}
        except Exception as e:
            return {"ok": False, "error": str(e), "text": ""}

    if ext in ("xlsx", "xls"):
        try:
            import openpyxl  # type: ignore
            import io
            wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True)
            sheets = []
            for ws in wb.worksheets:
                rows = [[str(cell.value or "") for cell in row] for row in ws.iter_rows()]
                headers = rows[0] if rows else []
                data = rows[1:] if len(rows) > 1 else []
                sheets.append({"name": ws.title, "headers": headers, "rows": data[:500]})
            return {"ok": True, "format": "xlsx", "sheets": sheets}
        except ImportError:
            return {"ok": False, "error": "openpyxl not installed", "sheets": []}
        except Exception as e:
            return {"ok": False, "error": str(e), "sheets": []}

    if ext in ("pptx",):
        try:
            from pptx import Presentation  # type: ignore
            import io
            prs = Presentation(io.BytesIO(file_bytes))
            slides = []
            for i, slide in enumerate(prs.slides):
                texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
                slides.append({"slide": i + 1, "content": texts})
            return {"ok": True, "format": "pptx", "slides": slides}
        except ImportError:
            return {"ok": False, "error": "python-pptx not installed", "slides": []}
        except Exception as e:
            return {"ok": False, "error": str(e), "slides": []}

    return {"ok": False, "error": f"Unsupported format: {ext}"}


# ── Document comparison / diff ────────────────────────────────────────────────

def diff_documents(text_a: str, text_b: str, context_lines: int = 3) -> dict:
    """Produce a unified diff and change summary between two text documents."""
    import difflib
    lines_a = text_a.splitlines(keepends=True)
    lines_b = text_b.splitlines(keepends=True)
    diff = list(difflib.unified_diff(lines_a, lines_b, lineterm="", n=context_lines))
    added = sum(1 for l in diff if l.startswith("+") and not l.startswith("+++"))
    removed = sum(1 for l in diff if l.startswith("-") and not l.startswith("---"))
    sm = difflib.SequenceMatcher(None, text_a, text_b)
    return {
        "ok": True,
        "unified_diff": "".join(diff),
        "added_lines": added,
        "removed_lines": removed,
        "similarity": round(sm.ratio(), 4),
        "changed": added > 0 or removed > 0,
    }
